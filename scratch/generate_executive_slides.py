import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_executive_deck():
    prs = Presentation()
    # 16:9 widescreen format
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Executive Color Palette
    PRIMARY = RGBColor(10, 37, 74)        # #0A254A - Deep Navy PTSC
    SECONDARY = RGBColor(0, 82, 204)      # #0052CC - Vibrant Enterprise Blue
    ACCENT_GREEN = RGBColor(16, 149, 93)  # #10955D - Success/Recommended Choice
    ACCENT_RED = RGBColor(209, 53, 53)    # #D13535 - Risk/Warning
    ACCENT_GOLD = RGBColor(224, 138, 0)   # #E08A00 - Priority Badge / Gold
    CARD_BG = RGBColor(255, 255, 255)
    CARD_BORDER = RGBColor(218, 226, 237)
    SECTION_BG = RGBColor(246, 248, 251)
    TEXT_DARK = RGBColor(26, 38, 57)
    TEXT_MUTED = RGBColor(100, 116, 139)
    WHITE = RGBColor(255, 255, 255)

    def add_header(slide, title_text, category_text="BÁO CÁO THỰC HIỆN CHỈ ĐẠO MỤC 9 CỦA BAN GIÁM ĐỐC"):
        # Top banner background
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.15))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = PRIMARY
        top_bar.line.fill.background()

        # Category text
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.12), Inches(11.5), Inches(0.28))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = RGBColor(186, 215, 245)

        # Title text
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.6))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(21)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Bottom gold stripe
        sub_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.15), Inches(13.333), Inches(0.04))
        sub_line.fill.solid()
        sub_line.fill.fore_color.rgb = ACCENT_GOLD
        sub_line.line.fill.background()

        # Footer
        footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.12), Inches(11.733), Inches(0.3))
        tf_foot = footer_box.text_frame
        p_foot = tf_foot.paragraphs[0]
        p_foot.text = "PTSC Quảng Ngãi | Báo cáo Phương án Nền tảng Dữ liệu & Trục tích hợp (Chỉ đạo số 9 BGĐ)"
        p_foot.font.size = Pt(9)
        p_foot.font.color.rgb = TEXT_MUTED

    def add_card(slide, left, top, width, height, title, subtitle="", border_color=CARD_BORDER, bg_color=CARD_BG, top_accent_color=None):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.2)

        if top_accent_color:
            accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.05), top + Inches(0.02), width - Inches(0.1), Inches(0.08))
            accent.fill.solid()
            accent.fill.fore_color.rgb = top_accent_color
            accent.line.fill.background()

        tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), height - Inches(0.25))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = PRIMARY
        
        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.font.size = Pt(10.5)
            p2.font.color.rgb = TEXT_MUTED
            p2.space_after = Pt(6)
            
        return tf

    # =========================================================================
    # SLIDE 1: COVER
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = PRIMARY
    bg1.line.fill.background()

    # Left decorative bar
    bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.6), Inches(0.15), Inches(4.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_GOLD
    bar.line.fill.background()

    # Text box cover
    tb_c = s1.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(11.2), Inches(2.5))
    tf_c = tb_c.text_frame
    tf_c.word_wrap = True

    p = tf_c.paragraphs[0]
    p.text = "BÁO CÁO BAN GIÁM ĐỐC CÔNG TY (THỰC HIỆN CHỈ ĐẠO MỤC 9)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(186, 215, 245)
    p.space_after = Pt(8)

    p2 = tf_c.add_paragraph()
    p2.text = "PHƯƠNG ÁN TRIỂN KHAI DATA PLATFORM\n& TRỤC TÍCH HỢP DỮ LIỆU CÔNG TY VS TCT"
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.space_after = Pt(14)

    p3 = tf_c.add_paragraph()
    p3.text = "• So sánh 2 Phương án: Thuê TCT vs Phát triển riêng đồng bộ TCT (có NCC tư vấn)\n• Mô hình Kiến trúc Dữ liệu tổng quan & Trục tích hợp phần mềm Cty vs TCT (Ưu tiên cao)\n• Khung Chi phí Dự toán trước khi đàm phán NCC & Kế hoạch triển khai chi tiết"
    p3.font.size = Pt(13)
    p3.font.color.rgb = RGBColor(220, 230, 242)

    # Info card bottom
    box_meta = s1.shapes.add_textbox(Inches(1.2), Inches(5.4), Inches(10.5), Inches(1.2))
    tf_meta = box_meta.text_frame
    pm1 = tf_meta.paragraphs[0]
    pm1.text = "• Đơn vị thực hiện: Tổ Công tác Chuyển đổi số & Bộ phận CNTT PTSC Quảng Ngãi"
    pm1.font.size = Pt(12.5)
    pm1.font.color.rgb = WHITE
    pm2 = tf_meta.add_paragraph()
    pm2.text = "• Thời gian báo cáo: Tháng 03/2026 | Kính trình: Ban Giám đốc Công ty xem xét và chỉ đạo"
    pm2.font.size = Pt(12.5)
    pm2.font.color.rgb = RGBColor(253, 224, 71)

    s1.notes_slide.notes_text_frame.text = (
        "KỊCH BẢN NÓI:\n"
        "Kính thưa Ban Giám đốc, thực hiện đúng nội dung chỉ đạo số 9 của Ban Giám đốc về việc đánh giá phương án triển khai "
        "Data Platform, hôm nay bộ phận CNTT xin báo cáo chi tiết 4 nội dung trọng tâm: Thứ nhất là so sánh 2 phương án (thuê TCT hay phát triển riêng); "
        "thứ hai là làm rõ Trục tích hợp dữ liệu giữa Cty và TCT; thứ ba là khung kiến trúc và chi phí dự toán trước khi đàm phán với NCC; "
        "và thứ tư là kế hoạch hành động cụ thể để xin phê duyệt chủ trương."
    )

    # =========================================================================
    # SLIDE 2: TÓM TẮT ĐIỀU HÀNH - 4 CÂU TRẢ LỜI CHO BAN GIÁM ĐỐC
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    add_header(s2, "TỔNG QUAN: TRẢ LỜI TRỰC TIẾP 4 YÊU CẦU TRONG CHỈ ĐẠO SỐ 9", "TÓM TẮT ĐIỀU HÀNH (EXECUTIVE SUMMARY)")

    card_data_s2 = [
        ("1. SO SÁNH 2 PHƯƠNG ÁN", 
         "Thuê TCT vs Phát triển riêng", 
         "• KHUYẾN NGHỊ: Chọn Phương án Phát triển riêng đồng bộ TCT (có NCC tư vấn).\n"
         "• Lý do: TCT chỉ làm hạ tầng lõi Hub, không đủ nguồn lực xuống cấu hình từng phần mềm con tại Quảng Ngãi; tự làm giúp chủ động tiến độ và đo ni đóng giày cho BGĐ.",
         Inches(0.8), ACCENT_GREEN),
        ("2. TRỤC TÍCH HỢP CÔNG TY VS TCT", 
         "Ưu tiên cao - Thiết kế kỹ thuật", 
         "• Mô hình Trạm trung chuyển (Data Gateway) đặt tại Quảng Ngãi.\n"
         "• Cơ chế: Trích xuất chỉ đọc (Read-only) từ Kế toán, Vật tư, Nhân sự -> Chuẩn hóa theo 29 Master Data của TCT -> Bơm qua VPN lên Cloud TCT an toàn 100%.",
         Inches(3.8), SECONDARY),
        ("3. KIẾN TRÚC & DỰ TOÁN CHI PHÍ", 
         "Chủ động trước khi gặp NCC", 
         "• Cty tự nắm khung kiến trúc lõi để không bị NCC 'vẽ' thêm tính năng tốn tiền.\n"
         "• Dự toán trọn gói: Khoảng 350 - 550 triệu VNĐ (tùy phạm vi báo cáo Power BI).\n"
         "• Tận dụng tối đa hạ tầng TCT để tiết kiệm tiền tỷ tiền máy chủ/bản quyền.",
         Inches(6.8), ACCENT_GOLD),
        ("4. KẾ HOẠCH THỰC HIỆN CỤ THỂ", 
         "Phân kỳ 3 Giai đoạn rõ ràng", 
         "• Giai đoạn 1 (Tháng 1-2): Khảo sát nội bộ & chuẩn hóa Master Data (0 VNĐ).\n"
         "• Giai đoạn 2 (Tháng 3-5): Thuê NCC triển khai Trục tích hợp & kết nối TCT.\n"
         "• Giai đoạn 3 (Tháng 6 trở đi): Nghiệm thu, bàn giao mã nguồn, đào tạo IT làm chủ.",
         Inches(9.8), PRIMARY)
    ]

    for title, sub, body, left_pos, color_top in card_data_s2:
        tf2 = add_card(s2, left_pos, Inches(1.5), Inches(2.8), Inches(5.2), title, sub, top_accent_color=color_top)
        p = tf2.add_paragraph()
        p.text = body
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(6)

    s2.notes_slide.notes_text_frame.text = (
        "KỊCH BẢN NÓI:\n"
        "Thưa các anh, slide này tóm tắt 4 câu trả lời cốt lõi: Về phương án, chúng tôi đề xuất phát triển riêng đồng bộ TCT có NCC hỗ trợ. "
        "Về trục tích hợp, sẽ dùng trạm trung chuyển gom dữ liệu an toàn. Về ngân sách, dự toán sơ bộ khoảng 350 - 550 triệu, phân kỳ theo nghiệm thu. "
        "Và về kế hoạch, giai đoạn 1 nội bộ tự làm hoàn toàn miễn phí để rà soát dữ liệu."
    )

    # =========================================================================
    # SLIDE 3: SO SÁNH CHI TIẾT 2 PHƯƠNG ÁN (TRỌNG TÂM 1)
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    add_header(s3, "1. SO SÁNH 2 PHƯƠNG ÁN: THUÊ TCT HAY PHÁT TRIỂN RIÊNG ĐỒNG BỘ TCT?", "CĂN CỨ KỸ THUẬT VÀ QUẢN TRỊ ĐỂ LỰA CHỌN PHƯƠNG ÁN TỐI ƯU")

    table_shape3 = s3.shapes.add_table(6, 3, Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.2))
    table3 = table_shape3.table
    table3.columns[0].width = Inches(2.4)
    table3.columns[1].width = Inches(4.6)
    table3.columns[2].width = Inches(4.733)

    headers3 = ["Tiêu chí đánh giá", "Phương án 1: Thuê dịch vụ trọn gói từ TCT", "Phương án 2: Phát triển riêng đồng bộ TCT (ĐỀ XUẤT)"]
    for i, h in enumerate(headers3):
        cell = table3.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY if i < 2 else ACCENT_GREEN
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.bold = True
        p.font.size = Pt(12.5)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    data3 = [
        ("Nguồn lực & Khả năng hỗ trợ từ TCT",
         "• BẤT KHẢ THI: Ban CNTT của TCT chỉ phụ trách hạ tầng lõi Hub chung toàn tổng.\n• TCT không có đủ nhân sự để xuống tận Quảng Ngãi khảo sát từng CSDL kế toán, vật tư riêng lẻ của mình.",
         "• KHẢ THI CAO: Thuê Nhà cung cấp (NCC) chuyên nghiệp làm việc trực tiếp tại Quảng Ngãi để đấu nối các phần mềm đặc thù của đơn vị."),
        ("Tính chủ động & Tiến độ triển khai",
         "• Rất bị động. Phải xếp hàng chờ TCT triển khai lần lượt cho hàng chục đơn vị thành viên khác.\n• Chắc chắn trễ hạn mốc năm 2026 do HĐQT giao.",
         "• Chủ động 100%. Đơn vị kiểm soát tiến độ, cam kết hoàn thành đường ống dữ liệu trong 3 - 4 tháng."),
        ("Đáp ứng nhu cầu quản trị của Ban Giám đốc QN",
         "• Kém linh hoạt: TCT chỉ xây báo cáo phục vụ cấp Tập đoàn (số liệu tài chính vĩ mô).\n• Không phục vụ được các bài toán quản trị chi tiết của QN (vật tư tồn xưởng, nhân công dự án...).",
         "• Tối ưu riêng: Vừa đồng bộ báo cáo cho TCT, vừa tùy biến xây dựng Dashboard quản trị chuyên sâu phục vụ đúng các bài toán điều hành của Ban Giám đốc QN."),
        ("Bản quyền & Chi phí duy trì",
         "• Phải trả phí dịch vụ duy trì hàng năm cho TCT hoặc gánh phân bổ chi phí phần mềm lớn từ TCT.",
         "• Tận dụng miễn phí không gian Cloud do TCT đã mua sẵn, chỉ chi trả 1 lần tiền triển khai ban đầu cho NCC."),
        ("KẾT LUẬN & ĐỀ XUẤT",
         "❌ KHÔNG KHẢ THI VỀ TIẾN ĐỘ VÀ NGUỒN LỰC TCT",
         "✔ CHỌN PHƯƠNG ÁN 2 (CẦN SỰ TƯ VẤN CỦA NCC)")
    ]

    for row_idx, row_data in enumerate(data3, start=1):
        for col_idx, text in enumerate(row_data):
            cell = table3.cell(row_idx, col_idx)
            cell.fill.solid()
            if row_idx == 5:
                cell.fill.fore_color.rgb = RGBColor(254, 226, 226) if col_idx == 1 else RGBColor(220, 252, 231)
            else:
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255) if col_idx > 0 else RGBColor(241, 245, 249)
            p = cell.text_frame.paragraphs[0]
            p.text = text
            p.font.size = Pt(10.5)
            p.font.bold = (col_idx == 0 or row_idx == 5)
            if row_idx == 5:
                p.font.color.rgb = ACCENT_RED if col_idx == 1 else ACCENT_GREEN
            else:
                p.font.color.rgb = PRIMARY if col_idx == 0 else TEXT_DARK

    s3.notes_slide.notes_text_frame.text = (
        "KỊCH BẢN NÓI (GIẢI THÍCH TRỌNG TÂM CÂU HỎI CỦA SẾP):\n"
        "Thưa Ban Giám đốc, nhiều người sẽ hỏi tại sao không nhờ luôn TCT làm cho tiện? "
        "Lý do thực tế là: Ban CNTT TCT chỉ xây cái 'sân ga trung tâm' và ban hành luật chơi. Họ không có người để đi vào từng ngõ ngách, "
        "ngồi bóc tách từng bảng CSDL kế toán hay vật tư riêng của Quảng Ngãi. Nếu chờ TCT thì Quảng Ngãi sẽ bị xếp hàng rất lâu và trễ hẹn 2026. "
        "Do đó, phương án tối ưu nhất là chúng ta tự thuê NCC ngoài dựng trạm kết nối riêng, vừa chuẩn theo TCT, vừa có báo cáo riêng cho Sếp."
    )

    # =========================================================================
    # SLIDE 4: MÔ HÌNH TỔNG QUAN KIẾN TRÚC DỮ LIỆU (TRỌNG TÂM 2)
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    add_header(s4, "2. MÔ HÌNH TỔNG QUAN KIẾN TRÚC DỮ LIỆU CÔNG TY", "KHUNG KIẾN TRÚC CHỦ ĐỘNG XÂY DỰNG TRƯỚC KHI LÀM VIỆC VỚI NCC")

    # 3 Layers of Architecture
    arch_cols = [
        ("TẦNG 1: NGUỒN DỮ LIỆU NỘI BỘ\n(Operational Source Data)",
         "Cơ sở dữ liệu các phần mềm hiện hữu",
         [
             "Phần mềm Kế toán - Tài chính (SQL Server/Oracle)",
             "Phần mềm Quản lý Vật tư & Thiết bị xưởng",
             "Phần mềm Quản lý Nhân sự - Chấm công",
             "Hệ thống Quản lý Hợp đồng & Dự án thi công",
             "Các file dữ liệu Excel quản lý chuyên biệt"
         ], Inches(0.8), PRIMARY),
        ("TẦNG 2: VÙNG XỬ LÝ & TRẠM TRUNG CHUYỂN\n(Integration & Staging Area)",
         "Trọng tâm phát triển riêng (Cần NCC)",
         [
             "Máy chủ ảo (VM) làm Trạm trung chuyển dữ liệu",
             "Module Ánh xạ: Quy đổi sang 29 Master Data của TCT",
             "Module Kiểm tra chất lượng dữ liệu (Data Quality)",
             "Đường ống tự động ETL/ELT (Trích xuất - Đóng gói)",
             "Lưu trữ vùng đệm an toàn (Staging Database)"
         ], Inches(4.8), ACCENT_GOLD),
        ("TẦNG 3: TẦNG PHỤC VỤ & ĐỒNG BỘ\n(Serving & Synchronization)",
         "Đầu ra phục vụ BGĐ và Tổng Công Ty",
         [
             "ĐẦU RA 1 - ĐỒNG BỘ TỔNG CÔNG TY:\n• Bơm qua VPN an toàn vào Workspace L3 trên Microsoft Fabric của TCT.",
             "ĐẦU RA 2 - PHỤC VỤ BAN GIÁM ĐỐC QN:\n• Bộ Dashboard Power BI quản trị theo thời gian thực (Doanh thu, dòng tiền, nhân lực, tồn kho)."
         ], Inches(8.8), SECONDARY)
    ]

    for title, sub, items, left_pos, col in arch_cols:
        tf4 = add_card(s4, left_pos, Inches(1.5), Inches(3.7), Inches(5.2), title, sub, top_accent_color=col)
        for it in items:
            p = tf4.add_paragraph()
            p.text = "• " + it if not it.startswith("ĐẦU RA") else it
            p.font.size = Pt(11)
            p.font.bold = it.startswith("ĐẦU RA")
            p.font.color.rgb = SECONDARY if it.startswith("ĐẦU RA") else TEXT_DARK
            p.space_after = Pt(6)

    s4.notes_slide.notes_text_frame.text = (
        "KỊCH BẢN NÓI:\n"
        "Thưa Ban Giám đốc, đây là mô hình kiến trúc do chúng tôi chủ động thiết kế trước khi gặp nhà cung cấp. "
        "Việc mình nắm chắc kiến trúc 3 tầng này giúp mình hoàn toàn làm chủ cuộc chơi, khi NCC vào mình là người ra đề bài "
        "chính xác, tránh trường hợp bị NCC tư vấn lan man hoặc bán thêm các phân hệ đắt tiền không cần thiết."
    )

    # =========================================================================
    # SLIDE 5: TRỤC TÍCH HỢP CÔNG TY VS TCT (ƯU TIÊN CAO - TRỌNG TÂM 3)
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    add_header(s5, "3. TRỤC TÍCH HỢP PHẦN MỀM CÔNG TY VS TCT HOẠT ĐỘNG THẾ NÀO?", "CHI TIẾT KỸ THUẬT KẾT NỐI - NỘI DUNG ƯU TIÊN CAO THEO CHỈ ĐẠO")

    # Left: 4 bước luồng tích hợp
    tf5_left = add_card(s5, Inches(0.8), Inches(1.5), Inches(7.5), Inches(5.2), "Cơ chế vận hành của Trục Tích hợp (Data Integration Hub)", "NGUYÊN LÝ 4 BƯỚC TỰ ĐỘNG HÓA HOÀN TOÀN", top_accent_color=SECONDARY)
    
    steps5 = [
        ("BƯỚC 1: KẾT NỐI AN TOÀN NỘI BỘ (Read-Only)",
         "Trục tích hợp chỉ kết nối với các phần mềm Kế toán, Nhân sự, Vật tư ở chế độ CHỈ ĐỌC (hoặc qua các View/API trung gian). Tuyệt đối không can thiệp hay làm chậm hệ thống đang sản xuất."),
        ("BƯỚC 2: BỘ LỌC VÀ QUY ĐỔI MASTER DATA (ETL Engine)",
         "Dữ liệu thô được đưa qua module chuẩn hóa: Tự động ánh xạ mã nhân viên, mã vật tư, mã dự án của QN khớp chính xác với 29 bộ danh mục Master Data do TCT quy định."),
        ("BƯỚC 3: MÃ HÓA VÀ TRUYỀN DẪN QUA VPN BẢO MẬT",
         "Gói dữ liệu đã làm sạch được mã hóa bảo mật chuẩn ngân hàng (AES-256) và đẩy qua đường hầm mạng riêng (VPN Site-to-Site) nối thẳng về Trung tâm Dữ liệu của TCT."),
        ("BƯỚC 4: TIẾP ĐẤT VÀ TỰ ĐỘNG CẬP NHẬT BÁO CÁO",
         "Dữ liệu đổ vào 'Workspace L3 Quảng Ngãi' trên Cloud TCT. Từ đây, Dashboard Power BI của Ban Giám đốc tự động làm mới (refresh) mỗi ngày/mỗi tuần theo lịch hẹn.")
    ]

    for st, sd in steps5:
        p = tf5_left.add_paragraph()
        r1 = p.add_run()
        r1.text = st + "\n"
        r1.font.bold = True
        r1.font.size = Pt(11.5)
        r1.font.color.rgb = PRIMARY
        
        r2 = p.add_run()
        r2.text = sd
        r2.font.size = Pt(10.5)
        r2.font.color.rgb = TEXT_DARK
        p.space_after = Pt(8)

    # Right: 3 Cam kết bảo mật then chốt
    tf5_right = add_card(s5, Inches(8.7), Inches(1.5), Inches(3.8), Inches(5.2), "3 Cam kết kỹ thuật then chốt", "AN TOÀN TUYỆT ĐỐI CHO DOANH NGHIỆP", top_accent_color=ACCENT_GREEN)
    
    safe_points = [
        ("1. Không lo sập phần mềm:", "Trục tích hợp chạy trên máy chủ ảo riêng biệt, lấy dữ liệu vào ban đêm hoặc giờ thấp điểm, không làm chậm phần mềm kế toán."),
        ("2. Không rò rỉ dữ liệu ra ngoài:", "Dữ liệu chỉ đi nội bộ qua đường truyền riêng VPN về TCT, không mở cổng internet công cộng, đáp ứng chuẩn an toàn thông tin TCT."),
        ("3. Tự động cảnh báo khi lỗi:", "Nếu có sai lệch mã hoặc nghẽn mạng, hệ thống tự gửi email cảnh báo cho IT xử lý ngay lập tức.")
    ]
    for sp_t, sp_d in safe_points:
        p = tf5_right.add_paragraph()
        p.text = sp_t
        p.font.bold = True
        p.font.size = Pt(11.5)
        p.font.color.rgb = SECONDARY
        p.space_after = Pt(2)
        pd = tf5_right.add_paragraph()
        pd.text = sp_d
        pd.font.size = Pt(10.5)
        pd.font.color.rgb = TEXT_DARK
        pd.space_after = Pt(8)

    s5.notes_slide.notes_text_frame.text = (
        "KỊCH BẢN NÓI (NỘI DUNG ƯU TIÊN CAO):\n"
        "Thưa Ban Giám đốc, đây là nội dung Ban Giám đốc chỉ đạo ưu tiên cao: Trục tích hợp giữa Cty và TCT hoạt động thế nào? "
        "Chúng tôi xin cam kết bằng 3 nguyên tắc: 1 là chế độ chỉ đọc không bao giờ làm hỏng dữ liệu phần mềm hiện tại; "
        "2 là tự động nắn dữ liệu của mình khớp với 29 bộ mã của TCT; và 3 là truyền bằng đường hầm bảo mật riêng, "
        "chỉ đi từ máy chủ QN sang Cloud của TCT, không hề lọt ra ngoài internet."
    )

    # =========================================================================
    # SLIDE 6: KHUNG DỰ TOÁN CHI PHÍ (TRỌNG TÂM 4)
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    add_header(s6, "4. KHUNG DỰ TOÁN CHI PHÍ SƠ BỘ TRƯỚC KHI LÀM VIỆC VỚI NCC", "CƠ SỞ NGÂN SÁCH ĐỂ ĐÀM PHÁN GIÁ & TỐI ƯU HÓA CHI PHÍ DOANH NGHIỆP")

    # Table of Budget
    table_shape6 = s6.shapes.add_table(6, 4, Inches(0.8), Inches(1.5), Inches(11.733), Inches(3.8))
    table6 = table_shape6.table
    table6.columns[0].width = Inches(0.6)
    table6.columns[1].width = Inches(4.3)
    table6.columns[2].width = Inches(4.833)
    table6.columns[3].width = Inches(2.0)

    headers6 = ["STT", "Hạng mục công việc", "Nội dung chi tiết & Phạm vi thực hiện", "Dự toán (VNĐ)"]
    for i, h in enumerate(headers6):
        cell = table6.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.bold = True
        p.font.size = Pt(11.5)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    cost_data = [
        ("1", "Hạ tầng máy chủ Trạm trung chuyển", "Tận dụng máy chủ hiện hữu của Công ty, tạo 1-2 Máy chủ ảo (VM)", "0 VNĐ (Tự có)"),
        ("2", "Thuê NCC: Khảo sát & Thiết kế chi tiết", "Khảo sát CSDL 4 phần mềm, thiết kế kiến trúc và bảng ánh xạ 29 Master Data", "80 - 120 Triệu"),
        ("3", "Thuê NCC: Xây dựng Trục tích hợp & ETL", "Lập trình đường ống trích xuất, làm sạch, mã hóa và bơm dữ liệu lên TCT qua VPN", "150 - 250 Triệu"),
        ("4", "Thuê NCC: Xây dựng Dashboard Power BI", "Xây dựng 3-5 Bảng báo cáo quản trị thông minh theo yêu cầu riêng của Ban Giám đốc QN", "100 - 150 Triệu"),
        ("5", "Đào tạo chuyển giao & Bảo hành 12 tháng", "Bàn giao mã nguồn, đào tạo kỹ sư IT quản trị làm chủ 100%, hỗ trợ kỹ thuật 1 năm", "50 - 80 Triệu")
    ]

    for row_idx, row_data in enumerate(cost_data, start=1):
        for col_idx, text in enumerate(row_data):
            cell = table6.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255) if col_idx != 3 else RGBColor(240, 249, 255)
            p = cell.text_frame.paragraphs[0]
            p.text = text
            p.font.size = Pt(10.5)
            p.font.bold = (col_idx == 3)
            p.font.color.rgb = SECONDARY if col_idx == 3 else TEXT_DARK
            if col_idx == 0 or col_idx == 3:
                p.alignment = PP_ALIGN.CENTER

    # Summary box
    box_total = s6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(5.5), Inches(11.733), Inches(1.35))
    box_total.fill.solid()
    box_total.fill.fore_color.rgb = RGBColor(241, 245, 249)
    box_total.line.color.rgb = ACCENT_GOLD
    box_total.line.width = Pt(1.5)

    tf_tot = box_total.text_frame
    p_t1 = tf_tot.paragraphs[0]
    p_t1.text = "TỔNG KHÁI TOÁN DỰ KIẾN: KHOẢNG 380 - 600 TRIỆU ĐỒNG (THANH TOÁN THEO MỐC NGHIỆM THU)"
    p_t1.font.bold = True
    p_t1.font.size = Pt(13)
    p_t1.font.color.rgb = PRIMARY
    p_t1.space_after = Pt(3)

    p_t2 = tf_tot.add_paragraph()
    p_t2.text = (
        "• Hiệu quả kinh tế: Giúp đơn vị tiết kiệm hàng tỷ đồng vì không phải đầu tư máy chủ Data Center và bản quyền phần mềm phân tích.\n"
        "• Cơ chế kiểm soát: Chỉ giải ngân theo từng mốc kết quả nghiệm thu (Giai đoạn 1 chuẩn bị: 0 VNĐ -> Giai đoạn 2 kết nối xong: 70% -> Giai đoạn 3 bàn giao: 30%)."
    )
    p_t2.font.size = Pt(10.5)
    p_t2.font.color.rgb = TEXT_DARK

    s6.notes_slide.notes_text_frame.text = (
        "KỊCH BẢN NÓI:\n"
        "Về chi phí, chúng tôi đã khảo sát mặt bằng thị trường để xây dựng khung dự toán trước khi tiếp xúc NCC: "
        "Tổng chi phí trọn gói dao động từ 380 đến 600 triệu đồng. Con số này rất hợp lý so với một dự án dữ liệu, "
        "bởi vì phần hạ tầng đắt nhất là Cloud và bản quyền Power BI thì TCT đã gánh cho chúng ta. "
        "Chúng ta chỉ trả tiền công cho NCC dựng đường ống kết nối và vẽ báo cáo quản trị cho Sếp."
    )

    # =========================================================================
    # SLIDE 7: KẾ HOẠCH THỰC HIỆN CỤ THỂ (ĐÍNH KÈM CHI TIẾT)
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    add_header(s7, "5. KẾ HOẠCH THỰC HIỆN CỤ THỂ (ĐÍNH KÈM CHI TIẾT)", "PHÂN KỲ 3 GIAI ĐOẠN - THAM CHIẾU KẾ HOẠCH CHI TIẾT ĐÍNH KÈM")

    plan_cards = [
        ("GIAI ĐOẠN 1: KHẢO SÁT & CHUẨN HÓA\n(Tháng 03 - Tháng 04/2026)",
         "NỘI BỘ TỰ CHỦ TRÌ (CHI PHÍ 0 VNĐ)",
         [
             "Thành lập Tổ công tác Data Platform nội bộ.",
             "Khảo sát hiện trạng CSDL 4 phần mềm nghiệp vụ (Kế toán, Vật tư, Nhân sự, Dự án).",
             "Rà soát và lập bảng ánh xạ với 29 Master Data của TCT.",
             "Lập Hồ sơ yêu cầu kỹ thuật (TOR) để mời các NCC gửi đề xuất kỹ thuật & báo giá.",
             "*(Chi tiết xem file Kế hoạch Giai đoạn 1 đính kèm)*"
         ], Inches(0.8), PRIMARY),
        ("GIAI ĐOẠN 2: LỰA CHỌN NCC & TRIỂN KHAI\n(Tháng 05 - Tháng 07/2026)",
         "THUÊ NCC CHUYÊN TRÁCH (GIAI ĐOẠN CHÍNH)",
         [
             "Đánh giá hồ sơ năng lực và lựa chọn NCC có giải pháp tối ưu, giá cạnh tranh nhất.",
             "NCC cài đặt Trạm trung chuyển dữ liệu trên máy chủ ảo nội bộ.",
             "Lập trình đường ống trích xuất dữ liệu tự động (ETL) và mã hóa kết nối VPN.",
             "Chạy kiểm thử luồng truyền dữ liệu từ Quảng Ngãi lên Workspace L3 trên TCT."
         ], Inches(4.8), SECONDARY),
        ("GIAI ĐOẠN 3: BÀN GIAO & KHAI THÁC\n(Tháng 08/2026 trở đi)",
         "BÀN GIAO, LÀM CHỦ & XÂY DASHBOARD",
         [
             "Nghiệm thu toàn bộ hệ thống Trục tích hợp dữ liệu đạt chuẩn TCT.",
             "NCC bàn giao toàn bộ mã nguồn và tài liệu kiến trúc.",
             "Đào tạo kỹ thuật cho đội ngũ CNTT Quảng Ngãi làm chủ vận hành 100%.",
             "Đưa vào khai thác các Dashboard báo cáo quản trị phục vụ Ban Giám đốc."
         ], Inches(8.8), ACCENT_GREEN)
    ]

    for title, sub, bullets, left_pos, col in plan_cards:
        tf7 = add_card(s7, left_pos, Inches(1.5), Inches(3.7), Inches(5.2), title, sub, top_accent_color=col)
        for b in bullets:
            p = tf7.add_paragraph()
            p.text = "• " + b if not b.startswith("*") else b
            p.font.size = Pt(10.5)
            p.font.italic = b.startswith("*")
            p.font.color.rgb = ACCENT_GOLD if b.startswith("*") else TEXT_DARK
            p.space_after = Pt(6)

    s7.notes_slide.notes_text_frame.text = (
        "KỊCH BẢN NÓI:\n"
        "Về kế hoạch thực hiện, chúng tôi chia thành 3 giai đoạn rõ ràng: "
        "Giai đoạn 1 trong tháng 3 và tháng 4, nội bộ tự làm, không tốn 1 đồng ngân sách nào để khảo sát và chuẩn bị đề bài; "
        "Sau đó mới tổ chức lựa chọn NCC để triển khai trong quý 2; "
        "và đến quý 3 là hoàn thành nghiệm thu, bàn giao toàn bộ cho anh em IT nội bộ quản lý. "
        "Chi tiết từng tuần công việc đã có trong tài liệu kế hoạch đính kèm theo báo cáo."
    )

    # ==========================================
    # SLIDE 8: KIẾN NGHỊ BAN GIÁM ĐỐC PHÊ DUYỆT
    # ==========================================
    s8 = prs.slides.add_slide(blank_layout)
    add_header(s8, "6. KIẾN NGHỊ & ĐỀ XUẤT BAN GIÁM ĐỐC PHÊ DUYỆT", "CÁC NỘI DUNG XIN CHỦ TRƯƠNG ĐỂ BẮT TAY TRIỂN KHAI NGAY")

    tf8 = add_card(s8, Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.2), "3 Đề xuất trọng tâm kính trình Ban Giám đốc phê duyệt", "CĂN CỨ THEO KẾ HOẠCH BÁO CÁO", top_accent_color=ACCENT_GREEN)

    proposals = [
        ("1. PHÊ DUYỆT CHỦ TRƯƠNG LỰA CHỌN PHƯƠNG ÁN 2:",
         "Chấp thuận chủ trương: Công ty chủ động phát triển Trục tích hợp dữ liệu riêng đồng bộ với Tổng công ty (có thuê Nhà cung cấp chuyên nghiệp tư vấn và triển khai)."),
        ("2. PHÊ DUYỆT THÀNH LẬP TỔ CÔNG TÁC DATA PLATFORM NỘI BỘ:",
         "Thành lập Tổ công tác gồm Bộ phận CNTT (chủ trì kỹ thuật) và các cán bộ phụ trách dữ liệu (Key users) từ Phòng Kế toán, Thương mại/Vật tư, HC-NS để bắt đầu ngay Giai đoạn 1 (Khảo sát và chuẩn hóa 29 danh mục Master Data)."),
        ("3. CHO PHÉP LÀM VIỆC VỚI CÁC NHÀ CUNG CẤP (NCC) ĐỂ KHẢO SÁT & BÁO GIÁ:",
         "Cho phép Tổ công tác tiếp xúc các NCC uy tín trên thị trường, truyền đạt đúng khung kiến trúc và yêu cầu kỹ thuật đã xây dựng để thu thập đề xuất kỹ thuật và báo giá cạnh tranh, hoàn thiện phương án tài chính trình Ban Giám đốc phê duyệt.")
    ]

    for title, detail in proposals:
        p = tf8.add_paragraph()
        p.text = title
        p.font.bold = True
        p.font.size = Pt(12.5)
        p.font.color.rgb = PRIMARY
        p.space_after = Pt(2)

        pd = tf8.add_paragraph()
        pd.text = detail
        pd.font.size = Pt(11)
        pd.font.color.rgb = TEXT_DARK
        pd.space_after = Pt(10)

    # Callout banner at bottom
    callout8 = s8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(5.75), Inches(10.9), Inches(0.68))
    callout8.fill.solid()
    callout8.fill.fore_color.rgb = RGBColor(220, 252, 231)
    callout8.line.color.rgb = ACCENT_GREEN
    callout8.line.width = Pt(1.5)
    tf_c8 = callout8.text_frame
    p_c8 = tf_c8.paragraphs[0]
    p_c8.text = "Kính trình Ban Giám đốc xem xét, thông qua chủ trương để kịp tiến độ kết nối năm 2026 của Tổng công ty!"
    p_c8.font.bold = True
    p_c8.font.size = Pt(12)
    p_c8.font.color.rgb = RGBColor(22, 101, 52)
    p_c8.alignment = PP_ALIGN.CENTER

    s8.notes_slide.notes_text_frame.text = (
        "KỊCH BẢN NÓI (KẾT THÚC BÁO CÁO):\n"
        "Kính thưa Ban Giám đốc, để đảm bảo tiến độ chuyển đổi số đúng chỉ đạo của TCT, bộ phận CNTT kính đề xuất Ban Giám đốc "
        "thông qua 3 chủ trương nêu trên. Việc phê duyệt sớm sẽ giúp chúng tôi hoàn thành Giai đoạn 1 nội bộ ngay trong tháng 4 "
        "và kịp triển khai đấu nối trong năm 2026. Chúng tôi xin trân trọng cảm ơn và lắng nghe ý kiến chỉ đạo của Ban Giám đốc!"
    )

    output_path = r"d:\My Profiles\DataPlatform\bao_cao_dataplatform_ptsc_qn_v2.pptx"
    prs.save(output_path)
    print(f"SUCCESS: Executive slide created at {output_path}")
    try:
        prs.save(r"d:\My Profiles\DataPlatform\bao_cao_dataplatform_ptsc_qn.pptx")
        print("Also updated v1.")
    except Exception as e:
        print("Note: Original v1 file was open in PowerPoint, saved to v2 instead.")

if __name__ == "__main__":
    create_executive_deck()
