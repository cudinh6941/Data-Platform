"""
BÁO CÁO DATA PLATFORM PTSC QUẢNG NGÃI - VERSION 3
Phong cách: Visual Storytelling, ít chữ, nhiều hình, kể chuyện từ từ
Bám sát Chỉ đạo số 9 của Ban Giám đốc
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE

# ─────────────────────── COLOR PALETTE ───────────────────────
NAVY       = RGBColor(10, 37, 74)
BLUE       = RGBColor(0, 95, 204)
BLUE_LIGHT = RGBColor(59, 130, 246)
GREEN      = RGBColor(16, 149, 96)
GREEN_SOFT = RGBColor(220, 252, 231)
RED        = RGBColor(209, 53, 53)
RED_SOFT   = RGBColor(254, 226, 226)
GOLD       = RGBColor(224, 138, 0)
GOLD_SOFT  = RGBColor(255, 243, 224)
WHITE      = RGBColor(255, 255, 255)
BLACK      = RGBColor(26, 38, 57)
GRAY       = RGBColor(100, 116, 139)
GRAY_LIGHT = RGBColor(241, 245, 249)
TEAL       = RGBColor(20, 184, 166)

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BL = prs.slide_layouts[6]  # blank

# ─────────────── HELPER FUNCTIONS ───────────────
def bg(slide, color):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    r.fill.solid(); r.fill.fore_color.rgb = color; r.line.fill.background()

def rect(slide, l, t, w, h, fill, line_c=None, line_w=None, radius=False):
    sh = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    r = slide.shapes.add_shape(sh, l, t, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = fill
    if line_c:
        r.line.color.rgb = line_c
        if line_w: r.line.width = line_w
    else:
        r.line.fill.background()
    return r

def txt(slide, l, t, w, h, text, size=14, bold=False, color=BLACK, align=PP_ALIGN.LEFT, font_name=None):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = color; p.alignment = align
    if font_name: p.font.name = font_name
    return tf

def header(slide, title, subtitle=""):
    rect(slide, 0, 0, W, Inches(1.1), NAVY)
    rect(slide, 0, Inches(1.1), W, Inches(0.04), GOLD)
    txt(slide, Inches(0.8), Inches(0.15), Inches(11.5), Inches(0.3), subtitle.upper(), 9.5, True, RGBColor(160, 195, 235))
    txt(slide, Inches(0.8), Inches(0.42), Inches(11.5), Inches(0.55), title, 20, True, WHITE)
    # footer
    txt(slide, Inches(0.8), Inches(7.12), Inches(11.7), Inches(0.3),
        "PTSC Quảng Ngãi | Báo cáo Data Platform — Chỉ đạo số 9 BGĐ", 8.5, color=GRAY)

def icon_box(slide, l, t, w, h, icon_text, label, sublabel="", fill_c=BLUE, icon_size=36, label_size=12):
    """A visual icon card: big emoji/icon on top, label below."""
    rect(slide, l, t, w, h, WHITE, RGBColor(218, 226, 237), Pt(1), radius=True)
    # icon
    txt(slide, l, t + Inches(0.15), w, Inches(0.6), icon_text, icon_size, align=PP_ALIGN.CENTER)
    # label
    tf = txt(slide, l + Inches(0.1), t + Inches(0.8), w - Inches(0.2), Inches(0.6), label, label_size, True, fill_c, PP_ALIGN.CENTER)
    if sublabel:
        p2 = tf.add_paragraph()
        p2.text = sublabel; p2.font.size = Pt(10); p2.font.color.rgb = GRAY; p2.alignment = PP_ALIGN.CENTER

def arrow_right(slide, l, t, w=Inches(0.6), h=Inches(0.35)):
    ar = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, w, h)
    ar.fill.solid(); ar.fill.fore_color.rgb = GOLD; ar.line.fill.background()

def arrow_down(slide, l, t, w=Inches(0.35), h=Inches(0.5)):
    ar = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, l, t, w, h)
    ar.fill.solid(); ar.fill.fore_color.rgb = GOLD; ar.line.fill.background()

def note(slide, text):
    slide.notes_slide.notes_text_frame.text = text

# ═══════════════════════════════════════════════════════════════
# SLIDE 1: TRANG BÌA
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, NAVY)
rect(s, Inches(0.8), Inches(1.5), Inches(0.15), Inches(4.5), GOLD)

txt(s, Inches(1.2), Inches(1.5), Inches(11), Inches(0.4),
    "BÁO CÁO BAN GIÁM ĐỐC  •  THỰC HIỆN CHỈ ĐẠO SỐ 9", 13, True, RGBColor(186, 215, 245))
txt(s, Inches(1.2), Inches(2.1), Inches(11), Inches(1.5),
    "DATA PLATFORM\n& TRỤC TÍCH HỢP DỮ LIỆU CÔNG TY – TCT", 32, True, WHITE)
txt(s, Inches(1.2), Inches(4.0), Inches(10), Inches(0.8),
    "So sánh Phương án  •  Kiến trúc Trục tích hợp  •  Dự toán Chi phí  •  Kế hoạch Triển khai", 14, color=RGBColor(220, 230, 242))
txt(s, Inches(1.2), Inches(5.5), Inches(10), Inches(0.4),
    "Đơn vị: Tổ Công tác CĐS & CNTT – PTSC Quảng Ngãi", 12.5, color=WHITE)
txt(s, Inches(1.2), Inches(5.95), Inches(10), Inches(0.4),
    "Kính trình: Ban Giám đốc Công ty", 12.5, True, RGBColor(253, 224, 71))

note(s, "Kính thưa Ban Giám đốc, hôm nay bộ phận CNTT xin báo cáo 4 nội dung trọng tâm "
     "theo đúng chỉ đạo số 9: So sánh phương án, kiến trúc trục tích hợp, dự toán chi phí, và kế hoạch triển khai.")

# ═══════════════════════════════════════════════════════════════
# SLIDE 2: VẤN ĐỀ — HIỆN TẠI CHÚNG TA ĐANG Ở ĐÂU?
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
header(s, "HIỆN TẠI: DỮ LIỆU ĐANG NẰM RẢI RÁC KHẮP NƠI", "BỐI CẢNH — VẤN ĐỀ CẦN GIẢI QUYẾT")

# 4 silo icons
silo_data = [
    ("💰", "Phần mềm\nKẾ TOÁN", "Dữ liệu tài chính"),
    ("👷", "Phần mềm\nNHÂN SỰ", "Dữ liệu nhân viên"),
    ("🔧", "Phần mềm\nVẬT TƯ", "Dữ liệu tồn kho"),
    ("📋", "Quản lý\nDỰ ÁN", "Dữ liệu hợp đồng"),
]
for i, (ic, lb, sub) in enumerate(silo_data):
    icon_box(s, Inches(0.8 + i * 2.9), Inches(1.6), Inches(2.5), Inches(1.6), ic, lb, sub, BLUE)

# Big red X connectors between silos
for i in range(3):
    txt(s, Inches(3.05 + i * 2.9), Inches(2.15), Inches(0.6), Inches(0.5), "✕", 24, True, RED, PP_ALIGN.CENTER)

# Pain point cards below
txt(s, Inches(0.8), Inches(3.6), Inches(11.7), Inches(0.4),
    "HẬU QUẢ NHÌN THẤY MỖI NGÀY:", 15, True, RED)

pain_data = [
    ("⏳", "Chậm trễ báo cáo", "Phòng ban mất 3-5 ngày\ntổng hợp số liệu thủ công"),
    ("⚠️", "Số liệu vênh nhau", "Kế toán báo 1 số,\nVật tư báo số khác"),
    ("📊", "Excel chồng Excel", "Copy-paste qua lại,\ndễ sai, khó kiểm soát"),
    ("🔒", "Không ai nhìn được\ntoàn cảnh", "BGĐ muốn xem tổng thể\nphải chờ rất lâu"),
]
for i, (ic, lb, sub) in enumerate(pain_data):
    x = Inches(0.8 + i * 3.05)
    icon_box(s, x, Inches(4.2), Inches(2.7), Inches(2.3), ic, lb, sub, RED, 30, 11.5)

note(s, "Thưa Ban Giám đốc, đây chính là bức tranh thực tế hàng ngày: 4 phần mềm hoạt động riêng rẽ, "
     "không ai nói chuyện được với ai. Hậu quả là BGĐ muốn xem số liệu tổng thể phải chờ các phòng ban "
     "ngồi gom file Excel, mất 3-5 ngày, và số liệu thường bị lệch giữa các phòng.")

# ═══════════════════════════════════════════════════════════════
# SLIDE 3: CÂU TRẢ LỜI — DATA PLATFORM LÀ GÌ? (ví von đơn giản)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
header(s, "DATA PLATFORM = \"HỆ THỐNG ĐƯỜNG ỐNG NƯỚC\" CHO DỮ LIỆU", "KHÁI NIỆM — GIẢI THÍCH ĐƠN GIẢN CHO LÃNH ĐẠO")

# Before vs After visual
rect(s, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.2), GRAY_LIGHT, RGBColor(220, 53, 53), Pt(2), True)
txt(s, Inches(0.8), Inches(1.55), Inches(5.5), Inches(0.4), "❌  TRƯỚC ĐÂY (Không có đường ống)", 14, True, RED, PP_ALIGN.CENTER)

# Before: messy icons
before_items = [
    (Inches(1.2), Inches(2.2), "💰", "Kế toán"),
    (Inches(3.4), Inches(2.2), "👷", "Nhân sự"),
    (Inches(5.0), Inches(2.2), "🔧", "Vật tư"),
    (Inches(2.3), Inches(4.3), "📋", "Dự án"),
]
for x, y, ic, lb in before_items:
    icon_box(s, x, y, Inches(1.4), Inches(1.2), ic, lb, "", RED, 24, 9)

# Messy lines text
txt(s, Inches(1.0), Inches(5.5), Inches(5.0), Inches(1.0),
    "Mỗi phòng tự \"xách xô\" (Excel)\nqua lại → Chậm, sai, mệt", 12, color=RED, align=PP_ALIGN.CENTER)

# Arrow between
arrow_right(s, Inches(6.5), Inches(3.6), Inches(0.7), Inches(0.4))

# After: clean flow
rect(s, Inches(7.4), Inches(1.5), Inches(5.1), Inches(5.2), GREEN_SOFT, GREEN, Pt(2), True)
txt(s, Inches(7.4), Inches(1.55), Inches(5.1), Inches(0.4), "✔  KHI CÓ DATA PLATFORM (Đường ống tự động)", 13, True, GREEN, PP_ALIGN.CENTER)

# After flow: Sources → Pipe → Dashboard
src_y = Inches(2.3)
for i, (ic, lb) in enumerate([("💰","Kế toán"), ("👷","Nhân sự"), ("🔧","Vật tư"), ("📋","Dự án")]):
    icon_box(s, Inches(7.7), src_y + Inches(i * 0.7), Inches(1.2), Inches(0.6), ic, lb, "", BLUE, 14, 8.5)

# Central pipe
rect(s, Inches(9.4), Inches(2.5), Inches(1.0), Inches(3.0), BLUE, radius=True)
txt(s, Inches(9.4), Inches(2.8), Inches(1.0), Inches(2.5),
    "🔄\n\nTự\nđộng\ngom\n&\nlọc\nsạch", 9, True, WHITE, PP_ALIGN.CENTER)

# arrows from sources to pipe
for i in range(4):
    arrow_right(s, Inches(8.95), src_y + Inches(i * 0.7) + Inches(0.12), Inches(0.4), Inches(0.25))

# arrow from pipe to dashboard
arrow_right(s, Inches(10.45), Inches(3.5), Inches(0.4), Inches(0.3))

# Dashboard
rect(s, Inches(11.0), Inches(2.6), Inches(1.3), Inches(2.2), WHITE, BLUE, Pt(1.5), True)
txt(s, Inches(11.0), Inches(2.7), Inches(1.3), Inches(2.0),
    "📊\n\nBan\nGiám\nđốc\nxem\nngay!", 9, True, BLUE, PP_ALIGN.CENTER)

note(s, "Nói đơn giản, Data Platform giống hệ thống đường ống nước tự động. "
     "Trước đây mỗi phòng ban là một cái giếng, muốn dùng nước phải xách xô sang nhau – "
     "đó chính là việc copy file Excel qua lại. Giờ Data Platform nối tất cả vào một bể lọc trung tâm, "
     "nước sạch tự chảy thẳng tới vòi của Ban Giám đốc mỗi ngày.")

# ═══════════════════════════════════════════════════════════════
# SLIDE 4: MÔ HÌNH HUB – SPOKE CỦA TCT (HÌNH VỊ THẾ)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
header(s, "VỊ THẾ CỦA QUẢNG NGÃI TRONG HỆ SINH THÁI TCT", "MÔ HÌNH — HUB & SPOKE (TRUNG TÂM – VỆ TINH)")

# Central Hub
rect(s, Inches(4.5), Inches(2.2), Inches(4.3), Inches(2.8), NAVY, radius=True)
txt(s, Inches(4.5), Inches(2.4), Inches(4.3), Inches(0.4), "🏢", 30, align=PP_ALIGN.CENTER)
txt(s, Inches(4.5), Inches(3.0), Inches(4.3), Inches(0.5), "TỔNG CÔNG TY (HUB)", 15, True, WHITE, PP_ALIGN.CENTER)
txt(s, Inches(4.5), Inches(3.5), Inches(4.3), Inches(1.2),
    "• Đầu tư máy chủ khổng lồ (Data Lakehouse)\n• Mua bản quyền Microsoft Fabric, Power BI\n• Ban hành 29 chuẩn Master Data", 10.5, color=RGBColor(186, 215, 245), align=PP_ALIGN.CENTER)

# Spoke: Quang Ngai (highlighted)
rect(s, Inches(0.5), Inches(3.0), Inches(3.5), Inches(1.6), GOLD_SOFT, GOLD, Pt(2.5), True)
txt(s, Inches(0.5), Inches(3.05), Inches(3.5), Inches(0.4), "🏭  PTSC QUẢNG NGÃI (L3)", 13, True, NAVY, PP_ALIGN.CENTER)
txt(s, Inches(0.5), Inches(3.5), Inches(3.5), Inches(1.0),
    "Chỉ cần dựng\n\"Trạm trung chuyển\" dữ liệu\nrồi bơm lên TCT", 11, color=BLACK, align=PP_ALIGN.CENTER)

# Arrow QN → Hub
arrow_right(s, Inches(4.05), Inches(3.5), Inches(0.4), Inches(0.3))

# Other spokes (faded)
spokes = [
    (Inches(9.5), Inches(3.0), "Đơn vị A"),
    (Inches(9.5), Inches(5.0), "Đơn vị B"),
    (Inches(5.0), Inches(5.5), "Đơn vị C"),
    (Inches(2.0), Inches(5.5), "Đơn vị D"),
]
for x, y, lb in spokes:
    rect(s, x, y, Inches(2.5), Inches(1.0), GRAY_LIGHT, RGBColor(200, 210, 220), Pt(1), True)
    txt(s, x, y + Inches(0.15), Inches(2.5), Inches(0.6), "🏭  " + lb, 10.5, color=GRAY, align=PP_ALIGN.CENTER)

# Key message at bottom
rect(s, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.5), GOLD_SOFT, GOLD, Pt(1.5))
txt(s, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.5),
    "💡 Quảng Ngãi KHÔNG phải tự mua máy chủ đắt tiền — chỉ cần xây \"đường dẫn kết nối\" lên TCT là được!",
    12, True, NAVY, PP_ALIGN.CENTER)

note(s, "Mô hình Hub-Spoke: TCT là trung tâm đã đầu tư hàng triệu đô xây hạ tầng. "
     "Quảng Ngãi là vệ tinh, chỉ cần xây đường ống nối lên là tận dụng được toàn bộ. "
     "Rất tiết kiệm — chúng ta không tốn tiền mua Server hay phần mềm phân tích.")

# ═══════════════════════════════════════════════════════════════
# SLIDE 5: ÁP LỰC — TIẾN ĐỘ TCT GIAO NĂM 2026
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
header(s, "TCT ĐÃ HOÀN THÀNH NỀN MÓNG — ĐẾN LƯỢT QUẢNG NGÃI!", "ÁP LỰC TIẾN ĐỘ — NGHỊ QUYẾT HĐQT")

# Timeline visual: 3 blocks horizontal
phases = [
    ("2024 – 2025", "ĐÃ HOÀN THÀNH ✔", "TCT xây xong\nData Lakehouse\n& ban hành\n29 Master Data", GREEN, GREEN_SOFT),
    ("2026 – 2027", "ĐANG THỰC HIỆN ⚡", "Quảng Ngãi\nphải kết nối!\nDựng trạm trung chuyển\n& bơm dữ liệu lên TCT", GOLD, GOLD_SOFT),
    ("2028 – 2030", "TƯƠNG LAI", "Khai thác toàn diện\nAI/ML phân tích\nbáo cáo thông minh", GRAY, GRAY_LIGHT),
]
for i, (period, status, desc, color, bg_c) in enumerate(phases):
    x = Inches(0.8 + i * 4.1)
    rect(s, x, Inches(1.6), Inches(3.7), Inches(0.7), color, radius=True)
    txt(s, x, Inches(1.65), Inches(3.7), Inches(0.35), period, 16, True, WHITE, PP_ALIGN.CENTER)
    txt(s, x, Inches(1.98), Inches(3.7), Inches(0.3), status, 10, True, WHITE, PP_ALIGN.CENTER)
    
    rect(s, x, Inches(2.4), Inches(3.7), Inches(2.8), bg_c, color, Pt(1.5), True)
    txt(s, x + Inches(0.15), Inches(2.6), Inches(3.4), Inches(2.4), desc, 13, color=BLACK, align=PP_ALIGN.CENTER)

    if i < 2:
        arrow_right(s, Inches(4.55 + i * 4.1), Inches(2.9), Inches(0.35), Inches(0.3))

# Callout: urgency
rect(s, Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.2), RED_SOFT, RED, Pt(2), True)
txt(s, Inches(1.2), Inches(5.7), Inches(11.0), Inches(0.35),
    "⏰  TẠI SAO PHẢI LÀM NGAY?", 15, True, RED)
tf = txt(s, Inches(1.2), Inches(6.1), Inches(11.0), Inches(0.6),
    "• Nghị quyết số 10/NQ-HĐQT-PTSC: Data Platform là trụ cột bắt buộc cho toàn bộ đơn vị thành viên", 11, color=BLACK)
p = tf.add_paragraph()
p.text = "• 2026-2027 là mốc deadline TCT giao — nếu trễ, đơn vị bị đánh giá chậm tiến độ CĐS"
p.font.size = Pt(11); p.font.color.rgb = BLACK

note(s, "TCT đã đi trước 2 năm rồi, xây xong sân ga trung tâm và mở sẵn làn đường cho mình. "
     "2026-2027 là deadline bắt buộc. Nếu chậm, đơn vị bị đánh giá KPI chuyển đổi số kém.")

# ═══════════════════════════════════════════════════════════════
# SLIDE 6: 3 NÚT THẮT TẠI QUẢNG NGÃI (Visual icons)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
header(s, "3 NÚT THẮT CẦN GIẢI QUYẾT TẠI QUẢNG NGÃI", "HIỆN TRẠNG — NHÌN THẲNG VÀO KHÓ KHĂN NỘI BỘ")

knots = [
    ("🔗", "DỮ LIỆU\nPHÂN MẢNH",
     "4 phần mềm hoạt động\nđộc lập, không kết nối.\nChuyển dữ liệu bằng\nExcel thủ công.",
     GOLD, GOLD_SOFT, Inches(0.8)),
    ("🔢", "LỆCH CHUẨN\nMASTER DATA",
     "Mã vật tư, nhân sự, dự án\ncủa QN chưa khớp\n29 chuẩn TCT ban hành.\nCần ánh xạ toàn bộ.",
     BLUE, RGBColor(219, 234, 254), Inches(4.8)),
    ("👨‍💻", "THIẾU KỸ SƯ DATA\nCHUYÊN SÂU",
     "IT giỏi hạ tầng mạng,\nnhưng chưa có kinh nghiệm\nData Engineering / ETL.\nTự mày mò = Trễ hạn TCT.",
     RED, RED_SOFT, Inches(8.8)),
]

for ic, title, desc, col, bg_c, x_pos in knots:
    rect(s, x_pos, Inches(1.5), Inches(3.7), Inches(5.0), bg_c, col, Pt(2), True)
    # Accent bar top
    rect(s, x_pos + Inches(0.05), Inches(1.52), Inches(3.6), Inches(0.08), col)
    txt(s, x_pos, Inches(1.8), Inches(3.7), Inches(0.6), ic, 40, align=PP_ALIGN.CENTER)
    txt(s, x_pos, Inches(2.6), Inches(3.7), Inches(0.7), title, 15, True, col, PP_ALIGN.CENTER)
    txt(s, x_pos + Inches(0.3), Inches(3.5), Inches(3.1), Inches(2.8), desc, 12.5, color=BLACK, align=PP_ALIGN.CENTER)

note(s, "3 nút thắt: 1 - phần mềm rời rạc không nói chuyện được; "
     "2 - mã danh mục chưa khớp chuẩn TCT nên đẩy lên sẽ bị lỗi; "
     "3 - IT mình chuyên mạng và hạ tầng, chưa từng làm Data Engineering, nếu tự mò mẫm sẽ trễ deadline.")

# ═══════════════════════════════════════════════════════════════
# SLIDE 7: SƠ ĐỒ LUỒNG DỮ LIỆU (TRẠM TRUNG CHUYỂN) — Diagram chính
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
header(s, "SƠ ĐỒ: DỮ LIỆU CHẢY TỪ QUẢNG NGÃI LÊN TCT NHƯ THẾ NÀO?", "KIẾN TRÚC — TRẠM TRUNG CHUYỂN DỮ LIỆU (ƯU TIÊN CAO)")

# Layer 1: Sources (left)
rect(s, Inches(0.4), Inches(1.4), Inches(2.6), Inches(5.4), RGBColor(219, 234, 254), BLUE, Pt(1.5), True)
txt(s, Inches(0.4), Inches(1.42), Inches(2.6), Inches(0.35), "PHẦN MỀM NỘI BỘ", 10, True, BLUE, PP_ALIGN.CENTER)

sources = [("💰", "Kế toán"), ("👷", "Nhân sự"), ("🔧", "Vật tư"), ("📋", "Dự án")]
for i, (ic, lb) in enumerate(sources):
    y = Inches(2.0 + i * 1.15)
    rect(s, Inches(0.7), y, Inches(2.0), Inches(0.85), WHITE, BLUE_LIGHT, Pt(1), True)
    txt(s, Inches(0.7), y + Inches(0.05), Inches(2.0), Inches(0.4), ic + "  " + lb, 11.5, True, NAVY, PP_ALIGN.CENTER)
    txt(s, Inches(0.7), y + Inches(0.45), Inches(2.0), Inches(0.35), "Cơ sở dữ liệu", 9, color=GRAY, align=PP_ALIGN.CENTER)
    # arrow from source
    arrow_right(s, Inches(2.75), y + Inches(0.25), Inches(0.4), Inches(0.25))

# Layer 2: Gateway (center)
rect(s, Inches(3.3), Inches(1.4), Inches(3.5), Inches(5.4), GOLD_SOFT, GOLD, Pt(2), True)
txt(s, Inches(3.3), Inches(1.42), Inches(3.5), Inches(0.35), "TRẠM TRUNG CHUYỂN (tại QN)", 10, True, GOLD, PP_ALIGN.CENTER)

gw_steps = [
    ("1️⃣  Hút dữ liệu", "Chế độ CHỈ ĐỌC\nKhông ảnh hưởng PM gốc"),
    ("2️⃣  Chuẩn hóa", "Ánh xạ theo\n29 Master Data TCT"),
    ("3️⃣  Kiểm tra", "Lọc sạch & kiểm\nchất lượng dữ liệu"),
    ("4️⃣  Mã hóa", "Đóng gói bảo mật\nAES-256"),
]
for i, (step_t, step_d) in enumerate(gw_steps):
    y = Inches(2.0 + i * 1.15)
    rect(s, Inches(3.6), y, Inches(2.9), Inches(0.85), WHITE, GOLD, Pt(1), True)
    txt(s, Inches(3.7), y + Inches(0.05), Inches(2.7), Inches(0.35), step_t, 11, True, NAVY)
    txt(s, Inches(3.7), y + Inches(0.4), Inches(2.7), Inches(0.4), step_d, 8.5, color=GRAY)

# Arrow from Gateway to VPN
arrow_right(s, Inches(6.9), Inches(3.5), Inches(0.5), Inches(0.3))

# Layer 3: VPN tunnel
rect(s, Inches(7.5), Inches(2.5), Inches(1.5), Inches(2.5), RGBColor(254, 243, 199), GOLD, Pt(1.5), True)
txt(s, Inches(7.5), Inches(2.6), Inches(1.5), Inches(2.3),
    "🔒\n\nĐường\nhầm\nVPN\nbảo mật", 10.5, True, NAVY, PP_ALIGN.CENTER)

# Arrow VPN → TCT
arrow_right(s, Inches(9.1), Inches(3.5), Inches(0.5), Inches(0.3))

# Layer 4: TCT destination (right)
rect(s, Inches(9.7), Inches(1.4), Inches(3.2), Inches(5.4), GREEN_SOFT, GREEN, Pt(1.5), True)
txt(s, Inches(9.7), Inches(1.42), Inches(3.2), Inches(0.35), "TỔNG CÔNG TY (CLOUD)", 10, True, GREEN, PP_ALIGN.CENTER)

# Workspace
rect(s, Inches(10.0), Inches(2.1), Inches(2.6), Inches(1.5), WHITE, GREEN, Pt(1), True)
txt(s, Inches(10.0), Inches(2.15), Inches(2.6), Inches(0.3), "☁️", 22, align=PP_ALIGN.CENTER)
txt(s, Inches(10.0), Inches(2.6), Inches(2.6), Inches(0.5), "Workspace L3\nQuảng Ngãi", 11, True, NAVY, PP_ALIGN.CENTER)
txt(s, Inches(10.0), Inches(3.1), Inches(2.6), Inches(0.4), "Microsoft Fabric", 9, color=GRAY, align=PP_ALIGN.CENTER)

arrow_down(s, Inches(11.15), Inches(3.65), Inches(0.3), Inches(0.35))

# Dashboard
rect(s, Inches(10.0), Inches(4.2), Inches(2.6), Inches(2.3), WHITE, BLUE, Pt(1.5), True)
txt(s, Inches(10.0), Inches(4.3), Inches(2.6), Inches(0.3), "📊", 24, align=PP_ALIGN.CENTER)
txt(s, Inches(10.0), Inches(4.8), Inches(2.6), Inches(0.5), "Dashboard\nPower BI", 12, True, BLUE, PP_ALIGN.CENTER)
txt(s, Inches(10.0), Inches(5.5), Inches(2.6), Inches(0.9),
    "BGĐ xem báo cáo\ntrực quan mỗi ngày\ntrên máy tính & điện thoại", 9.5, color=BLACK, align=PP_ALIGN.CENTER)

note(s, "Đây là sơ đồ ưu tiên cao theo chỉ đạo: Trục tích hợp hoạt động thế nào. "
     "Dữ liệu đi 4 bước: Hút → Chuẩn hóa → Kiểm tra → Mã hóa rồi bơm qua VPN sang Cloud TCT. "
     "Quan trọng: Chỉ đọc bản sao, không bao giờ đụng vào phần mềm gốc. An toàn tuyệt đối.")

# ═══════════════════════════════════════════════════════════════
# SLIDE 8: 3 CAM KẾT AN TOÀN KHI VẬN HÀNH
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
header(s, "3 CAM KẾT AN TOÀN KHI VẬN HÀNH TRỤC TÍCH HỢP", "AN TOÀN — BAN GIÁM ĐỐC YÊN TÂM 100%")

safe_cards = [
    ("🛡️", "KHÔNG SỢ\nSẬP PHẦN MỀM",
     "Trạm trung chuyển chạy\ntrên máy chủ ảo RIÊNG BIỆT.\n\nChỉ đọc bản sao dữ liệu,\nkhông can thiệp hệ thống gốc.\n\nChạy vào ban đêm / giờ vắng,\nkhông tranh tài nguyên.",
     BLUE, RGBColor(219, 234, 254), Inches(0.8)),
    ("🔐", "KHÔNG LỘ\nDỮ LIỆU RA NGOÀI",
     "Dữ liệu chỉ đi trong\nđường hầm VPN riêng\ncủa Tổng công ty.\n\nKhông mở cổng internet.\nKhông đi qua bên thứ 3.\n\nMã hóa AES-256\n(chuẩn ngân hàng).",
     GREEN, GREEN_SOFT, Inches(4.8)),
    ("🔔", "TỰ ĐỘNG\nCẢNH BÁO KHI LỖI",
     "Nếu nghẽn mạng hoặc\nsai lệch dữ liệu,\nhệ thống tự gửi email\ncảnh báo cho IT ngay.\n\nKhông cần người trực 24/7.\nMọi thứ được ghi log.",
     GOLD, GOLD_SOFT, Inches(8.8)),
]

for ic, title, desc, col, bg_c, x_pos in safe_cards:
    rect(s, x_pos, Inches(1.5), Inches(3.7), Inches(5.0), bg_c, col, Pt(2), True)
    rect(s, x_pos + Inches(0.05), Inches(1.52), Inches(3.6), Inches(0.08), col)
    txt(s, x_pos, Inches(1.8), Inches(3.7), Inches(0.6), ic, 42, align=PP_ALIGN.CENTER)
    txt(s, x_pos, Inches(2.6), Inches(3.7), Inches(0.7), title, 14, True, col, PP_ALIGN.CENTER)
    txt(s, x_pos + Inches(0.3), Inches(3.5), Inches(3.1), Inches(2.8), desc, 11.5, color=BLACK, align=PP_ALIGN.CENTER)

note(s, "Ban Giám đốc hay lo nhất 3 điều: sập phần mềm, lộ dữ liệu, có lỗi ai xử lý. "
     "Cả 3 đều được xử lý triệt để: chạy máy riêng, đi đường riêng, tự cảnh báo.")

# ═══════════════════════════════════════════════════════════════
# SLIDE 9: SO SÁNH 2 PHƯƠNG ÁN — THUÊ TCT VS PHÁT TRIỂN RIÊNG
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
header(s, "SO SÁNH: THUÊ TCT LÀM HỘ  VS  TỰ PHÁT TRIỂN RIÊNG (CÓ NCC)", "PHƯƠNG ÁN — CÂU HỎI TRỌNG TÂM CỦA BAN GIÁM ĐỐC")

# Two columns visual comparison
# Left: PA1 (red tint)
rect(s, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.2), RED_SOFT, RED, Pt(2), True)
rect(s, Inches(0.8), Inches(1.5), Inches(5.6), Inches(0.6), RED)
txt(s, Inches(0.8), Inches(1.5), Inches(5.6), Inches(0.6),
    "❌  PHƯƠNG ÁN 1: THUÊ DỊCH VỤ CỦA TCT", 13, True, WHITE, PP_ALIGN.CENTER)

pa1_items = [
    ("⛔  TCT chỉ xây Hub lõi chung", "Không đủ người xuống QN khảo sát\ntừng phần mềm kế toán, vật tư"),
    ("⏳  Xếp hàng chờ lượt", "TCT phải triển khai cho hàng chục\nđơn vị thành viên → Chắc chắn trễ 2026"),
    ("📉  Không có báo cáo riêng", "TCT chỉ xây Dashboard tài chính\nvĩ mô, không phục vụ bài toán\nquản trị xưởng/vật tư/nhân công của QN"),
    ("💸  Phải trả phí duy trì hàng năm", "Gánh phân bổ chi phí phần mềm\ntừ TCT, không kiểm soát được"),
]
for i, (t, d) in enumerate(pa1_items):
    y = Inches(2.3 + i * 1.05)
    txt(s, Inches(1.1), y, Inches(5.0), Inches(0.3), t, 11, True, RED)
    txt(s, Inches(1.1), y + Inches(0.28), Inches(5.0), Inches(0.7), d, 9.5, color=BLACK)

# Right: PA2 (green tint)
rect(s, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.2), GREEN_SOFT, GREEN, Pt(2), True)
rect(s, Inches(6.8), Inches(1.5), Inches(5.7), Inches(0.6), GREEN)
txt(s, Inches(6.8), Inches(1.5), Inches(5.7), Inches(0.6),
    "✔  PHƯƠNG ÁN 2: PHÁT TRIỂN RIÊNG + NCC  (ĐỀ XUẤT)", 12.5, True, WHITE, PP_ALIGN.CENTER)

pa2_items = [
    ("✅  NCC làm việc trực tiếp tại QN", "Đấu nối đúng từng phần mềm đặc thù\nmà công ty đang dùng hàng ngày"),
    ("⚡  Hoàn thành trong 3-4 tháng", "Đơn vị chủ động 100% tiến độ,\nđảm bảo kịp mốc hạn TCT giao"),
    ("📊  Có Dashboard riêng cho BGĐ", "Vừa đồng bộ chuẩn cho TCT,\nvừa có báo cáo tùy biến\nphục vụ điều hành tại đơn vị"),
    ("💰  Chi 1 lần, tận dụng Cloud TCT", "Tiết kiệm tiền tỷ nhờ dùng miễn phí\nhạ tầng Cloud mà TCT đã mua"),
]
for i, (t, d) in enumerate(pa2_items):
    y = Inches(2.3 + i * 1.05)
    txt(s, Inches(7.1), y, Inches(5.1), Inches(0.3), t, 11, True, GREEN)
    txt(s, Inches(7.1), y + Inches(0.28), Inches(5.1), Inches(0.7), d, 9.5, color=BLACK)

note(s, "Đây là câu hỏi trọng tâm nhất của Ban Giám đốc: Tại sao không nhờ TCT làm luôn? "
     "Câu trả lời rất rõ: TCT chỉ xây sân ga, không có người đi vào từng ngõ ngách phần mềm của QN. "
     "Nếu chờ TCT thì trễ hạn 2026. Phương án tối ưu là tự thuê NCC dựng trạm theo chuẩn TCT.")

# ═══════════════════════════════════════════════════════════════
# SLIDE 10: BAN GIÁM ĐỐC ĐƯỢC GÌ? (ĐẦU RA CỤ THỂ — DASHBOARD)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
header(s, "BAN GIÁM ĐỐC SẼ NHẬN ĐƯỢC GÌ KHI HOÀN THÀNH?", "KẾT QUẢ — ĐẦU RA NHÌN THẤY ĐƯỢC")

outputs = [
    ("📊", "DASHBOARD\nĐIỀU HÀNH",
     "Báo cáo trực quan\ncập nhật mỗi ngày.\nXem trên máy tính\nhoặc điện thoại.",
     BLUE, RGBColor(219, 234, 254)),
    ("🔄", "DỮ LIỆU\nTỰ ĐỘNG",
     "Hết phụ thuộc Excel.\nPhòng ban không cần\nxuất file thủ công.\nSố liệu tự chảy về.",
     TEAL, RGBColor(204, 251, 241)),
    ("✅", "MỘT NGUỒN\nSỰ THẬT",
     "Kế toán, Vật tư,\nNhân sự đều nhìn\ncùng 1 con số.\nHết tranh cãi lệch số.",
     GREEN, GREEN_SOFT),
    ("🏆", "ĐẠT CHỈ TIÊU\nCĐS TCT GIAO",
     "Hoàn thành đúng hạn\nnhiệm vụ kết nối\nData Platform theo\nNghị quyết HĐQT.",
     GOLD, GOLD_SOFT),
]

for i, (ic, title, desc, col, bg_c) in enumerate(outputs):
    x = Inches(0.8 + i * 3.05)
    rect(s, x, Inches(1.5), Inches(2.7), Inches(5.0), bg_c, col, Pt(2), True)
    rect(s, x + Inches(0.05), Inches(1.52), Inches(2.6), Inches(0.08), col)
    txt(s, x, Inches(1.8), Inches(2.7), Inches(0.6), ic, 38, align=PP_ALIGN.CENTER)
    txt(s, x, Inches(2.6), Inches(2.7), Inches(0.7), title, 13, True, col, PP_ALIGN.CENTER)
    txt(s, x + Inches(0.2), Inches(3.5), Inches(2.3), Inches(2.8), desc, 12, color=BLACK, align=PP_ALIGN.CENTER)

note(s, "Đầu ra rõ ràng: 1 - Ban Giám đốc có Dashboard xem tức thời trên điện thoại; "
     "2 - Hết phụ thuộc Excel thủ công; 3 - Một con số thống nhất toàn công ty; "
     "4 - Hoàn thành đúng chỉ tiêu KPI chuyển đổi số TCT giao.")

# ═══════════════════════════════════════════════════════════════
# SLIDE 11: DỰ TOÁN CHI PHÍ — NGẮN GỌN, RÕ RÀNG
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
header(s, "DỰ TOÁN CHI PHÍ SƠ BỘ (XÂY DỰNG TRƯỚC KHI GẶP NCC)", "NGÂN SÁCH — CHỦ ĐỘNG ĐÀM PHÁN, KHÔNG ĐỂ NCC DẪN DẮT")

# Budget breakdown as visual bars
items = [
    ("Hạ tầng máy chủ (VM)", "Tận dụng server hiện có", "0 đ", Inches(0.3), GREEN, "MIỄN PHÍ"),
    ("Khảo sát & Thiết kế kiến trúc", "NCC khảo sát CSDL + ánh xạ 29 Master Data", "80 – 120 Tr", Inches(1.8), BLUE, ""),
    ("Xây dựng Trục tích hợp & ETL", "Lập trình đường ống tự động hóa dữ liệu", "150 – 250 Tr", Inches(3.3), BLUE, "HẠNG MỤC LỚN NHẤT"),
    ("Dashboard Power BI", "3-5 bảng báo cáo cho Ban Giám đốc", "100 – 150 Tr", Inches(1.7), BLUE, ""),
    ("Đào tạo + Bảo hành 12 tháng", "Bàn giao mã nguồn, IT làm chủ 100%", "50 – 80 Tr", Inches(1.0), BLUE, ""),
]

y_start = Inches(1.6)
max_bar_w = Inches(5.5)
for i, (name, desc, cost, bar_w, col, badge) in enumerate(items):
    y = y_start + Inches(i * 0.95)
    # Label
    txt(s, Inches(0.8), y, Inches(3.5), Inches(0.3), name, 11.5, True, NAVY)
    txt(s, Inches(0.8), y + Inches(0.3), Inches(3.5), Inches(0.3), desc, 9, color=GRAY)
    # Bar
    if bar_w > 0:
        rect(s, Inches(4.5), y + Inches(0.05), bar_w, Inches(0.45), col if i > 0 else GREEN, radius=True)
    # Cost text
    txt(s, Inches(4.5) + bar_w + Inches(0.15), y + Inches(0.05), Inches(2.5), Inches(0.45), cost, 13, True, col if i > 0 else GREEN)
    # Badge
    if badge:
        bx = Inches(10.5)
        rect(s, bx, y + Inches(0.07), Inches(2.2), Inches(0.38), GOLD_SOFT if i > 0 else GREEN_SOFT, GOLD if i > 0 else GREEN, Pt(1), True)
        txt(s, bx, y + Inches(0.07), Inches(2.2), Inches(0.38), badge, 8.5, True, GOLD if i > 0 else GREEN, PP_ALIGN.CENTER)

# Total box
rect(s, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.9), NAVY, radius=True)
txt(s, Inches(1.2), Inches(6.05), Inches(6.0), Inches(0.4),
    "TỔNG DỰ TOÁN KHÁI TOÁN", 14, True, WHITE)
txt(s, Inches(7.5), Inches(6.0), Inches(4.8), Inches(0.45),
    "380 – 600 TRIỆU VNĐ", 22, True, RGBColor(253, 224, 71), PP_ALIGN.RIGHT)
txt(s, Inches(1.2), Inches(6.45), Inches(10.5), Inches(0.4),
    "Thanh toán phân kỳ theo mốc nghiệm thu  •  Tiết kiệm tiền tỷ nhờ tận dụng Cloud TCT đã mua", 10, color=RGBColor(186, 215, 245))

note(s, "Tổng chi phí khoảng 380 đến 600 triệu. Nhìn có vẻ lớn nhưng so với việc tự mua Cloud Server "
     "và bản quyền phần mềm phân tích thì tiết kiệm hàng tỷ đồng. "
     "Quan trọng: Giải ngân theo mốc nghiệm thu, không trả 1 cục.")

# ═══════════════════════════════════════════════════════════════
# SLIDE 12: LỘ TRÌNH 3 GIAI ĐOẠN — VISUAL TIMELINE
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
header(s, "LỘ TRÌNH TRIỂN KHAI 3 GIAI ĐOẠN", "KẾ HOẠCH — TỪNG BƯỚC RÕ RÀNG, KIỂM SOÁT CHẶT CHẼ")

timeline_data = [
    ("GIAI ĐOẠN 1", "Tháng 3 – 4 / 2026", "KHẢO SÁT NỘI BỘ",
     "NỘI BỘ TỰ LÀM  •  0 VNĐ",
     ["Thành lập Tổ công tác", "Khảo sát CSDL 4 phần mềm", "Ánh xạ 29 Master Data", "Lập đề bài cho NCC"],
     BLUE, RGBColor(219, 234, 254)),
    ("GIAI ĐOẠN 2", "Tháng 5 – 7 / 2026", "THUÊ NCC TRIỂN KHAI",
     "THUÊ NCC  •  GIAI ĐOẠN CHÍNH",
     ["Lựa chọn NCC tối ưu", "Dựng Trạm trung chuyển", "Viết đường ống ETL", "Kiểm thử kết nối TCT"],
     GOLD, GOLD_SOFT),
    ("GIAI ĐOẠN 3", "Tháng 8 / 2026 →", "BÀN GIAO & KHAI THÁC",
     "IT LÀM CHỦ 100%",
     ["Nghiệm thu hệ thống", "NCC bàn giao mã nguồn", "Đào tạo IT vận hành", "Xây Dashboard cho BGĐ"],
     GREEN, GREEN_SOFT),
]

for i, (phase, period, title, tag, items, col, bg_c) in enumerate(timeline_data):
    x = Inches(0.8 + i * 4.1)
    
    # Phase card
    rect(s, x, Inches(1.5), Inches(3.7), Inches(5.0), bg_c, col, Pt(2), True)
    rect(s, x + Inches(0.05), Inches(1.52), Inches(3.6), Inches(0.08), col)
    
    # Phase label
    txt(s, x, Inches(1.75), Inches(3.7), Inches(0.35), phase, 10, True, col, PP_ALIGN.CENTER)
    txt(s, x, Inches(2.05), Inches(3.7), Inches(0.35), period, 12, True, NAVY, PP_ALIGN.CENTER)
    
    # Title
    txt(s, x, Inches(2.55), Inches(3.7), Inches(0.4), title, 14, True, col, PP_ALIGN.CENTER)
    
    # Tag
    tag_w = Inches(3.0)
    rect(s, x + Inches(0.35), Inches(3.05), tag_w, Inches(0.35), col, radius=True)
    txt(s, x + Inches(0.35), Inches(3.05), tag_w, Inches(0.35), tag, 9, True, WHITE, PP_ALIGN.CENTER)
    
    # Checklist
    for j, item in enumerate(items):
        txt(s, x + Inches(0.3), Inches(3.6 + j * 0.5), Inches(3.1), Inches(0.4),
            "☐  " + item, 11, color=BLACK)
    
    # Arrow between phases
    if i < 2:
        arrow_right(s, Inches(4.55 + i * 4.1), Inches(3.5), Inches(0.35), Inches(0.3))

note(s, "Lộ trình chia 3 bước rõ ràng: Bước 1 nội bộ tự làm hoàn toàn miễn phí trong 2 tháng; "
     "Bước 2 mới bắt đầu giải ngân thuê NCC dựng trạm; Bước 3 nghiệm thu và đào tạo để IT làm chủ.")

# ═══════════════════════════════════════════════════════════════
# SLIDE 13: KIẾN NGHỊ — CALL TO ACTION (3 hộp lớn, ít chữ)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
header(s, "KÍNH ĐỀ NGHỊ BAN GIÁM ĐỐC PHÊ DUYỆT", "KIẾN NGHỊ — 3 NỘI DUNG XIN CHỦ TRƯƠNG")

proposals = [
    ("1️⃣", "CHỌN PHƯƠNG ÁN 2", "Phát triển riêng\nđồng bộ TCT\n(có thuê NCC tư vấn)", BLUE, RGBColor(219, 234, 254)),
    ("2️⃣", "THÀNH LẬP\nTỔ CÔNG TÁC", "IT chủ trì +\nKey Users các Phòng ban\nbắt đầu Giai đoạn 1 ngay", GOLD, GOLD_SOFT),
    ("3️⃣", "CHO PHÉP\nLIÊN HỆ NCC", "Khảo sát thị trường,\nlấy báo giá cạnh tranh\ntheo khung kiến trúc\nđã xây dựng", GREEN, GREEN_SOFT),
]

for i, (num, title, desc, col, bg_c) in enumerate(proposals):
    x = Inches(0.8 + i * 4.1)
    rect(s, x, Inches(1.5), Inches(3.7), Inches(4.2), bg_c, col, Pt(2.5), True)
    rect(s, x + Inches(0.05), Inches(1.52), Inches(3.6), Inches(0.08), col)
    
    txt(s, x, Inches(1.8), Inches(3.7), Inches(0.5), num, 30, align=PP_ALIGN.CENTER)
    txt(s, x, Inches(2.5), Inches(3.7), Inches(0.7), title, 16, True, col, PP_ALIGN.CENTER)
    txt(s, x + Inches(0.3), Inches(3.4), Inches(3.1), Inches(2.0), desc, 13, color=BLACK, align=PP_ALIGN.CENTER)

# Bottom callout
rect(s, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.8), GREEN_SOFT, GREEN, Pt(2), True)
txt(s, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.8),
    "Kính trình Ban Giám đốc xem xét, thông qua chủ trương\nđể đơn vị kịp tiến độ kết nối năm 2026 của Tổng công ty!",
    13, True, RGBColor(22, 101, 52), PP_ALIGN.CENTER)

note(s, "Kính thưa Ban Giám đốc, xin phép tóm gọn 3 đề xuất: "
     "Một là chọn phương án 2 phát triển riêng; Hai là thành lập tổ công tác để bắt đầu khảo sát ngay; "
     "Ba là cho phép chúng tôi liên hệ các NCC uy tín lấy báo giá trình Sếp duyệt trước khi ký kết. "
     "Rất mong Ban Giám đốc ủng hộ!")

# ═══════════════════════════════════════════════════════════════
# SLIDE 14: TRANG KẾT THÚC
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, NAVY)
rect(s, Inches(0.8), Inches(2.5), Inches(0.15), Inches(2.5), GOLD)

txt(s, Inches(1.2), Inches(2.5), Inches(11), Inches(1.0),
    "XIN TRÂN TRỌNG CẢM ƠN\nBAN GIÁM ĐỐC!", 30, True, WHITE, PP_ALIGN.LEFT)
txt(s, Inches(1.2), Inches(4.0), Inches(11), Inches(0.5),
    "Kính mời Ban Giám đốc cho ý kiến chỉ đạo.", 16, color=RGBColor(220, 230, 242))
txt(s, Inches(1.2), Inches(5.0), Inches(11), Inches(0.8),
    "Tài liệu đính kèm: Kế hoạch chi tiết Giai đoạn 1 & 2 & 3\nĐơn vị thực hiện: Tổ Công tác CĐS – Bộ phận CNTT PTSC Quảng Ngãi",
    12, color=RGBColor(186, 215, 245))

note(s, "Chúng tôi xin trân trọng cảm ơn và lắng nghe ý kiến chỉ đạo của Ban Giám đốc!")

# ═══════════════ SAVE ═══════════════
output = r"d:\My Profiles\DataPlatform\bao_cao_dataplatform_ptsc_qn_v3.pptx"
prs.save(output)
print(f"SUCCESS: {output}")
print(f"Total slides: {len(prs.slides)}")
