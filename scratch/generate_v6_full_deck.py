"""
BÁO CÁO DATA PLATFORM PTSC QUẢNG NGÃI - VERSION 6 (FULL DECK)
Cấu trúc 4 phần logic hoàn chỉnh theo chỉ đạo của Lãnh đạo:
- Phần 1 (Slide 1-4): Bản chất Data Platform của TCT (Tại sao làm, Nó thực ra là gì - 3 tầng, TCT đã đầu tư gì)
- Phần 2 (Slide 5-8): Vị thế & Chi tiết tại Quảng Ngãi (L3 Spoke, Hiện trạng 4 phần mềm, Chủ quyền Data Owner cấp 3, Bảo mật On-prem vs Hub)
- Phần 3 (Slide 9-11): Cách thức kết nối (Sơ đồ luồng trục tích hợp, So sánh 2 phương án, Cơ cấu chi phí PA3)
- Phần 4 (Slide 12-15): Kế hoạch hành động & Việc cần làm ở Giai đoạn tới (Quick-Win 4-6 tuần, Kế hoạch 3 GĐ, Cam kết nguồn lực & Bảo mật NCC, Kiến nghị)
"""
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Corporate Color Palette
NAVY       = RGBColor(12, 35, 68)       # #0C2344
NAVY_MID   = RGBColor(24, 59, 108)     # #183B6C
BLUE       = RGBColor(30, 100, 190)    # #1E64BE
BLUE_SOFT  = RGBColor(230, 240, 252)   # #E6F0FC
GREEN      = RGBColor(21, 128, 80)     # #158050
GREEN_SOFT = RGBColor(232, 248, 239)   # #E8F8EF
RED        = RGBColor(180, 40, 40)     # #B42828
RED_SOFT   = RGBColor(252, 235, 235)   # #FCEBEB
GOLD       = RGBColor(180, 115, 0)     # #B47300
GOLD_SOFT  = RGBColor(255, 245, 225)   # #FFF5E1
WHITE      = RGBColor(255, 255, 255)
BLACK      = RGBColor(28, 35, 48)
GRAY       = RGBColor(108, 118, 132)
GRAY_LIGHT = RGBColor(243, 245, 248)
GRAY_MED   = RGBColor(200, 208, 218)
TEAL       = RGBColor(13, 148, 136)
PURPLE     = RGBColor(109, 40, 217)

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BL = prs.slide_layouts[6]

def shape(slide, l, t, w, h, fill, line_c=None, line_w=None, rounded=False):
    sh_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    r = slide.shapes.add_shape(sh_type, l, t, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = fill
    if line_c:
        r.line.color.rgb = line_c
        if line_w: r.line.width = line_w
    else:
        r.line.fill.background()
    return r

def text(slide, l, t, w, h, content, size=14, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = content
    p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = color; p.alignment = align
    return tf

def header(slide, title, subtitle="", slide_num=""):
    shape(slide, 0, 0, W, Inches(1.05), NAVY)
    shape(slide, 0, Inches(1.05), W, Inches(0.035), GOLD)
    if subtitle:
        text(slide, Inches(0.8), Inches(0.1), Inches(10), Inches(0.25), subtitle.upper(), 9, True, RGBColor(155, 185, 220))
    text(slide, Inches(0.8), Inches(0.38), Inches(11), Inches(0.55), title, 18, True, WHITE)
    if slide_num:
        text(slide, Inches(12.0), Inches(0.38), Inches(0.8), Inches(0.5), slide_num, 12, True, RGBColor(155, 185, 220), PP_ALIGN.RIGHT)
    shape(slide, 0, Inches(7.2), W, Inches(0.008), GRAY_MED)
    text(slide, Inches(0.8), Inches(7.22), Inches(11.7), Inches(0.25),
         "PTSC Quảng Ngãi  |  Báo cáo Data Platform & Trục tích hợp  |  Thực hiện Chỉ đạo số 9 BGĐ", 8, color=GRAY)

def numbered_label(slide, l, t, num, color=BLUE):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, Inches(0.4), Inches(0.4))
    c.fill.solid(); c.fill.fore_color.rgb = color; c.line.fill.background()
    tf = c.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = str(num)
    p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER

def arrow_shape(slide, l, t, direction="right", w=Inches(0.5), h=Inches(0.3)):
    sh_type = MSO_SHAPE.RIGHT_ARROW if direction == "right" else MSO_SHAPE.DOWN_ARROW
    if direction == "down": w, h = h, w
    ar = slide.shapes.add_shape(sh_type, l, t, w, h)
    ar.fill.solid(); ar.fill.fore_color.rgb = GOLD; ar.line.fill.background()

def note(slide, txt_content):
    slide.notes_slide.notes_text_frame.text = txt_content

print("Building 15 slides...")

# =========================================================================
# SLIDE 1: COVER
# =========================================================================
s1 = prs.slides.add_slide(BL)
shape(s1, 0, 0, W, H, NAVY)
shape(s1, 0, Inches(1.4), W, Inches(0.04), GOLD)
shape(s1, 0, Inches(5.8), W, Inches(0.015), RGBColor(60, 90, 140))

text(s1, Inches(1.0), Inches(1.8), Inches(11), Inches(0.4),
     "BÁO CÁO BAN GIÁM ĐỐC  —  THỰC HIỆN CHỈ ĐẠO MỤC 9", 13, True, RGBColor(180, 200, 225))
text(s1, Inches(1.0), Inches(2.5), Inches(11), Inches(1.6),
     "PHƯƠNG ÁN TRIỂN KHAI DATA PLATFORM\nVÀ TRỤC TÍCH HỢP DỮ LIỆU CÔNG TY – TCT", 30, True, WHITE)
shape(s1, Inches(1.0), Inches(4.4), Inches(2.5), Inches(0.04), GOLD)

scope_items = [
    "Phần 1: Bản chất Data Platform của TCT — Tại sao xây dựng và Nó thực ra là gì?",
    "Phần 2: Vị thế & Hiện trạng tại Quảng Ngãi (L3 Spoke, 4 Ốc đảo phần mềm, Chủ quyền Data Owner cấp 3)",
    "Phần 3: Cách thức kết nối Trục tích hợp (MinIO + Fabric) và Cơ cấu chi phí tối ưu",
    "Phần 4: Kế hoạch Giai đoạn tới (Bài toán Quick-Win 4-6 tuần, Kế hoạch 3 Giai đoạn, Kiến nghị)"
]
for i, item in enumerate(scope_items):
    text(s1, Inches(1.0), Inches(4.65 + i * 0.28), Inches(11), Inches(0.28), "—  " + item, 11, color=RGBColor(200, 215, 235))

text(s1, Inches(1.0), Inches(6.05), Inches(11), Inches(0.3),
     "Đơn vị thực hiện:  Tổ Công tác CĐS & CNTT  –  PTSC Quảng Ngãi", 11.5, color=WHITE)
text(s1, Inches(1.0), Inches(6.4), Inches(11), Inches(0.3),
     "Kính trình:  Ban Giám đốc Công ty PTSC Quảng Ngãi", 11.5, True, RGBColor(253, 224, 71))

note(s1, "Kính thưa Ban Giám đốc, hôm nay bộ phận CNTT xin báo cáo toàn diện phương án triển khai Data Platform theo chỉ đạo số 9 của BGĐ. Bài báo cáo đi từ bản chất giải pháp của TCT, hiện trạng đơn vị, cách kết nối đến kế hoạch hành động cụ thể ở giai đoạn tới.")

# =========================================================================
# SLIDE 2: TẠI SAO TCT XÂY DỰNG DATA PLATFORM? (BỨC TRANH TOÀN TỔNG)
# =========================================================================
s2 = prs.slides.add_slide(BL)
header(s2, "BỨC TRANH TOÀN TỔNG: TẠI SAO TỔNG CÔNG TY XÂY DỰNG DATA PLATFORM?", "PHẦN 1: BẢN CHẤT GIẢI PHÁP TCT — NGUYÊN NHÂN RA ĐỜI", "01")

# Left box: Vấn đề của TCT
shape(s2, Inches(0.8), Inches(1.4), Inches(5.6), Inches(5.4), RED_SOFT, RED, Pt(1.5), rounded=True)
shape(s2, Inches(0.8), Inches(1.4), Inches(5.6), Inches(0.5), RED)
text(s2, Inches(0.8), Inches(1.42), Inches(5.6), Inches(0.45), "NGUYÊN NHÂN TCT PHẢI LÀM DATA PLATFORM", 12, True, WHITE, PP_ALIGN.CENTER)

tct_pains = [
    ("80% dữ liệu nằm ở Đơn vị thành viên", "Toàn bộ số liệu SXKD, nhân lực, công trình nằm ở các ĐVTV (Quảng Ngãi, M&C, POS...). Nếu không gom lại, TCT như 'người mù số liệu'."),
    ("Tồn tại hàng chục 'Ốc đảo phần mềm'", "Mỗi đơn vị dùng một phần mềm khác nhau (FAST, Bravo, phần mềm tự viết). Không có chuẩn chung để nói chuyện với nhau."),
    ("Báo cáo hợp nhất chậm và sai lệch", "TCT phải mất 2-3 tuần gửi công văn xin file Excel, ghép số thủ công để nộp Tập đoàn PVN. Số liệu vênh nhau liên tục."),
    ("Áp lực tuân thủ pháp lý mới", "Luật Dữ liệu, Luật BVDLCN và Petrovietnam ban hành Quy chế Quản trị Dữ liệu bắt buộc toàn tổng phải chuẩn hóa thống nhất.")
]
for i, (title, desc) in enumerate(tct_pains):
    y = Inches(2.1 + i * 1.15)
    shape(s2, Inches(1.1), y, Inches(5.0), Inches(1.0), WHITE, GRAY_MED, Pt(0.5), rounded=True)
    text(s2, Inches(1.3), y + Inches(0.08), Inches(4.6), Inches(0.3), f"• {title}", 11.5, True, RED)
    text(s2, Inches(1.3), y + Inches(0.35), Inches(4.6), Inches(0.6), desc, 10.5, color=BLACK)

arrow_shape(s2, Inches(6.6), Inches(3.8))

# Right box: Mục tiêu Data Platform TCT
shape(s2, Inches(7.3), Inches(1.4), Inches(5.2), Inches(5.4), GREEN_SOFT, GREEN, Pt(1.5), rounded=True)
shape(s2, Inches(7.3), Inches(1.4), Inches(5.2), Inches(0.5), GREEN)
text(s2, Inches(7.3), Inches(1.42), Inches(5.2), Inches(0.45), "MỤC TIÊU CHIẾN LƯỢC CỦA TCT", 12, True, WHITE, PP_ALIGN.CENTER)

tct_goals = [
    ("Tạo 'Hồ dữ liệu dùng chung' (Data Lakehouse)", "Gom toàn bộ dữ liệu của các đơn vị về một kho trung tâm an toàn."),
    ("Ban hành 29 Danh mục Master Data", "Quy định một chuẩn mã thống nhất toàn tổng: 1 mã nhân viên, 1 mã vật tư, 1 mã dự án duy nhất."),
    ("Tự động hóa báo cáo hợp nhất toàn Tổng", "Chuyển từ ghép file Excel thủ công sang Dashboard cập nhật tự động hàng ngày."),
    ("Nghị quyết 10/NQ-HĐQT bắt buộc thực hiện", "Mốc 2026-2027 là thời hạn bắt buộc các ĐVTV (như Quảng Ngãi) phải hoàn thành kết nối.")
]
for i, (title, desc) in enumerate(tct_goals):
    y = Inches(2.1 + i * 1.15)
    shape(s2, Inches(7.5), y, Inches(4.8), Inches(1.0), WHITE, GRAY_MED, Pt(0.5), rounded=True)
    text(s2, Inches(7.7), y + Inches(0.08), Inches(4.4), Inches(0.3), f"✔ {title}", 11.5, True, GREEN)
    text(s2, Inches(7.7), y + Inches(0.35), Inches(4.4), Inches(0.6), desc, 10.5, color=BLACK)

note(s2, "Kính thưa Ban Giám đốc: Tại sao TCT lại rầm rộ làm Data Platform? Vì 80% dữ liệu nằm ở các đơn vị như mình. TCT không thể điều hành nếu cứ chờ gửi file Excel thủ công. Do đó TCT xây Data Platform để gom dữ liệu toàn tổng về một mối.")

# =========================================================================
# SLIDE 3: DATA PLATFORM THỰC RA LÀ GÌ? (ĐỊNH NGHĨA & 3 TẦNG BẢN CHẤT)
# =========================================================================
s3 = prs.slides.add_slide(BL)
header(s3, "DATA PLATFORM THỰC RA LÀ GÌ? — HIỂU ĐÚNG BẢN CHẤT TRONG 10 GIÂY", "PHẦN 1: BẢN CHẤT GIẢI PHÁP TCT — ĐỊNH NGHĨA DÂN DÃ", "02")

# Top Banner: Định nghĩa phủ định (Giải tỏa tâm lý)
shape(s3, Inches(0.8), Inches(1.35), Inches(11.7), Inches(0.8), GOLD_SOFT, GOLD, Pt(1.5), rounded=True)
text(s3, Inches(1.1), Inches(1.42), Inches(11.1), Inches(0.65),
     "★ KHẲNG ĐỊNH CỐT LÕI: Data Platform KHÔNG PHẢI là phần mềm mới bắt nhân viên nhập liệu thêm!\n"
     "Nó là HỆ THỐNG ĐƯỜNG ỐNG NGẦM TỰ ĐỘNG: Hút bản sao dữ liệu ➔ Lọc sạch, đổi mã ➔ Bơm lên Dashboard cho Sếp.",
     11.5, True, NAVY, PP_ALIGN.CENTER)

# 3 Columns: 3 Tầng kiến trúc
layers = [
    ("TẦNG 1: THU THẬP DỮ LIỆU", "(Data Ingestion Pipeline)", BLUE, BLUE_SOFT,
     "Bản chất: ĐƯỜNG ỐNG HÚT TỰ ĐỘNG",
     ["• Nối vào các phần mềm hiện có (Kế toán, HRM, Vật tư, Dự án).",
      "• Hút bản sao dữ liệu (Read-Only), tuyệt đối không sửa dữ liệu gốc.",
      "• Chạy ngầm tự động theo lịch hẹn (đêm hoặc định kỳ hàng giờ).",
      "➔ Không làm gián đoạn công việc của CBNV."]),
    ("TẦNG 2: HỒ CHỨA & CHUẨN HÓA", "(Data Lakehouse & MDM)", TEAL, RGBColor(230, 248, 246),
     "Bản chất: NHÀ MÁY LỌC NƯỚC SẠCH",
     ["• Gom toàn bộ dữ liệu thô vào Hồ chứa tập trung (Lakehouse).",
      "• Khử trùng lặp, phát hiện dữ liệu rác, thiếu sót.",
      "• Tự động ĐỔI MÃ theo 29 danh mục Master Data của TCT.",
      "➔ Tạo ra 'Một nguồn sự thật duy nhất' tin cậy 100%."]),
    ("TẦNG 3: BÁO CÁO THÔNG MINH", "(Analytics & Power BI)", GREEN, GREEN_SOFT,
     "Bản chất: VÒI NƯỚC SẠCH CHO LÃNH ĐẠO",
     ["• Xuất ra màn hình Dashboard trực quan (Power BI).",
      "• Mở xem tức thì trên Laptop, iPad, Điện thoại di động.",
      "• Xem ngay Doanh thu, Chi phí, Dòng tiền, Nhân sự mỗi sáng.",
      "➔ Chấm dứt cảnh chờ đợi 3-5 ngày tổng hợp Excel."])
]

for i, (title, sub, col, bg_c, role, items) in enumerate(layers):
    x = Inches(0.8 + i * 4.05)
    shape(s3, x, Inches(2.35), Inches(3.7), Inches(4.6), bg_c, col, Pt(1.5), rounded=True)
    shape(s3, x, Inches(2.35), Inches(3.7), Inches(0.7), col)
    text(s3, x, Inches(2.38), Inches(3.7), Inches(0.3), title, 11.5, True, WHITE, PP_ALIGN.CENTER)
    text(s3, x, Inches(2.68), Inches(3.7), Inches(0.25), sub, 9.5, False, RGBColor(220, 235, 255), PP_ALIGN.CENTER)
    
    # Sub-role banner
    shape(s3, x + Inches(0.2), Inches(3.15), Inches(3.3), Inches(0.4), WHITE, col, Pt(1.0), rounded=True)
    text(s3, x + Inches(0.2), Inches(3.2), Inches(3.3), Inches(0.3), role, 10, True, col, PP_ALIGN.CENTER)
    
    # Items
    for j, item in enumerate(items):
        text(s3, x + Inches(0.25), Inches(3.7 + j * 0.75), Inches(3.2), Inches(0.7), item, 10.5, color=BLACK)

note(s3, "Giải thích cho Lãnh đạo hiểu: Data Platform không phải phần mềm mới bắt nhân viên gõ máy tính. Nó gồm 3 tầng: Ống hút tự động từ phần mềm cũ -> Trạm lọc sạch và đổi mã TCT -> Vòi nước sạch là màn hình Dashboard cho Ban Giám đốc xem mỗi sáng.")

# =========================================================================
# SLIDE 4: TCT ĐÃ ĐẦU TƯ NHỮNG GÌ? (NỀN TẢNG CỐT LÕI ĐÃ CÓ)
# =========================================================================
s4 = prs.slides.add_slide(BL)
header(s4, "TỔNG CÔNG TY ĐÃ ĐẦU TƯ NHỮNG GÌ? — NỀN TẢNG HYBRID ĐÃ HOÀN THÀNH", "PHẦN 1: BẢN CHẤT GIẢI PHÁP TCT — TÀI NGUYÊN SẴN CÓ", "03")

# Left Column: Cloud (Microsoft Fabric)
shape(s4, Inches(0.8), Inches(1.4), Inches(5.6), Inches(4.3), BLUE_SOFT, BLUE, Pt(1.5), rounded=True)
shape(s4, Inches(0.8), Inches(1.4), Inches(5.6), Inches(0.55), BLUE)
text(s4, Inches(0.8), Inches(1.45), Inches(5.6), Inches(0.4), "HẠ TẦNG CLOUD — TCT ĐÃ MUA BẢN QUYỀN", 12, True, WHITE, PP_ALIGN.CENTER)

cloud_items = [
    ("Microsoft Fabric OneLake (Khởi điểm 20 TB)", "Kho lưu trữ khổng lồ trên đám mây Azure, được tối ưu cho phân tích dữ liệu lớn và AI/ML."),
    ("Cấp sẵn Workspace riêng cho Quảng Ngãi", "Mỗi đơn vị có một phân vùng độc lập (Tenant/Workspace L3). Dữ liệu QN được cô lập bảo mật 100%."),
    ("Bản quyền Power BI Enterprise", "Đã bao gồm sẵn license cho lãnh đạo và cán bộ khai thác báo cáo thông minh."),
    ("Microsoft Purview", "Hệ thống quản trị danh mục dữ liệu, phân loại nhãn bảo mật tự động.")
]
for i, (title, desc) in enumerate(cloud_items):
    y = Inches(2.1 + i * 0.88)
    shape(s4, Inches(1.0), y, Inches(5.2), Inches(0.78), WHITE, GRAY_MED, Pt(0.5), rounded=True)
    text(s4, Inches(1.15), y + Inches(0.05), Inches(4.9), Inches(0.25), f"✔ {title}", 11, True, BLUE)
    text(s4, Inches(1.15), y + Inches(0.28), Inches(4.9), Inches(0.45), desc, 9.5, color=BLACK)

# Right Column: On-Premise (Datacenter TCT)
shape(s4, Inches(6.9), Inches(1.4), Inches(5.6), Inches(4.3), GRAY_LIGHT, NAVY, Pt(1.5), rounded=True)
shape(s4, Inches(6.9), Inches(1.4), Inches(5.6), Inches(0.55), NAVY)
text(s4, Inches(6.9), Inches(1.45), Inches(5.6), Inches(0.4), "HẠ TẦNG ON-PREMISE — TẠI DATACENTER TCT", 12, True, WHITE, PP_ALIGN.CENTER)

onprem_items = [
    ("Hồ dữ liệu MinIO Lakehouse", "Cụm lưu trữ on-premise an toàn đặt tại Tầng 10 CQTCT, lưu trữ dữ liệu gốc theo Luật Dữ liệu VN."),
    ("Trục tích hợp doanh nghiệp (ESB)", "Cầu nối truyền nhận thông điệp thời gian thực giữa các hệ thống phần mềm."),
    ("Hệ thống MDM (Master Data Management)", "Bộ động cơ chuẩn hóa 'Hồ sơ vàng' (Golden Record) dựa trên 29 danh mục chuẩn."),
    ("Trung tâm An ninh thông tin SIEM & SOC", "Giám sát an toàn thông tin 24/7 theo tiêu chuẩn bảo mật 8 Zones của TCT.")
]
for i, (title, desc) in enumerate(onprem_items):
    y = Inches(2.1 + i * 0.88)
    shape(s4, Inches(7.1), y, Inches(5.2), Inches(0.78), WHITE, GRAY_MED, Pt(0.5), rounded=True)
    text(s4, Inches(7.25), y + Inches(0.05), Inches(4.9), Inches(0.25), f"✔ {title}", 11, True, NAVY)
    text(s4, Inches(7.25), y + Inches(0.28), Inches(4.9), Inches(0.45), desc, 9.5, color=BLACK)

# Bottom Banner: Ý nghĩa với Quảng Ngãi
shape(s4, Inches(0.8), Inches(5.9), Inches(11.7), Inches(1.0), GOLD_SOFT, GOLD, Pt(1.5), rounded=True)
text(s4, Inches(1.1), Inches(5.98), Inches(11.1), Inches(0.85),
     "👉 KẾT LUẬN QUAN TRỌNG VỚI QUẢNG NGÃI:\n"
     "TCT đã bỏ ra hàng chục tỷ đầu tư 'Đường cao tốc' và mua sẵn bản quyền. Quảng Ngãi KHÔNG PHẢI MUA MÁY CHỦ ĐẮT TIỀN!\n"
     "Việc duy nhất của Quảng Ngãi là: Xây dựng 'Đường nhánh' (Trạm trung chuyển) để đấu nối vào hạ tầng TCT.",
     11.5, True, NAVY, PP_ALIGN.CENTER)

note(s4, "TCT đã xây xong nền tảng Hybrid hoàn chỉnh: Cloud có Fabric 20TB, On-prem có MinIO và MDM. TCT đã mở sẵn làn đường cho Quảng Ngãi. Chúng ta chỉ cần làm đường nhánh đấu nối vào, tiết kiệm hàng tỷ đồng.")

# =========================================================================
# SLIDE 5: VỊ THẾ QUẢNG NGÃI (LEVEL 3 SPOKE)
# =========================================================================
s5 = prs.slides.add_slide(BL)
header(s5, "VỊ THẾ QUẢNG NGÃI: LEVEL 3 TRONG MÔ HÌNH HUB-SPOKE CỦA TCT", "PHẦN 2: CHI TIẾT TẠI QUẢNG NGÃI — PHÂN LOẠI ĐƠN VỊ", "04")

levels = [
    ("LEVEL 1", "Chi nhánh", "Vận hành như 1 ban TCT. Dùng chung toàn bộ hệ thống.", GRAY, GRAY_LIGHT, False),
    ("LEVEL 2", "Đơn vị nhỏ", "Dùng chung Tenant trên Hub. Kéo dữ liệu qua API nhẹ.", GRAY, GRAY_LIGHT, False),
    ("LEVEL 3", "ĐƠN VỊ LỚN\n(PTSC QUẢNG NGÃI)", "• Cấp Workspace L3 riêng biệt trên Hub TCT.\n• KHÔNG CẦN mua máy chủ dHCI tiền tỷ.\n• Tận dụng VM sẵn có làm trạm trung chuyển.\n• Độc lập xây Dashboard quản trị riêng.", BLUE, BLUE_SOFT, True),
    ("LEVEL 4", "Đặc biệt lớn\n(PTSC M&C)", "Tự đầu tư cụm máy chủ lưu trữ dHCI riêng tại chỗ.", GRAY, GRAY_LIGHT, False),
]

for i, (lvl, name, desc, col, bg_c, highlight) in enumerate(levels):
    x = Inches(0.8 + i * 2.95)
    border_w = Pt(2.5) if highlight else Pt(1.0)
    shape(s5, x, Inches(1.5), Inches(2.75), Inches(4.3), bg_c, col, border_w, rounded=True)
    shape(s5, x, Inches(1.5), Inches(2.75), Inches(0.6), col)
    text(s5, x, Inches(1.55), Inches(2.75), Inches(0.5), lvl, 13, True, WHITE, PP_ALIGN.CENTER)
    text(s5, x + Inches(0.15), Inches(2.2), Inches(2.45), Inches(0.6), name, 12, True, col, PP_ALIGN.CENTER)
    shape(s5, x + Inches(0.2), Inches(2.9), Inches(2.35), Inches(0.01), GRAY_MED)
    text(s5, x + Inches(0.15), Inches(3.05), Inches(2.45), Inches(2.5), desc, 10.5, color=BLACK)

# Bottom conclusion
shape(s5, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.9), GREEN_SOFT, GREEN, Pt(1.5), rounded=True)
text(s5, Inches(1.1), Inches(6.1), Inches(11.1), Inches(0.7),
     "✔ TIẾT KIỆM TỐI ĐA NGÂN SÁCH ĐẦU TƯ:\n"
     "Là Level 3 Spoke, PTSC Quảng Ngãi vừa có không gian lưu trữ riêng độc lập, vừa không phải gánh chi phí đầu tư phần cứng.",
     11.5, True, GREEN, PP_ALIGN.CENTER)

note(s5, "Theo quy hoạch của TCT, Quảng Ngãi là Level 3 Spoke. Chúng ta có phân vùng riêng hoàn toàn, nhưng không phải mua máy chủ dHCI đắt tiền như cấp L4.")

# =========================================================================
# SLIDE 6: HIỆN TRẠNG DỮ LIỆU TẠI QUẢNG NGÃI (3 NÚT THẮT)
# =========================================================================
s6 = prs.slides.add_slide(BL)
header(s6, "HIỆN TRẠNG DỮ LIỆU TẠI QUẢNG NGÃI: 3 NÚT THẮT CẦN THÁO GỠ", "PHẦN 2: CHI TIẾT TẠI QUẢNG NGÃI — NGUYÊN NHÂN CHẬM TRỄ", "05")

# Top 4 Silo boxes
silos = ["Phần mềm\nKẾ TOÁN", "Phần mềm\nNHÂN SỰ", "Phần mềm\nVẬT TƯ", "Quản lý\nDỰ ÁN"]
for i, name in enumerate(silos):
    x = Inches(0.8 + i * 2.95)
    shape(s6, x, Inches(1.4), Inches(2.55), Inches(1.1), WHITE, GRAY_MED, Pt(0.8), rounded=True)
    shape(s6, x + Inches(0.04), Inches(1.42), Inches(0.08), Inches(1.06), BLUE)
    text(s6, x + Inches(0.25), Inches(1.5), Inches(2.1), Inches(0.8), name, 12, True, NAVY, PP_ALIGN.CENTER)
    if i < 3:
        text(s6, Inches(3.1 + i * 2.95), Inches(1.65), Inches(0.5), Inches(0.4), "✕", 18, True, RED, PP_ALIGN.CENTER)

text(s6, Inches(0.8), Inches(2.7), Inches(11.7), Inches(0.35),
     "3 Nút thắt lớn cản trở công tác quản trị điều hành của Ban Giám đốc:", 13, True, RED)

knots = [
    ("DỮ LIỆU PHÂN MẢNH (SILO)", RED, RED_SOFT,
     "4 phần mềm hoạt động riêng rẽ.\nChuyển số liệu bằng file Excel thủ công.\nMất 3–5 ngày tổng hợp báo cáo tuần/tháng.\nSố liệu dễ lệch nhau giữa các phòng."),
    ("LỆCH 29 CHUẨN MASTER DATA", GOLD, GOLD_SOFT,
     "Mã nhân sự, vật tư, đối tác của QN\nchưa khớp 29 danh mục chuẩn của TCT.\nNếu đẩy trực tiếp lên Hub sẽ bị lỗi\nvà hệ thống TCT từ chối tiếp nhận."),
    ("THIẾU KỸ SƯ DATA PIPELINE", BLUE, BLUE_SOFT,
     "IT nội bộ giỏi mạng, máy chủ và hỗ trợ,\nnhưng chưa có kinh nghiệm viết đường ống (ETL).\nTự mày mò từ đầu sẽ mất 8–12 tháng\nvà nguy cơ trễ mốc 2026 của TCT."),
]
for i, (title, col, bg_c, desc) in enumerate(knots):
    x = Inches(0.8 + i * 4.05)
    shape(s6, x, Inches(3.15), Inches(3.7), Inches(3.8), bg_c, col, Pt(1.5), rounded=True)
    shape(s6, x + Inches(0.04), Inches(3.17), Inches(0.08), Inches(3.74), col)
    numbered_label(s6, x + Inches(0.3), Inches(3.35), i + 1, col)
    text(s6, x + Inches(0.85), Inches(3.35), Inches(2.6), Inches(0.4), title, 12.5, True, col)
    shape(s6, x + Inches(0.3), Inches(3.9), Inches(3.1), Inches(0.01), GRAY_MED)
    text(s6, x + Inches(0.3), Inches(4.1), Inches(3.1), Inches(2.6), desc, 11, color=BLACK)

note(s6, "Nhìn thẳng vào thực tế Quảng Ngãi: 4 phần mềm rời rạc, mã chưa khớp chuẩn TCT, và anh em IT nội bộ chưa từng làm Data Pipeline. Cần phương án chuyên nghiệp để vượt qua 3 nút thắt này.")

# =========================================================================
# SLIDE 7: QUY CHẾ QUẢN TRỊ DỮ LIỆU & CHỦ QUYỀN ĐƠN VỊ
# =========================================================================
s7 = prs.slides.add_slide(BL)
header(s7, "QUY CHẾ QUẢN TRỊ DỮ LIỆU TCT: CHỦ QUYỀN CỦA QUẢNG NGÃI", "PHẦN 2: CHI TIẾT TẠI QUẢNG NGÃI — QUYỀN TỰ QUYẾT", "06")

gov_levels = [
    ("CẤP 1 & 2: HỘI ĐỒNG DỮ LIỆU TCT", "Ban hành chính sách khung, chiến lược dữ liệu toàn tổng.", GRAY, GRAY_LIGHT, False),
    ("CẤP 3: CHỦ QUẢN DỮ LIỆU (DATA OWNER)\nLÃNH ĐẠO ĐƠN VỊ THÀNH VIÊN (PTSC QUẢNG NGÃI)",
     "• NẮM QUYỀN SỞ HỮU NGHIỆP VỤ (BUSINESS OWNERSHIP) DUY NHẤT.\n"
     "• TOÀN QUYỀN PHÊ DUYỆT: Dữ liệu nào được phép chia sẻ lên TCT, dữ liệu nào BẮT BUỘC GIỮ LẠI NỘI BỘ.\n"
     "• Chịu trách nhiệm giải trình về tính chính xác và an toàn dữ liệu của đơn vị.",
     GREEN, GREEN_SOFT, True),
    ("CẤP 4: QUẢN TRỊ MIỀN DỮ LIỆU (DATA STEWARDS)", "Key users các phòng ban: Rà soát chất lượng và ánh xạ mã kỹ thuật.", GRAY, GRAY_LIGHT, False),
    ("CẤP 5: BAN NCPT & CĐS TỔNG CÔNG TY",
     "ĐƠN VỊ VẬN HÀNH KỸ THUẬT NỀN TẢNG — TUYỆT ĐỐI KHÔNG SỞ HỮU DỮ LIỆU NGHIỆP VỤ CỦA ĐƠN VỊ!", RED, RED_SOFT, False),
]

y_pos = Inches(1.4)
for title, desc, col, bg_c, highlight in gov_levels:
    h = Inches(2.1) if highlight else Inches(1.05)
    border_w = Pt(2.5) if highlight else Pt(1.0)
    shape(s7, Inches(0.8), y_pos, Inches(11.7), h, bg_c, col, border_w, rounded=True)
    shape(s7, Inches(0.84), y_pos + Inches(0.04), Inches(0.1), h - Inches(0.08), col)
    
    if highlight:
        text(s7, Inches(1.2), y_pos + Inches(0.12), Inches(11.0), Inches(0.55), title, 12.5, True, col)
        text(s7, Inches(1.2), y_pos + Inches(0.7), Inches(11.0), Inches(1.3), desc, 11, color=BLACK)
    else:
        text(s7, Inches(1.2), y_pos + Inches(0.1), Inches(11.0), Inches(0.3), title, 11.5, True, col)
        text(s7, Inches(1.2), y_pos + Inches(0.42), Inches(11.0), Inches(0.55), desc, 10.5, color=BLACK)
    y_pos += h + Inches(0.15)

note(s7, "Theo Quy chế Quản trị Dữ liệu chính thức của PTSC, Lãnh đạo Đơn vị là Data Owner Cấp 3. Sếp có quyền tối cao quyết định dữ liệu nào đi, dữ liệu nào ở lại. Ban CNTT TCT ở Cấp 5 chỉ là thợ kỹ thuật vận hành máy.")

# =========================================================================
# SLIDE 8: NGUYÊN TẮC BẢO MẬT: DỮ LIỆU NÀO Ở LẠI, NÀO ĐI?
# =========================================================================
s8 = prs.slides.add_slide(BL)
header(s8, "NGUYÊN TẮC BẢO MẬT: DỮ LIỆU NÀO Ở LẠI, DỮ LIỆU NÀO CHUYỂN ĐI?", "PHẦN 2: CHI TIẾT TẠI QUẢNG NGÃI — AN TOÀN THÔNG TIN", "07")

# Left Column: Ở lại Quảng Ngãi
shape(s8, Inches(0.8), Inches(1.4), Inches(5.6), Inches(5.4), RED_SOFT, RED, Pt(1.5), rounded=True)
shape(s8, Inches(0.8), Inches(1.4), Inches(5.6), Inches(0.55), RED)
text(s8, Inches(0.8), Inches(1.45), Inches(5.6), Inches(0.4), "LƯU GIỮ 100% TẠI NỘI BỘ QUẢNG NGÃI (ON-PREMISE)", 11.5, True, WHITE, PP_ALIGN.CENTER)

stay_items = [
    ("Bí mật kinh doanh & Giá thầu", "Định mức đơn giá thầu, chi phí riêng từng dự án, biên lợi nhuận, chiến lược thương mại."),
    ("Dữ liệu vận hành chi tiết", "Nhật trình thi công, nhật ký xưởng, thông số bảo dưỡng thiết bị chi tiết hàng ngày."),
    ("Dữ liệu cá nhân chưa xử lý", "Bảng lương chi tiết, thông tin tài chính cá nhân, hồ sơ sức khỏe CBNV theo Luật BVDLCN."),
    ("Toàn bộ CSDL sản xuất gốc", "Cơ sở dữ liệu đang vận hành đặt tại phòng máy chủ của Quảng Ngãi, TCT không truy cập trực tiếp.")
]
for i, (title, desc) in enumerate(stay_items):
    y = Inches(2.15 + i * 1.15)
    shape(s8, Inches(1.0), y, Inches(5.2), Inches(1.0), WHITE, GRAY_MED, Pt(0.5), rounded=True)
    text(s8, Inches(1.2), y + Inches(0.08), Inches(4.8), Inches(0.3), f"🔒 {title}", 11.5, True, RED)
    text(s8, Inches(1.2), y + Inches(0.35), Inches(4.8), Inches(0.6), desc, 10.5, color=BLACK)

# Right Column: Chuyển về TCT
shape(s8, Inches(6.9), Inches(1.4), Inches(5.6), Inches(5.4), GREEN_SOFT, GREEN, Pt(1.5), rounded=True)
shape(s8, Inches(6.9), Inches(1.4), Inches(5.6), Inches(0.55), GREEN)
text(s8, Inches(6.9), Inches(1.45), Inches(5.6), Inches(0.4), "ĐƯỢC PHÉP ĐỒNG BỘ LÊN TCT HUB", 11.5, True, WHITE, PP_ALIGN.CENTER)

send_items = [
    ("29 Danh mục Master Data", "Mã phòng ban, mã chức danh, mã vật tư chuẩn... đã được quy đổi thống nhất."),
    ("Báo cáo số liệu tổng hợp", "Số liệu Doanh thu tổng, Sản lượng tổng, Số lượng nhân sự hợp nhất phục vụ báo cáo toàn tổng."),
    ("Dữ liệu đã che mờ (Data Masking)", "Số CCCD, SĐT cá nhân được mã hóa/làm mờ tự động trước khi truyền qua mạng theo NĐ 13."),
    ("Tenant L3 cách ly độc lập", "Dữ liệu chuyển vào Workspace riêng của QN, các đơn vị thành viên khác không thể xem được.")
]
for i, (title, desc) in enumerate(send_items):
    y = Inches(2.15 + i * 1.15)
    shape(s8, Inches(7.1), y, Inches(5.2), Inches(1.0), WHITE, GRAY_MED, Pt(0.5), rounded=True)
    text(s8, Inches(7.3), y + Inches(0.08), Inches(4.8), Inches(0.3), f"✔ {title}", 11.5, True, GREEN)
    text(s8, Inches(7.3), y + Inches(0.35), Inches(4.8), Inches(0.6), desc, 10.5, color=BLACK)

note(s8, "Nguyên tắc bảo mật rất rõ ràng: Dữ liệu chi tiết nhạy cảm và bí mật kinh doanh lưu giữ tại Quảng Ngãi. Chỉ dữ liệu tổng hợp và danh mục dùng chung mới đồng bộ về TCT sau khi đã che mờ. Quảng Ngãi kiểm soát 100%.")

# =========================================================================
# SLIDE 9: KIẾN TRÚC TRỤC TÍCH HỢP HYBRID (DỮ LIỆU ĐI NHƯ THẾ NÀO?)
# =========================================================================
s9 = prs.slides.add_slide(BL)
header(s9, "KIẾN TRÚC TRỤC TÍCH HỢP: DỮ LIỆU ĐI TỪ QUẢNG NGÃI LÊN TCT NHƯ THẾ NÀO?", "PHẦN 3: CÁCH THỨC KẾT NỐI — SƠ ĐỒ KỸ THUẬT", "08")

pipe_steps = [
    ("BƯỚC 1: NGUỒN DỮ LIỆU", "NỘI BỘ QUẢNG NGÃI", BLUE, BLUE_SOFT,
     ["4 Phần mềm nghiệp vụ:", "• Kế toán (FAST/Bravo)", "• Nhân sự (HRM)", "• Quản lý Vật tư", "• Quản lý Dự án", "Giữ nguyên hiện trạng!"]),
    ("BƯỚC 2: TRẠM TRUNG CHUYỂN", "MÁY CHỦ ẢO (VM) TẠI QN", TEAL, RGBColor(230, 248, 246),
     ["Tự động hóa ngầm:", "• Hút bản sao (Read-Only)", "• Lọc rác, khử trùng", "• Ánh xạ 29 Master Data", "• Masking dữ liệu cá nhân", "• Mã hóa chuẩn AES-256"]),
    ("BƯỚC 3: KÊNH TRUYỀN", "BẢO MẬT TUYỆT ĐỐI", GOLD, GOLD_SOFT,
     ["Hạ tầng mạng an toàn:", "• Đường hầm VPN Site-to-Site", "• Nối Datacenter QN về TCT", "• Tuân thủ 8 Zones ATTT", "• Tự động đồng bộ định kỳ", "• Không lộ IP ra Internet"]),
    ("BƯỚC 4: HẠ TẦNG TCT HUB", "L3 WORKSPACE & POWER BI", GREEN, GREEN_SOFT,
     ["Tiếp nhận & Khai thác:", "• Lưu trữ MinIO On-prem", "• Fabric OneLake Cloud", "• Tự động kích hoạt", "  Dashboard Power BI", "• Ban Giám đốc xem mỗi sáng"])
]

for i, (st_name, st_sub, col, bg_c, lines) in enumerate(pipe_steps):
    x = Inches(0.8 + i * 3.0)
    shape(s9, x, Inches(1.5), Inches(2.65), Inches(5.2), bg_c, col, Pt(1.5), rounded=True)
    shape(s9, x, Inches(1.5), Inches(2.65), Inches(0.7), col)
    text(s9, x, Inches(1.55), Inches(2.65), Inches(0.3), st_name, 11, True, WHITE, PP_ALIGN.CENTER)
    text(s9, x, Inches(1.85), Inches(2.65), Inches(0.25), st_sub, 9, False, RGBColor(220, 235, 255), PP_ALIGN.CENTER)
    
    for j, l in enumerate(lines):
        bold = True if j in [0, len(lines)-1] else False
        c = col if bold else BLACK
        text(s9, x + Inches(0.15), Inches(2.4 + j * 0.45), Inches(2.35), Inches(0.4), l, 10.5, bold, c)
    
    if i < 3:
        arrow_shape(s9, Inches(3.52 + i * 3.0), Inches(3.9), "right", Inches(0.35), Inches(0.25))

note(s9, "Sơ đồ luồng 4 bước: Nguồn dữ liệu tại QN -> Trích xuất bản sao Read-only qua VM làm sạch, đổi mã -> Đẩy qua đường hầm VPN bảo mật -> Tiếp đất vào Workspace L3 trên MinIO và Fabric của TCT -> Ra Dashboard cho Sếp.")

# =========================================================================
# SLIDE 10: SO SÁNH 2 PHƯƠNG ÁN THEO CHỈ ĐẠO BAN GIÁM ĐỐC
# =========================================================================
s10 = prs.slides.add_slide(BL)
header(s10, "SO SÁNH 2 PHƯƠNG ÁN THEO CHỈ ĐẠO BAN GIÁM ĐỐC", "PHẦN 3: CÁCH THỨC KẾT NỐI — LỰA CHỌN TỐI ƯU", "09")

# PA 1: Thuê TCT
shape(s10, Inches(0.8), Inches(1.4), Inches(5.6), Inches(5.4), RED_SOFT, RED, Pt(1.5), rounded=True)
shape(s10, Inches(0.8), Inches(1.4), Inches(5.6), Inches(0.55), RED)
text(s10, Inches(0.8), Inches(1.45), Inches(5.6), Inches(0.4), "PHƯƠNG ÁN 1: THUÊ TỔNG CÔNG TY LÀM HỘ", 12, True, WHITE, PP_ALIGN.CENTER)

pa1_criteria = [
    ("Khả năng đáp ứng của TCT", "BẤT KHẢ THI", RED, "Ban CNTT TCT chỉ làm Hub dùng chung, không đủ nhân lực xuống cắm chốt tại QN bóc tách từng phần mềm."),
    ("Tiến độ thực hiện", "NGUY CƠ TRỄ HẠN CAO", RED, "Xếp hàng sau hơn 10 đơn vị thành viên khác. Chắc chắn không kịp mốc đánh giá KPI năm 2026."),
    ("Mức độ đáp ứng nghiệp vụ", "CHỈ PHỤC VỤ TCT", RED, "TCT chỉ xây báo cáo vĩ mô phục vụ TCT, không xây dựng Dashboard chi tiết đo ni đóng giày cho Sếp QN."),
    ("Chi phí & Tính chủ động", "BỊ ĐỘNG HOÀN TOÀN", RED, "Phụ thuộc lịch trình và quyết định của TCT. QN không chủ động được mã nguồn và nâng cấp.")
]
for i, (crit, val, v_col, desc) in enumerate(pa1_criteria):
    y = Inches(2.15 + i * 1.15)
    shape(s10, Inches(1.0), y, Inches(5.2), Inches(1.0), WHITE, GRAY_MED, Pt(0.5), rounded=True)
    text(s10, Inches(1.15), y + Inches(0.06), Inches(2.5), Inches(0.25), crit, 10.5, True, BLACK)
    text(s10, Inches(3.6), y + Inches(0.06), Inches(2.4), Inches(0.25), f"[{val}]", 10.5, True, v_col, PP_ALIGN.RIGHT)
    text(s10, Inches(1.15), y + Inches(0.35), Inches(4.9), Inches(0.6), desc, 10, color=BLACK)

# PA 2: Thuê NCC
shape(s10, Inches(6.9), Inches(1.4), Inches(5.6), Inches(5.4), GREEN_SOFT, GREEN, Pt(2.0), rounded=True)
shape(s10, Inches(6.9), Inches(1.4), Inches(5.6), Inches(0.55), GREEN)
text(s10, Inches(6.9), Inches(1.45), Inches(5.6), Inches(0.4), "PHƯƠNG ÁN 2 (KHUYẾN NGHỊ): THUÊ NCC ĐỒNG BỘ TCT", 12, True, WHITE, PP_ALIGN.CENTER)

pa2_criteria = [
    ("Khả năng đáp ứng", "CHUYÊN NGHIỆP & TẬN NƠI", GREEN, "NCC chuyên nghiệp về Data cử kỹ sư cắm chốt tại QN, khảo sát kỹ thuật và xây dựng đường ống trực tiếp."),
    ("Tiến độ thực hiện", "3–4 THÁNG (ĐẢM BẢO 2026)", GREEN, "Chủ động kiểm soát tiến độ hoàn toàn. Đảm bảo hoàn thành đúng hạn KPI năm 2026 của TCT."),
    ("Mức độ đáp ứng nghiệp vụ", "ĐO NI ĐÓNG GIÀY CHO QN", GREEN, "Vừa đồng bộ dữ liệu lên TCT, vừa xây dựng Dashboard quản trị phục vụ trực tiếp Ban Giám đốc QN."),
    ("Làm chủ công nghệ", "BÀN GIAO MÃ NGUỒN 100%", GREEN, "Đào tạo IT Quảng Ngãi làm chủ vận hành. Tận dụng miễn phí hạ tầng Cloud & License của TCT.")
]
for i, (crit, val, v_col, desc) in enumerate(pa2_criteria):
    y = Inches(2.15 + i * 1.15)
    shape(s10, Inches(7.1), y, Inches(5.2), Inches(1.0), WHITE, GRAY_MED, Pt(0.5), rounded=True)
    text(s10, Inches(7.25), y + Inches(0.06), Inches(2.5), Inches(0.25), crit, 10.5, True, BLACK)
    text(s10, Inches(9.6), y + Inches(0.06), Inches(2.5), Inches(0.25), f"[{val}]", 10.5, True, v_col, PP_ALIGN.RIGHT)
    text(s10, Inches(7.25), y + Inches(0.35), Inches(4.9), Inches(0.6), desc, 10, color=BLACK)

note(s10, "So sánh 2 phương án: TCT chỉ làm Hub chung, không làm thay được cho mình. Thuê NCC chuyên nghiệp là phương án khả thi duy nhất để vừa kịp tiến độ 2026, vừa có Dashboard phục vụ đúng nhu cầu Ban Giám đốc.")

# =========================================================================
# SLIDE 11: CƠ CẤU ĐẦU TƯ & CHI PHÍ (THEO PHƯƠNG ÁN 3 CỦA TCT)
# =========================================================================
s11 = prs.slides.add_slide(BL)
header(s11, "CƠ CẤU ĐẦU TƯ & CẤU TRÚC CHI PHÍ (THEO PHƯƠNG ÁN 3 CỦA TCT)", "PHẦN 3: CÁCH THỨC KẾT NỐI — DỰ TOÁN NGÂN SÁCH", "10")

cost_rows = [
    ("1", "Hạ tầng máy chủ tại Đơn vị", "Tận dụng máy chủ ảo (VM) sẵn có tại QN", "0 VNĐ (Tự có)", GREEN),
    ("2", "Hạ tầng Cloud Hub & License", "Fabric OneLake 20TB, MinIO, PowerBI", "TCT ĐÃ ĐẦU TƯ", BLUE),
    ("3", "Dịch vụ: Khảo sát & Thiết kế kiến trúc", "Khảo sát CSDL 4 phần mềm + Ánh xạ 29 Master Data", "[Chờ NCC báo giá GĐ1]", GOLD),
    ("4", "Dịch vụ: Xây dựng Trục tích hợp (ETL)", "Viết đường ống trích xuất, làm sạch, mã hóa VPN", "[Chờ NCC báo giá GĐ1]", GOLD),
    ("5", "Dịch vụ: Xây dựng Dashboard Power BI", "Xây dựng các bảng Dashboard điều hành cho Sếp QN", "[Chờ NCC báo giá GĐ1]", GOLD),
    ("6", "Đào tạo chuyển giao & Bảo hành", "Bàn giao mã nguồn, đào tạo IT QN làm chủ 100%", "[Chờ NCC báo giá GĐ1]", GOLD),
    ("7", "Phí duy trì nền tảng hàng năm", "TCT phân bổ chi phí theo mức độ sử dụng (Usage)", "Theo cơ chế TCT", GRAY),
]

shape(s11, Inches(0.8), Inches(1.4), Inches(11.7), Inches(0.4), NAVY)
text(s11, Inches(0.9), Inches(1.45), Inches(0.6), Inches(0.3), "STT", 10, True, WHITE, PP_ALIGN.CENTER)
text(s11, Inches(1.6), Inches(1.45), Inches(3.2), Inches(0.3), "CẤU PHẦN CHI PHÍ (CHUẨN TCT)", 10, True, WHITE)
text(s11, Inches(4.9), Inches(1.45), Inches(4.5), Inches(0.3), "NỘI DUNG & TRÁCH NHIỆM THỰC HIỆN", 10, True, WHITE)
text(s11, Inches(9.5), Inches(1.45), Inches(2.9), Inches(0.3), "DỰ TOÁN NGÂN SÁCH", 10, True, WHITE, PP_ALIGN.CENTER)

for i, (stt, comp, desc, cost, c_col) in enumerate(cost_rows):
    y = Inches(1.85 + i * 0.6)
    bg = WHITE if i % 2 == 0 else GRAY_LIGHT
    shape(s11, Inches(0.8), y, Inches(11.7), Inches(0.55), bg, GRAY_MED, Pt(0.5))
    text(s11, Inches(0.9), y + Inches(0.1), Inches(0.6), Inches(0.35), stt, 10, True, BLACK, PP_ALIGN.CENTER)
    text(s11, Inches(1.6), y + Inches(0.1), Inches(3.2), Inches(0.35), comp, 10, True, BLACK)
    text(s11, Inches(4.9), y + Inches(0.1), Inches(4.5), Inches(0.35), desc, 9.5, color=BLACK)
    text(s11, Inches(9.5), y + Inches(0.1), Inches(2.9), Inches(0.35), cost, 10, True, c_col, PP_ALIGN.CENTER)

# Bottom note
shape(s11, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.7), GOLD_SOFT, GOLD, Pt(1.0), rounded=True)
text(s11, Inches(1.1), Inches(6.25), Inches(11.1), Inches(0.55),
     "📌 NGUYÊN TẮC QUẢN TRỊ TÀI CHÍNH: Tiết kiệm hàng tỷ đồng phần cứng/license nhờ TCT.\n"
     "Kinh phí thuê NCC sẽ được xác định chính xác sau Giai đoạn 1 (khảo sát kỹ thuật & lấy báo giá cạnh tranh).",
     10.5, True, NAVY, PP_ALIGN.CENTER)

note(s11, "Về chi phí, chúng ta bám sát Phương án 3 của TCT: TCT đã gánh phần Cloud và phần mềm khung. Phần chi phí thuê NCC triển khai tại chỗ, chúng tôi xin phép để trống và sẽ xác định chính xác sau khi khảo sát kỹ thuật ở GĐ1.")

# =========================================================================
# SLIDE 12: ĐẦU RA CHO BAN GIÁM ĐỐC & BÀI TOÁN QUICK-WIN (4-6 TUẦN)
# =========================================================================
s12 = prs.slides.add_slide(BL)
header(s12, "ĐẦU RA CHO BAN GIÁM ĐỐC & BÀI TOÁN 'THẮNG NHANH' (QUICK-WIN)", "PHẦN 4: KẾ HOẠCH HÀNH ĐỘNG — KẾT QUẢ THỰC TẾ", "11")

# Left Column: 4 Giá trị cụ thể
shape(s12, Inches(0.8), Inches(1.4), Inches(5.6), Inches(5.4), BLUE_SOFT, BLUE, Pt(1.5), rounded=True)
shape(s12, Inches(0.8), Inches(1.4), Inches(5.6), Inches(0.55), BLUE)
text(s12, Inches(0.8), Inches(1.45), Inches(5.6), Inches(0.4), "4 GIÁ TRỊ BAN GIÁM ĐỐC NHẬN ĐƯỢC", 11.5, True, WHITE, PP_ALIGN.CENTER)

outcomes = [
    ("Dashboard điều hành Real-Time", "Mở xem số liệu SXKD, doanh thu, dòng tiền trên điện thoại/laptop mỗi sáng trước giờ giao ban."),
    ("Giải phóng phụ thuộc Excel thủ công", "Không còn cảnh chờ đợi 3–5 ngày để các phòng ban xuất và ghép file báo cáo."),
    ("Một nguồn sự thật duy nhất (Single Truth)", "Số liệu Kế toán, Vật tư, Dự án khớp nhau 100%, không còn vênh số liệu khi giải trình."),
    ("Hoàn thành xuất sắc chỉ tiêu của TCT", "Đáp ứng đúng tiến độ Nghị quyết 10 HĐQT PTSC và tuân thủ Luật Bảo vệ DLCN mới.")
]
for i, (title, desc) in enumerate(outcomes):
    y = Inches(2.15 + i * 1.15)
    shape(s12, Inches(1.0), y, Inches(5.2), Inches(1.0), WHITE, GRAY_MED, Pt(0.5), rounded=True)
    text(s12, Inches(1.2), y + Inches(0.08), Inches(4.8), Inches(0.3), f"✔ {title}", 11, True, BLUE)
    text(s12, Inches(1.2), y + Inches(0.35), Inches(4.8), Inches(0.6), desc, 10, color=BLACK)

# Right Column: Quick-Win Pilot
shape(s12, Inches(6.9), Inches(1.4), Inches(5.6), Inches(5.4), GOLD_SOFT, GOLD, Pt(2.0), rounded=True)
shape(s12, Inches(6.9), Inches(1.4), Inches(5.6), Inches(0.55), GOLD)
text(s12, Inches(6.9), Inches(1.45), Inches(5.6), Inches(0.4), "BÀI TOÁN THÍ ĐIỂM THẮNG NHANH (QUICK-WIN IN 4-6 WEEKS)", 11.5, True, WHITE, PP_ALIGN.CENTER)

shape(s12, Inches(7.1), Inches(2.15), Inches(5.2), Inches(1.15), WHITE, GOLD, Pt(1.0), rounded=True)
text(s12, Inches(7.3), Inches(2.22), Inches(4.8), Inches(0.3), "★ ĐỀ XUẤT THÍ ĐIỂM: QUẢN LÝ DOANH THU & CHI PHÍ DỰ ÁN", 11, True, NAVY)
text(s12, Inches(7.3), Inches(2.55), Inches(4.8), Inches(0.7),
     "Không dàn trải cả 4 phần mềm cùng lúc. Chọn ngay mảng Dự án trọng điểm & Chi phí Kế toán để xây dựng Dashboard đầu tiên trong 4–6 tuần.", 10, color=BLACK)

qw_steps = [
    ("Tuần 1–2", "Trích xuất CSDL Dự án & Kế toán (Read-Only)"),
    ("Tuần 3–4", "Viết luồng ETL làm sạch & Ánh xạ mã công trình/hợp đồng"),
    ("Tuần 5–6", "Bàn giao Dashboard đầu tiên cho Ban Giám đốc trải nghiệm")
]
for i, (w, d) in enumerate(qw_steps):
    y = Inches(3.45 + i * 0.75)
    shape(s12, Inches(7.1), y, Inches(5.2), Inches(0.65), WHITE, GRAY_MED, Pt(0.5), rounded=True)
    text(s12, Inches(7.3), y + Inches(0.12), Inches(1.2), Inches(0.35), f"[{w}]", 10.5, True, GOLD)
    text(s12, Inches(8.5), y + Inches(0.12), Inches(3.6), Inches(0.35), d, 10, color=BLACK)

shape(s12, Inches(7.1), Inches(5.85), Inches(5.2), Inches(0.8), GREEN_SOFT, GREEN, Pt(1.0), rounded=True)
text(s12, Inches(7.3), Inches(5.9), Inches(4.8), Inches(0.7),
     "👉 LỢI ÍCH: Ban Giám đốc thấy ngay sản phẩm thực tế chỉ sau hơn 1 tháng, tạo niềm tin tuyệt đối trước khi mở rộng toàn diện.", 10, True, GREEN)

note(s12, "Điểm mới bổ sung: Chúng ta đề xuất 1 bài toán Quick-Win làm trước trong 4-6 tuần: Dashboard Quản lý Doanh thu & Chi phí Dự án. Sếp sẽ thấy ngay sản phẩm thực tế trên điện thoại chỉ sau hơn 1 tháng.")

# =========================================================================
# SLIDE 13: KẾ HOẠCH HÀNH ĐỘNG 3 GIAI ĐOẠN (TRỌNG TÂM GIAI ĐOẠN 1)
# =========================================================================
s13 = prs.slides.add_slide(BL)
header(s13, "KẾ HOẠCH HÀNH ĐỘNG 3 GIAI ĐOẠN: TRỌNG TÂM GIAI ĐOẠN 1 (T3-T4/2026)", "PHẦN 4: KẾ HOẠCH HÀNH ĐỘNG — LỘ TRÌNH THỰC HIỆN", "12")

phases_detail = [
    ("GIAI ĐOẠN 1 (T3 – T4/2026)", "KHẢO SÁT & ĐỀ BÀI KỸ THUẬT", "NỘI BỘ TỰ LÀM — CHI PHÍ 0 VNĐ", GOLD, GOLD_SOFT,
     ["1. Thành lập Tổ công tác Data Platform nội bộ (CNTT + Key Users).",
      "2. Lập danh mục CSDL 4 phần mềm (Data Catalog).",
      "3. Lập bảng đối chiếu ánh xạ với 29 Master Data của TCT.",
      "4. Lập bài toán thí điểm Quick-Win.",
      "5. Hoàn thiện Hồ sơ kỹ thuật (TOR) mời các NCC gửi báo giá cạnh tranh."]),
    ("GIAI ĐOẠN 2 (T5 – T7/2026)", "TRIỂN KHAI TRỤC TÍCH HỢP", "THUÊ NCC CHUYÊN NGHIỆP", BLUE, BLUE_SOFT,
     ["1. Lựa chọn NCC tối ưu thông qua xét duyệt chào giá cạnh tranh.",
      "2. Cấp phát máy chủ ảo (VM) làm trạm trung chuyển dữ liệu tại QN.",
      "3. Lập trình luồng ETL tự động hút, làm sạch và mã hóa dữ liệu.",
      "4. Cấu hình VPN bảo mật kết nối lên Workspace L3 trên Hub TCT.",
      "5. Thiết kế và bàn giao các Dashboard quản trị Power BI."]),
    ("GIAI ĐOẠN 3 (T8/2026 TRỞ ĐI)", "BÀN GIAO & KHAI THÁC", "IT QUẢNG NGÃI LÀM CHỦ", GREEN, GREEN_SOFT,
     ["1. Nghiệm thu kỹ thuật và vận hành thử nghiệm toàn hệ thống.",
      "2. Bàn giao 100% mã nguồn (Source code) và tài liệu kiến trúc.",
      "3. Đào tạo chuyển giao công nghệ cho đội ngũ IT Quảng Ngãi.",
      "4. Đưa vào sử dụng chính thức phục vụ các phiên họp giao ban của BGĐ.",
      "5. Báo cáo TCT hoàn thành chỉ tiêu KPI Chuyển đổi số năm 2026."])
]

for i, (p_title, p_name, p_cost, col, bg_c, tasks) in enumerate(phases_detail):
    x = Inches(0.8 + i * 4.05)
    shape(s13, x, Inches(1.45), Inches(3.7), Inches(5.4), bg_c, col, Pt(1.5), rounded=True)
    shape(s13, x, Inches(1.45), Inches(3.7), Inches(0.8), col)
    text(s13, x, Inches(1.5), Inches(3.7), Inches(0.28), p_title, 11.5, True, WHITE, PP_ALIGN.CENTER)
    text(s13, x, Inches(1.78), Inches(3.7), Inches(0.25), p_name, 10, True, RGBColor(220, 235, 255), PP_ALIGN.CENTER)
    
    # Cost badge
    shape(s13, x + Inches(0.2), Inches(2.35), Inches(3.3), Inches(0.35), WHITE, col, Pt(1.0), rounded=True)
    text(s13, x + Inches(0.2), Inches(2.4), Inches(3.3), Inches(0.25), p_cost, 9.5, True, col, PP_ALIGN.CENTER)
    
    # Task list
    for j, t in enumerate(tasks):
        text(s13, x + Inches(0.2), Inches(2.85 + j * 0.78), Inches(3.3), Inches(0.72), t, 10, color=BLACK)

note(s13, "Lộ trình 3 giai đoạn rất rõ: Giai đoạn 1 nội bộ tự làm hoàn toàn miễn phí trong 2 tháng để ra bài toán chuẩn; sau đó mới tổ chức mời thầu NCC trong Giai đoạn 2; và Giai đoạn 3 là nghiệm thu, bàn giao để IT nội bộ làm chủ 100%.")

# =========================================================================
# SLIDE 14: CAM KẾT NGUỒN LỰC CÁC PHÒNG BAN & AN TOÀN KHI THUÊ NCC
# =========================================================================
s14 = prs.slides.add_slide(BL)
header(s14, "CAM KẾT NGUỒN LỰC CÁC PHÒNG BAN & AN TOÀN KHI THUÊ NCC", "PHẦN 4: KẾ HOẠCH HÀNH ĐỘNG — GIẢI TỎA MỌI LO NGẠI", "13")

# Left Column: Cam kết các phòng ban
shape(s14, Inches(0.8), Inches(1.4), Inches(5.6), Inches(5.4), GREEN_SOFT, GREEN, Pt(1.5), rounded=True)
shape(s14, Inches(0.8), Inches(1.4), Inches(5.6), Inches(0.55), GREEN)
text(s14, Inches(0.8), Inches(1.45), Inches(5.6), Inches(0.4), "CAM KẾT VỚI CÁC PHÒNG BAN (KHÔNG GÂY QUÁ TẢI)", 11.5, True, WHITE, PP_ALIGN.CENTER)

dept_commitments = [
    ("Không phát sinh công việc nhập liệu", "Hệ thống tự động đọc dữ liệu từ phần mềm hiện tại. Cán bộ các phòng KHÔNG PHẢI GÕ THÊM BẤT KỲ DÒNG NÀO."),
    ("Thời lượng tham gia cực ít (2–3 buổi)", "Key Users mỗi phòng chỉ cần tham gia 2–3 buổi phỏng vấn (1-2 giờ/buổi) để giải thích cấu trúc dữ liệu và kiểm tra số liệu."),
    ("90% công việc kỹ thuật do CNTT & NCC gánh", "Toàn bộ việc viết code, làm sạch dữ liệu, dựng trạm trung chuyển do IT và kỹ sư NCC đảm nhiệm 100%."),
    ("Trực tiếp hưởng lợi sau khi hoàn thành", "Các phòng ban được cấp tài khoản xem báo cáo tự động của phòng mình, chấm dứt việc cặm cụi làm Excel báo cáo sếp.")
]
for i, (title, desc) in enumerate(dept_commitments):
    y = Inches(2.15 + i * 1.15)
    shape(s14, Inches(1.0), y, Inches(5.2), Inches(1.0), WHITE, GRAY_MED, Pt(0.5), rounded=True)
    text(s14, Inches(1.2), y + Inches(0.08), Inches(4.8), Inches(0.3), f"✔ {title}", 11, True, GREEN)
    text(s14, Inches(1.2), y + Inches(0.35), Inches(4.8), Inches(0.6), desc, 10, color=BLACK)

# Right Column: An toàn khi thuê NCC
shape(s14, Inches(6.9), Inches(1.4), Inches(5.6), Inches(5.4), BLUE_SOFT, BLUE, Pt(1.5), rounded=True)
shape(s14, Inches(6.9), Inches(1.4), Inches(5.6), Inches(0.55), BLUE)
text(s14, Inches(6.9), Inches(1.45), Inches(5.6), Inches(0.4), "CAM KẾT AN TOÀN & CHỐNG PHỤ THUỘC NCC", 11.5, True, WHITE, PP_ALIGN.CENTER)

ncc_commitments = [
    ("Chỉ cấp quyền Đọc (Read-Only Replica)", "NCC chỉ được tiếp cận bản sao dữ liệu. Tuyệt đối không cấp quyền vào Database gốc đang chạy sản xuất."),
    ("Ký thỏa thuận bảo mật thông tin (NDA)", "Ràng buộc trách nhiệm pháp lý chặt chẽ theo Luật BVDLCN mới. Nghiêm cấm sao chép dữ liệu ra ngoài."),
    ("Bàn giao 100% mã nguồn (Source Code)", "Hợp đồng bắt buộc bàn giao toàn bộ script ETL, từ điển dữ liệu để IT Quảng Ngãi hoàn toàn làm chủ."),
    ("Cơ chế 'Hai nấc độc lập' với TCT", "Hệ thống nội bộ phục vụ Sếp QN chạy độc lập. Nếu TCT chậm trễ, tiến độ của Quảng Ngãi vẫn không bị ảnh hưởng.")
]
for i, (title, desc) in enumerate(ncc_commitments):
    y = Inches(2.15 + i * 1.15)
    shape(s14, Inches(7.1), y, Inches(5.2), Inches(1.0), WHITE, GRAY_MED, Pt(0.5), rounded=True)
    text(s14, Inches(7.3), y + Inches(0.08), Inches(4.8), Inches(0.3), f"🔒 {title}", 11, True, BLUE)
    text(s14, Inches(7.3), y + Inches(0.35), Inches(4.8), Inches(0.6), desc, 10, color=BLACK)

note(s14, "Slide này giải tỏa 2 mối lo lớn nhất: Các phòng ban yên tâm không bị đè thêm việc, và Ban Giám đốc hoàn toàn yên tâm về an toàn thông tin khi thuê NCC.")

# =========================================================================
# SLIDE 15: KIẾN NGHỊ BAN GIÁM ĐỐC PHÊ DUYỆT (CALL TO ACTION)
# =========================================================================
s15 = prs.slides.add_slide(BL)
header(s15, "KIẾN NGHỊ BAN GIÁM ĐỐC PHÊ DUYỆT ĐỂ TRIỂN KHAI", "PHẦN 4: KẾ HOẠCH HÀNH ĐỘNG — ĐỀ XUẤT PHÊ DUYỆT", "14")

proposals = [
    ("1", "PHÊ DUYỆT CHỦ TRƯƠNG PHƯƠNG ÁN 2",
     "Đồng ý chủ trương: Công ty chủ động phát triển Trục tích hợp dữ liệu riêng đồng bộ với TCT (thuê NCC chuyên nghiệp tư vấn & triển khai) để đảm bảo tiến độ KPI năm 2026 và phục vụ điều hành nội bộ.",
     BLUE),
    ("2", "THÀNH LẬP TỔ CÔNG TÁC DATA PLATFORM NỘI BỘ",
     "Gồm Bộ phận CNTT chủ trì kỹ thuật và Key Users các phòng ban (Kế toán, Nhân sự, Vật tư, Dự án) để triển khai ngay Giai đoạn 1 (Khảo sát & Lập bảng ánh xạ 29 Master Data — Chi phí 0 VNĐ).",
     GOLD),
    ("3", "CHO PHÉP TIẾP XÚC NCC ĐỂ LẤY BÁO GIÁ CẠNH TRANH",
     "Cho phép Tổ công tác tiếp xúc các NCC chuyên nghiệp để lấy đề xuất giải pháp và báo giá cạnh tranh. Hoàn thiện dự toán chi tiết trình Ban Giám đốc xem xét phê duyệt trước khi ký kết hợp đồng.",
     GREEN)
]

for i, (num, title, desc, col) in enumerate(proposals):
    y = Inches(1.5 + i * 1.75)
    shape(s15, Inches(0.8), y, Inches(11.7), Inches(1.55), WHITE, col, Pt(1.5), rounded=True)
    shape(s15, Inches(0.8), y, Inches(0.6), Inches(1.55), col)
    numbered_label(s15, Inches(0.9), y + Inches(0.55), num, col)
    text(s15, Inches(1.7), y + Inches(0.12), Inches(10.5), Inches(0.35), title, 13, True, col)
    shape(s15, Inches(1.7), y + Inches(0.5), Inches(10.5), Inches(0.01), GRAY_MED)
    text(s15, Inches(1.7), y + Inches(0.6), Inches(10.5), Inches(0.85), desc, 11, color=BLACK)

note(s15, "Kính thưa Ban Giám đốc, bộ phận CNTT kính đề xuất Ban Giám đốc thông qua 3 chủ trương nêu trên để chúng tôi kịp triển khai ngay Giai đoạn 1 trong tháng 3 và đảm bảo tiến độ chung của Tổng công ty. Xin trân trọng cảm ơn Ban Giám đốc!")

# Save to pptx
out_path = r"d:\My Profiles\DataPlatform\bao_cao_dataplatform_ptsc_qn_v6.pptx"
prs.save(out_path)
print(f"Presentation v6 saved successfully to: {out_path}")
