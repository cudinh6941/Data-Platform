import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    # 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6] # blank layout

    # Color Palette
    PRIMARY = RGBColor(11, 44, 86)       # #0B2C56 - Deep Navy Corporate
    SECONDARY = RGBColor(21, 93, 168)    # #155DA8 - Professional Blue
    ACCENT = RGBColor(220, 53, 69)       # #DC3545 - Accent Crimson/Red
    ACCENT_GOLD = RGBColor(217, 119, 6)  # #D97706 - Warning/Highlight Gold
    BG_LIGHT = RGBColor(245, 247, 250)   # #F5F7FA - Clean light gray
    CARD_BG = RGBColor(255, 255, 255)    # White
    CARD_BORDER = RGBColor(218, 225, 233)
    TEXT_DARK = RGBColor(30, 41, 59)     # Charcoal slate
    TEXT_MUTED = RGBColor(100, 116, 139) # Cool gray
    WHITE = RGBColor(255, 255, 255)

    def add_header(slide, title_text, category_text="PTSC QUẢNG NGÃI - CHIẾN LƯỢC CHUYỂN ĐỔI SỐ"):
        # Top banner background
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.15))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = PRIMARY
        top_bar.line.fill.background()

        # Category text
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.12), Inches(11.5), Inches(0.3))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = RGBColor(186, 215, 245)

        # Title text
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.6))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Bottom subtle line on banner
        sub_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.15), Inches(13.333), Inches(0.04))
        sub_line.fill.solid()
        sub_line.fill.fore_color.rgb = ACCENT_GOLD
        sub_line.line.fill.background()

        # Footer
        footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.1), Inches(11.733), Inches(0.3))
        tf_foot = footer_box.text_frame
        p_foot = tf_foot.paragraphs[0]
        p_foot.text = "Báo cáo Ban Giám đốc | Nền tảng Dữ liệu (Data Platform) & Trạm trung chuyển dữ liệu"
        p_foot.font.size = Pt(9)
        p_foot.font.color.rgb = TEXT_MUTED

    def add_card(slide, left, top, width, height, title, subtitle="", border_color=CARD_BORDER, bg_color=CARD_BG):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.2)

        tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), height - Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = PRIMARY
        
        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.font.size = Pt(11)
            p2.font.color.rgb = TEXT_MUTED
            p2.space_after = Pt(8)
            
        return tf

    # ==========================================
    # SLIDE 1: COVER
    # ==========================================
    s1 = prs.slides.add_slide(blank_layout)
    bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = PRIMARY
    bg1.line.fill.background()

    # Decorative accent card
    accent_bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(0.15), Inches(3.8))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT_GOLD
    accent_bar.line.fill.background()

    # Main title
    t_box = s1.shapes.add_textbox(Inches(1.2), Inches(1.7), Inches(11.0), Inches(2.2))
    tf1 = t_box.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "CHIẾN LƯỢC CHUYỂN ĐỔI SỐ PTSC"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(186, 215, 245)
    p.space_after = Pt(10)

    p2 = tf1.add_paragraph()
    p2.text = "XÂY DỰNG NỀN TẢNG DỮ LIỆU (DATA PLATFORM)\n& TRẠM TRUNG CHUYỂN DỮ LIỆU TẠI PTSC QUẢNG NGÃI"
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.space_after = Pt(16)

    p3 = tf1.add_paragraph()
    p3.text = "Báo cáo Toàn cảnh, Thực trạng kết nối Tổng công ty và Đề xuất Phương án Thuê ngoài triển khai"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(220, 230, 242)

    # Info card
    info_box = s1.shapes.add_textbox(Inches(1.2), Inches(5.2), Inches(10.5), Inches(1.2))
    tf_info = info_box.text_frame
    p_i1 = tf_info.paragraphs[0]
    p_i1.text = "• Đơn vị báo cáo: Bộ phận IT & Chuyển đổi số - PTSC Quảng Ngãi"
    p_i1.font.size = Pt(13)
    p_i1.font.color.rgb = WHITE
    p_i2 = tf_info.add_paragraph()
    p_i2.text = "• Kính trình: Ban Giám đốc Công ty PTSC Quảng Ngãi"
    p_i2.font.size = Pt(13)
    p_i2.font.color.rgb = RGBColor(253, 224, 71)

    s1.notes_slide.notes_text_frame.text = (
        "KỊCH BẢN NÓI (SPEAKER NOTES):\n"
        "Kính thưa Ban Giám đốc, hôm nay phòng IT xin phép trình bày bức tranh toàn cảnh về kế hoạch triển khai "
        "Nền tảng dữ liệu (Data Platform) tại đơn vị. Buổi báo cáo này sẽ làm rõ vị thế của Quảng Ngãi trong hệ sinh thái "
        "dùng chung của Tổng công ty, thực trạng hiện tại của chúng ta, và đề xuất phương án tối ưu nhất về tiến độ, chất lượng "
        "thông qua việc kết hợp chuyên gia thuê ngoài."
    )

    # ==========================================
    # SLIDE 2: TỔNG QUAN VỀ DATA PLATFORM
    # ==========================================
    s2 = prs.slides.add_slide(blank_layout)
    add_header(s2, "1. TỔNG QUAN: DATA PLATFORM LÀ GÌ?", "HIỂU ĐÚNG BẢN CHẤT DƯỚI GÓC NHÌN QUẢN TRỊ")

    # Card 1: Khái niệm đơn giản
    tf2_1 = add_card(s2, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.2), "Bản chất: 'Hệ tuần hoàn máu' Doanh nghiệp")
    p = tf2_1.add_paragraph()
    p.text = "Data Platform không đơn thuần là một phần mềm mới, mà là TRỤC KẾT NỐI VÀ LÀM SẠCH toàn bộ dữ liệu công ty:"
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_DARK
    p.space_after = Pt(10)

    points = [
        ("Thu gom tự động:", " Hút dữ liệu từ Kế toán, Vật tư, Nhân sự, Hợp đồng... mà không cần người xuất file."),
        ("Chuẩn hóa tập trung:", " Quy đổi toàn bộ dữ liệu về một ngôn ngữ chuẩn mực chung."),
        ("Cung cấp tức thời:", " Bơm dữ liệu lên hệ thống báo cáo Dashboard thông minh cho Ban Giám đốc.")
    ]
    for bold_prefix, text in points:
        p = tf2_1.add_paragraph()
        run1 = p.add_run()
        run1.text = "• " + bold_prefix
        run1.font.bold = True
        run1.font.color.rgb = SECONDARY
        run2 = p.add_run()
        run2.text = text
        run2.font.color.rgb = TEXT_DARK
        p.space_after = Pt(8)

    # Card 2: Sự khác biệt Trước - Sau
    tf2_2 = add_card(s2, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.2), "Sự chuyển dịch mô hình quản trị số")
    
    comp = [
        ("TRƯỚC ĐÂY (MÔ HÌNH ỐC ĐẢO):", [
            "Dữ liệu nằm phân tán ở các máy chủ phần mềm riêng biệt.",
            "Muốn số liệu, nhân viên phải xuất Excel, copy-paste ghép thủ công.",
            "Số liệu lệch pha: Kế toán báo một kiểu, Vật tư báo một kiểu."
        ], RGBColor(192, 41, 43)),
        ("KHI CÓ DATA PLATFORM:", [
            "Một nguồn sự thật duy nhất (Single Source of Truth).",
            "Dữ liệu cập nhật liên tục theo ngày/giờ, không phụ thuộc người làm báo cáo.",
            "Tự động cảnh báo rủi ro về tồn kho, dòng tiền, trễ hạn dự án."
        ], RGBColor(39, 174, 96))
    ]

    for title, items, col in comp:
        p = tf2_2.add_paragraph()
        p.text = title
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = col
        p.space_after = Pt(4)
        for it in items:
            pi = tf2_2.add_paragraph()
            pi.text = "  - " + it
            pi.font.size = Pt(12)
            pi.font.color.rgb = TEXT_DARK
            pi.space_after = Pt(3)
        tf2_2.add_paragraph().space_after = Pt(6)

    s2.notes_slide.notes_text_frame.text = (
        "KỊCH BẢN NÓI:\n"
        "Thưa các anh/chị, để dễ hình dung, Data Platform giống như một hệ thống đường ống nước tự động. "
        "Trước đây mỗi phòng ban là một giếng nước độc lập, muốn dùng nước thì phải xách xô (chính là file Excel thủ công). "
        "Nay Data Platform nối tất cả lại, qua hệ thống lọc sạch và bơm thẳng lên vòi nước tại phòng Ban Giám đốc."
    )

    # ==========================================
    # SLIDE 3: TẦM QUAN TRỌNG ĐỐI VỚI LÃNH ĐẠO
    # ==========================================
    s3 = prs.slides.add_slide(blank_layout)
    add_header(s3, "2. TẦM QUAN TRỌNG: BAN GIÁM ĐỐC & ĐƠN VỊ ĐƯỢC GÌ?", "LỢI ÍCH TRỰC TIẾP CHO CÔNG TÁC ĐIỀU HÀNH")

    cards_s3 = [
        ("QUẢN TRỊ TỨC THỜI (REAL-TIME)", 
         "Không còn độ trễ thông tin", 
         [
             "Thay thế việc chờ đợi 3-5 ngày để các phòng ban tổng hợp báo cáo tuần/tháng.",
             "Theo dõi số liệu sản xuất kinh doanh, doanh thu, dòng tiền trực quan trên di động & máy tính.",
             "Phát hiện ngay nút thắt công việc khi có số liệu bất thường."
         ], Inches(0.8)),
        ("MINH BẠCH & CHUẨN XÁC", 
         "Một nguồn sự thật duy nhất", 
         [
             "Chấm dứt hoàn toàn hiện tượng 'lệch số' giữa báo cáo Kế toán, Vật tư và Dự án.",
             "Số liệu được lấy thẳng từ cơ sở dữ liệu gốc, loại bỏ hoàn toàn sai lệch do can thiệp chủ quan.",
             "Tăng độ tin cậy khi giải trình với Tổng công ty và Kiểm toán."
         ], Inches(4.8)),
        ("TIẾT KIỆM TỐI ĐA NGÂN SÁCH", 
         "Đứng trên vai người khổng lồ", 
         [
             "Tổng công ty đã đầu tư hàng triệu USD mua hạ tầng Microsoft Fabric và bản quyền Power BI.",
             "Quảng Ngãi được dùng MIỄN PHÍ tài nguyên lưu trữ và công cụ đắt tiền này.",
             "Chúng ta chỉ cần xây 'đường ống kết nối' tại chỗ mà không cần mua Server Data khổng lồ."
         ], Inches(8.8))
    ]

    for title, subtitle, bullets, left_pos in cards_s3:
        tf3 = add_card(s3, left_pos, Inches(1.5), Inches(3.7), Inches(5.2), title, subtitle)
        for b in bullets:
            p = tf3.add_paragraph()
            p.text = "• " + b
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT_DARK
            p.space_after = Pt(8)

    s3.notes_slide.notes_text_frame.text = (
        "KỊCH BẢN NÓI:\n"
        "Lợi ích lớn nhất của dự án này trước hết là phục vụ chính công tác điều hành của Ban Giám đốc: "
        "Thứ nhất, có báo cáo tức thì mà không phải chờ đợi. Thứ hai, chuẩn hóa dữ liệu minh bạch, số liệu khớp từng đồng. "
        "Và thứ ba, chúng ta tiết kiệm được ngân sách rất lớn vì chỉ việc tận dụng hạ tầng đám mây khủng mà TCT đã thanh toán."
    )

    # ==========================================
    # SLIDE 4: TOÀN CẢNH TỔNG CÔNG TY PTSC
    # ==========================================
    s4 = prs.slides.add_slide(blank_layout)
    add_header(s4, "3. HIỆN TRẠNG TCT: 'QUẢ BÓNG' ĐÃ CHUYỀN VỀ QUẢNG NGÃI", "LỘ TRÌNH VÀ YÊU CẦU BẮT BUỘC TỪ NGHỊ QUYẾT HĐQT")

    # Left: TCT đã hoàn thành gì
    tf4_left = add_card(s4, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.2), "Tổng Công Ty đã chuẩn bị xong nền móng")
    points_tct = [
        ("Nghị quyết số 10/NQ-HĐQT-PTSC:", " Ban hành Chiến lược CĐS toàn diện, trong đó Nền tảng Dữ liệu là trụ cột bắt buộc đối với toàn bộ các đơn vị thành viên."),
        ("Giai đoạn 1 (2024 - 2025) ĐÃ XONG:", " TCT đã xây dựng hoàn chỉnh Trung tâm Dữ liệu tập trung (Data Lakehouse) trên nền tảng Microsoft Fabric."),
        ("Ban hành 29 Danh mục Master Data:", " Bộ quy chuẩn chuẩn hóa mã nhân sự, mã khách hàng, mã thiết bị, mã chi phí... áp dụng thống nhất toàn TCT."),
        ("Cấp sẵn Workspace riêng:", " TCT đã mở sẵn phân vùng 'Workspace L3 Quảng Ngãi' trên Cloud để sẵn sàng đón nhận dữ liệu từ chúng ta.")
    ]
    for bp, txt in points_tct:
        p = tf4_left.add_paragraph()
        r1 = p.add_run()
        r1.text = "✔ " + bp
        r1.font.bold = True
        r1.font.color.rgb = SECONDARY
        r2 = p.add_run()
        r2.text = txt
        r2.font.color.rgb = TEXT_DARK
        p.space_after = Pt(8)

    # Right: Mốc thời gian 2026 - 2027
    tf4_right = add_card(s4, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.2), "Áp lực thời gian: Giai đoạn 2 (2026 - 2027)", "THỜI ĐIỂM QUẢNG NGÃI PHẢI KẾT NỐI", ACCENT)
    
    stages = [
        ("2024 - 2025 (Đã hoàn thành)", "Trọng tâm tại TCT: Xây lõi Lakehouse và ban hành quy chuẩn 29 Master Data.", RGBColor(100, 116, 139)),
        ("2026 - 2027 (THỜI ĐIỂM HIỆN TẠI)", "Trọng tâm tại Đơn vị thành viên (Quảng Ngãi):\n• Chuẩn hóa dữ liệu nội bộ khớp 29 danh mục.\n• Dựng trạm trung chuyển và kết nối đường ống bơm dữ liệu lên TCT.\n• Đây là tiêu chí đánh giá KPI chuyển đổi số của đơn vị!", ACCENT),
        ("2028 - 2030 (Tương lai)", "Khai thác toàn diện: Áp dụng AI/ML dự báo sản xuất kinh doanh nâng cao.", RGBColor(100, 116, 139))
    ]
    for period, desc, col in stages:
        p = tf4_right.add_paragraph()
        p.text = "▶ " + period
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = col
        p.space_after = Pt(2)
        
        pd = tf4_right.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(11)
        pd.font.color.rgb = TEXT_DARK
        pd.space_after = Pt(8)

    s4.notes_slide.notes_text_frame.text = (
        "KỊCH BẢN NÓI:\n"
        "Thưa Ban Giám đốc, Tổng công ty đã đi trước một bước rất xa. Họ đã xong giai đoạn 1, xây xong kho tổng và phát lệnh "
        "bắt buộc trong 2026-2027 các đơn vị phải kết nối dữ liệu. Nói một cách ví von, TCT đã xây xong đại lộ cao tốc và mở sẵn làn đường "
        "cho Quảng Ngãi, nhiệm vụ của chúng ta là phải xây ngay con đường nhánh từ công ty mình đấu nối vào cao tốc đó."
    )

    # ==========================================
    # SLIDE 5: HIỆN TRẠNG TẠI PTSC QUẢNG NGÃI
    # ==========================================
    s5 = prs.slides.add_slide(blank_layout)
    add_header(s5, "4. HIỆN TRẠNG QUẢNG NGÃI: 3 'NÚT THẮT CỔ CHAI'", "NHỮNG KHÓ KHĂN THỰC TẾ CẦN GIẢI QUYẾT NGAY")

    issues = [
        ("NÚT THẮT 1: DỮ LIỆU PHÂN MẢNH", 
         "Hệ thống ứng dụng rời rạc", 
         [
             "Các phần mềm Kế toán, Quản lý Nhân sự, Vật tư, Dự án hoạt động độc lập.",
             "Cơ sở dữ liệu đặt trên các máy chủ riêng lẻ, chưa có cầu nối thông suốt.",
             "Vẫn phải dùng con người làm trung gian luân chuyển số liệu qua file Excel."
         ], Inches(0.8), ACCENT_GOLD),
        ("NÚT THẮT 2: LỆCH CHUẨN MASTER DATA", 
         "Chưa đồng bộ với 29 chuẩn TCT", 
         [
             "Mã danh mục của Quảng Ngãi (mã nhân viên, vật tư, đối tác...) chưa khớp chuẩn TCT.",
             "Nếu không 'Ánh xạ' (Mapping) lại thì dữ liệu đẩy lên TCT sẽ bị lỗi và từ chối tiếp nhận.",
             "Khối lượng rà soát nghiệp vụ rất lớn, cần sự tham gia của các phòng ban."
         ], Inches(4.8), ACCENT_GOLD),
        ("NÚT THẮT 3: THIẾU KỸ SƯ DATA CHUYÊN SÂU", 
         "Năng lực nội bộ chưa đáp ứng", 
         [
             "Phòng IT hiện tại làm rất tốt công tác quản trị hạ tầng, mạng và hỗ trợ người dùng.",
             "Tuy nhiên, thiếu hụt chuyên gia Data Engineering chuyên sâu về ETL/ELT, Fabric, API bảo mật.",
             "Nếu ép nội bộ tự nghiên cứu: Tốn 6-12 tháng, rủi ro lỗi kiến trúc và trễ hạn TCT."
         ], Inches(8.8), ACCENT)
    ]

    for title, sub, bullets, left_pos, border_col in issues:
        tf5 = add_card(s5, left_pos, Inches(1.5), Inches(3.7), Inches(5.2), title, sub, border_col)
        for b in bullets:
            p = tf5.add_paragraph()
            p.text = "• " + b
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT_DARK
            p.space_after = Pt(8)

    s5.notes_slide.notes_text_frame.text = (
        "KỊCH BẢN NÓI:\n"
        "Nhìn thẳng vào thực tế Quảng Ngãi, chúng ta có 3 nút thắt: 1 là phần mềm rời rạc, 2 là mã số liệu chưa khớp chuẩn TCT, "
        "và đặc biệt là nút thắt thứ 3: Anh em IT nhà mình chuyên về mạng, máy chủ và vận hành ứng dụng, chứ chưa từng làm chuyên sâu "
        "về Data Engineering đường ống lớn. Nếu để tự mày mò làm từ đầu, nguy cơ trễ hạn tiến độ giao nộp cho TCT là rất cao."
    )

    # ==========================================
    # SLIDE 6: MÔ HÌNH TRẠM TRUNG CHUYỂN DỮ LIỆU
    # ==========================================
    s6 = prs.slides.add_slide(blank_layout)
    add_header(s6, "5. KIẾN TRÚC ĐỀ XUẤT: TRẠM TRUNG CHUYỂN DỮ LIỆU (L3 SPOKE)", "MÔ HÌNH LUỒNG DỮ LIỆU TỪ QUẢNG NGÃI LÊN TỔNG CÔNG TY")

    # Left: Kiến trúc luồng
    tf6_left = add_card(s6, Inches(0.8), Inches(1.5), Inches(7.5), Inches(5.2), "Sơ đồ nguyên lý hoạt động")
    
    flow_steps = [
        ("BƯỚC 1: HỆ THỐNG NỘI BỘ (QUẢNG NGÃI)", "Phần mềm Kế toán + Quản lý Nhân sự + Vật tư + Hợp đồng."),
        ("BƯỚC 2: TRẠM TRUNG CHUYỂN (DATA GATEWAY)", "Cài đặt trên 1 máy chủ ảo tại Quảng Ngãi -> Tự động trích xuất, ánh xạ dữ liệu theo chuẩn 29 Master Data của TCT -> Mã hóa bảo mật."),
        ("BƯỚC 3: ĐƯỜNG TRUYỀN BẢO MẬT (VPN / SD-WAN)", "Dữ liệu được đẩy qua kênh truyền riêng biệt, đảm bảo an toàn tuyệt đối."),
        ("BƯỚC 4: WORKSPACE L3 (TỔNG CÔNG TY)", "Dữ liệu tiếp đất vào kho riêng của Quảng Ngãi trên đám mây TCT."),
        ("BƯỚC 5: BÁO CÁO ĐIỀU HÀNH (POWER BI)", "Tự động hiển thị lên bảng Dashboard điều hành cho Ban Giám đốc xem mỗi ngày.")
    ]
    for step_title, step_desc in flow_steps:
        p = tf6_left.add_paragraph()
        r1 = p.add_run()
        r1.text = "➜ " + step_title + ": "
        r1.font.bold = True
        r1.font.color.rgb = PRIMARY
        r2 = p.add_run()
        r2.text = step_desc
        r2.font.color.rgb = TEXT_DARK
        p.space_after = Pt(6)

    # Right: Điểm cốt lõi
    tf6_right = add_card(s6, Inches(8.7), Inches(1.5), Inches(3.8), Inches(5.2), "Ưu điểm của mô hình này", "GỌN NHẸ & AN TOÀN")
    advs = [
        ("Không động chạm dữ liệu gốc:", " Trạm trung chuyển chỉ 'đọc' bản sao, không gây ảnh hưởng hay chậm hệ thống kế toán đang chạy."),
        ("Bảo mật cao:", " Không phơi bày máy chủ nội bộ ra internet, đi qua đường truyền riêng VPN của TCT."),
        ("Chi phí thấp:", " Chỉ cần 1-2 máy chủ ảo (VM) cấu hình vừa phải đặt tại phòng máy chủ của đơn vị.")
    ]
    for at, ad in advs:
        p = tf6_right.add_paragraph()
        r = p.add_run()
        r.text = "★ " + at
        r.font.bold = True
        r.font.color.rgb = SECONDARY
        p.space_after = Pt(2)
        pd = tf6_right.add_paragraph()
        pd.text = ad
        pd.font.size = Pt(11)
        pd.font.color.rgb = TEXT_DARK
        pd.space_after = Pt(8)

    s6.notes_slide.notes_text_frame.text = (
        "KỊCH BẢN NÓI:\n"
        "Đây là mô hình kiến trúc kỹ thuật mà chúng ta sẽ làm. Ban Giám đốc hoàn toàn yên tâm: "
        "Trạm trung chuyển dữ liệu này chỉ hoạt động ở chế độ đọc, tuyệt đối không làm ảnh hưởng hay làm chậm phần mềm kế toán, nhân sự. "
        "Nó chỉ làm nhiệm vụ lấy bản sao dữ liệu, đóng gói theo chuẩn của TCT và gửi qua đường hầm bảo mật lên kho trên mây."
    )

    # ==========================================
    # SLIDE 7: PHÂN TÍCH PHƯƠNG ÁN: TỰ LÀM VS THUÊ NGOÀI
    # ==========================================
    s7 = prs.slides.add_slide(blank_layout)
    add_header(s7, "6. SO SÁNH PHƯƠNG ÁN: NỘI BỘ TỰ LÀM VS THUÊ NGOÀI", "CĂN CỨ ĐỂ ĐỀ XUẤT THUÊ ĐỐI TÁC CHUYÊN NGHIỆP")

    # Table comparison
    table_shape = s7.shapes.add_table(5, 3, Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.2))
    table = table_shape.table
    table.columns[0].width = Inches(2.3)
    table.columns[1].width = Inches(4.7)
    table.columns[2].width = Inches(4.733)

    # Header Row
    headers = ["Tiêu chí đánh giá", "Phương án 1: Tự làm 100% nội bộ", "Phương án 2: Kết hợp Thuê ngoài (ĐỀ XUẤT)"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY if i < 2 else SECONDARY
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    data = [
        ("Tiến độ & Thời gian",
         "• Rất chậm (mất 8 - 12 tháng do phải vừa làm vừa học hỏi công nghệ mới).\n• Rất dễ trễ hạn deadline 2026 của TCT.",
         "• Rất nhanh (3 - 4 tháng là xong toàn bộ luồng kết nối).\n• Đảm bảo 100% đúng tiến độ Nghị quyết TCT."),
        ("Chất lượng & Chuẩn kỹ thuật",
         "• Tiềm ẩn rủi ro lỗi pipeline, nghẽn mạng hoặc bảo mật do thiếu kinh nghiệm thực chiến.",
         "• Đạt chuẩn doanh nghiệp, tối ưu luồng dữ liệu, chuẩn bảo mật của Microsoft & TCT ban hành."),
        ("Nguồn lực & Rủi ro",
         "• Nhân sự IT quá tải, bỏ bê các nhiệm vụ hỗ trợ vận hành thường nhật.\n• Nếu nhân sự nghỉ việc, dự án bị đứt gãy hoàn toàn.",
         "• Có cam kết SLA, chuyển giao tài liệu chuẩn mực, bảo hành bảo trì dài hạn từ đơn vị chuyên nghiệp."),
        ("Vai trò của Nhân sự Quảng Ngãi",
         "• Bị động, lúng túng trong khâu kỹ thuật đường ống dữ liệu phức tạp.",
         "• Chủ động: Nội bộ tập trung vào làm chủ dữ liệu nghiệp vụ, giám sát chất lượng và tiếp nhận chuyển giao.")
    ]

    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255) if col_idx > 0 else RGBColor(241, 245, 249)
            p = cell.text_frame.paragraphs[0]
            p.text = text
            p.font.size = Pt(11)
            p.font.bold = (col_idx == 0)
            p.font.color.rgb = PRIMARY if col_idx == 0 else TEXT_DARK

    s7.notes_slide.notes_text_frame.text = (
        "KỊCH BẢN NÓI (TRỌNG TÂM THUYẾT PHỤC):\n"
        "Thưa Ban Giám đốc, đây là slide then chốt trong báo cáo hôm nay. Chúng tôi đã cân nhắc rất kỹ bài toán tự làm và thuê ngoài. "
        "Nếu tự làm, công ty có thể tiết kiệm một khoản chi phí ban đầu, nhưng cái giá phải trả là thời gian kéo dài cả năm và nguy cơ "
        "bị TCT phê bình vì chậm trễ. "
        "Thuê ngoài ở đây KHÔNG PHẢI là giao trắng cho người ta, mà là thuê 'thợ chuyên nghiệp làm đường ống', còn cán bộ Quảng Ngãi "
        "giữ vai trò kiểm soát nguồn nước, giám sát và tiếp nhận toàn bộ công nghệ."
    )

    # ==========================================
    # SLIDE 8: PHẠM VI THUÊ NGOÀI & LỘ TRÌNH TRIỂN KHAI
    # ==========================================
    s8 = prs.slides.add_slide(blank_layout)
    add_header(s8, "7. PHẠM VI CÔNG VIỆC & LỘ TRÌNH 3 GIAI ĐOẠN", "KẾ HOẠCH HÀNH ĐỘNG CỤ THỂ, KIỂM SOÁT CHẶT CHẼ")

    phases = [
        ("GIAI ĐOẠN 1 (Tháng 1 - 2)", 
         "NỘI BỘ CHỦ TRÌ (0 Đồng)", 
         [
             "Thành lập Tổ công tác Data Platform nội bộ.",
             "Khảo sát toàn bộ cơ sở dữ liệu các phần mềm nghiệp vụ.",
             "Thực hiện Ánh xạ (Mapping) dữ liệu nội bộ với 29 danh mục Master Data của TCT.",
             "Lập bảng yêu cầu kỹ thuật chi tiết để chuẩn bị mời chào giá thuê ngoài."
         ], Inches(0.8), SECONDARY),
        ("GIAI ĐOẠN 2 (Tháng 3 - 5)", 
         "THUÊ ĐỐI TÁC TRIỂN KHAI (Outsource)", 
         [
             "Cài đặt và cấu hình Trạm trung chuyển dữ liệu trên máy chủ ảo.",
             "Viết các đường ống ETL tự động trích xuất, làm sạch và đẩy dữ liệu lên TCT.",
             "Cấu hình đường truyền bảo mật VPN/SD-WAN kết nối về Hub TCT.",
             "Kiểm thử an toàn thông tin và nghiệm thử thông luồng dữ liệu thành công."
         ], Inches(4.8), ACCENT),
        ("GIAI ĐOẠN 3 (Tháng 6 trở đi)", 
         "BÀN GIAO & KHAI THÁC", 
         [
             "Đối tác bàn giao toàn bộ mã nguồn, tài liệu hướng dẫn vận hành.",
             "Đào tạo chuyển giao công nghệ cho đội ngũ IT Quảng Ngãi làm chủ 100%.",
             "Cùng phòng ban xây dựng Dashboard báo cáo quản trị phục vụ Ban Giám đốc.",
             "Bảo hành, hỗ trợ kỹ thuật định kỳ."
         ], Inches(8.8), RGBColor(39, 174, 96))
    ]

    for title, sub, bullets, left_pos, color_bar in phases:
        tf8 = add_card(s8, left_pos, Inches(1.5), Inches(3.7), Inches(5.2), title, sub, color_bar)
        for b in bullets:
            p = tf8.add_paragraph()
            p.text = "• " + b
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT_DARK
            p.space_after = Pt(8)

    s8.notes_slide.notes_text_frame.text = (
        "KỊCH BẢN NÓI:\n"
        "Lộ trình được chia làm 3 giai đoạn rất rành mạch: "
        "Giai đoạn 1 nội bộ tự làm để hiểu cặn kẽ số liệu và tiết kiệm chi phí; "
        "Giai đoạn 2 mới giải ngân thuê đối tác kỹ thuật dựng trạm và kiểm thử kết nối; "
        "Giai đoạn 3 là nghiệm thu, nhận bàn giao toàn bộ tài liệu để mình tự quản lý vận hành lâu dài."
    )

    # ==========================================
    # SLIDE 9: KIẾN NGHỊ BAN GIÁM ĐỐC PHÊ DUYỆT
    # ==========================================
    s9 = prs.slides.add_slide(blank_layout)
    add_header(s9, "8. KIẾN NGHỊ & ĐỀ XUẤT BAN GIÁM ĐỐC PHÊ DUYỆT", "CÁC NỘI DUNG CẦN THÔNG QUA ĐỂ BẮT TAY THỰC HIỆN")

    tf9 = add_card(s9, Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.2), "3 Kiến nghị trọng tâm kính trình Ban Giám đốc", "CĂN CỨ VÀ HÀNH ĐỘNG TIẾP THEO")

    proposals = [
        ("1. PHÊ DUYỆT CHỦ TRƯƠNG TRIỂN KHAI & THÀNH LẬP TỔ CÔNG TÁC DATA PLATFORM:",
         "Chấp thuận thành lập Tổ công tác gồm Phòng IT (chủ trì kỹ thuật) và các cán bộ phụ trách dữ liệu (Key users) của Phòng Kế toán, Phòng Thương mại/Vật tư, Phòng HC-NS để tiến hành khảo sát và chuẩn hóa dữ liệu Master Data ngay trong Tháng tới."),
        ("2. PHÊ DUYỆT CẤP PHÁT TÀI NGUYÊN HẠ TẦNG NỘI BỘ:",
         "Cho phép phòng IT khởi tạo 1 - 2 Máy chủ ảo (Virtual Machine) trên hạ tầng server hiện hữu của đơn vị và phối hợp cùng TCT thiết lập đường truyền mạng bảo mật (VPN) phục vụ làm Trạm trung chuyển dữ liệu."),
        ("3. PHÊ DUYỆT CHỦ TRƯƠNG THUÊ NGOÀI VÀ CHO PHÉP LẬP DỰ TOÁN CHI PHÍ:",
         "Cho phép tổ công tác liên hệ, khảo sát thị trường và lấy báo giá cạnh tranh từ các đơn vị cung cấp giải pháp Data Engineering uy tín để lập khái toán/dự toán chi tiết trình Ban Giám đốc xem xét phê duyệt trước khi ký kết hợp đồng.")
    ]

    for title, detail in proposals:
        p = tf9.add_paragraph()
        p.text = title
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = PRIMARY
        p.space_after = Pt(3)

        pd = tf9.add_paragraph()
        pd.text = detail
        pd.font.size = Pt(11)
        pd.font.color.rgb = TEXT_DARK
        pd.space_after = Pt(12)

    # Callout bottom box
    callout = s9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(5.8), Inches(10.9), Inches(0.65))
    callout.fill.solid()
    callout.fill.fore_color.rgb = RGBColor(254, 243, 199)
    callout.line.color.rgb = ACCENT_GOLD
    tf_c = callout.text_frame
    p_c = tf_c.paragraphs[0]
    p_c.text = "Kính trình Ban Giám đốc xem xét, cho ý kiến chỉ đạo để Đơn vị kịp thời bắt nhịp tiến độ chung của TCT!"
    p_c.font.bold = True
    p_c.font.size = Pt(12)
    p_c.font.color.rgb = RGBColor(146, 64, 14)
    p_c.alignment = PP_ALIGN.CENTER

    s9.notes_slide.notes_text_frame.text = (
        "KỊCH BẢN NÓI:\n"
        "Kính thưa Ban Giám đốc, để không bị trễ nhịp với Tổng công ty, phòng IT kính đề xuất Ban Giám đốc thông qua 3 nội dung: "
        "Một là cho phép thành lập Tổ công tác nội bộ; Hai là cấp tài nguyên máy chủ ảo; và Ba là phê duyệt chủ trương cho phép chúng tôi "
        "đi khảo sát báo giá các đơn vị ngoài để lập dự toán trình Sếp. "
        "Rất mong nhận được sự ủng hộ và chỉ đạo của Ban Giám đốc. Chúng tôi xin trân trọng cảm ơn!"
    )

    output_path = r"d:\My Profiles\DataPlatform\bao_cao_dataplatform_ptsc_qn.pptx"
    prs.save(output_path)
    print(f"SUCCESS: Slide created at {output_path}")

if __name__ == "__main__":
    create_deck()
