"""
BÁO CÁO DATA PLATFORM PTSC QUẢNG NGÃI - VERSION 5
Phong cách: Corporate Executive — nghiêm túc, chuyên nghiệp, chuẩn mực doanh nghiệp dầu khí.
Tích hợp đầy đủ thông tin thực tế từ tài liệu Tổng công ty:
- Nền tảng Hybrid IMIP (MinIO On-prem + Fabric Cloud 20TB, MDM, ESB, 50 API).
- Quy chế Quản trị Dữ liệu 5 Cấp & Quyền tự quyết của Đơn vị (Data Owner Cấp 3).
- Nguyên tắc Chủ quyền dữ liệu: Dữ liệu nhạy cảm ở lại On-prem, chỉ chuyển dữ liệu tổng hợp.
- Cơ chế tài chính Phương án 3 của TCT & Cấu trúc chi phí (để trống phần tiền chờ NCC báo giá).
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Corporate Color Palette
NAVY       = RGBColor(12, 35, 68)
NAVY_MID   = RGBColor(24, 59, 108)
BLUE       = RGBColor(30, 100, 190)
BLUE_SOFT  = RGBColor(230, 240, 252)
GREEN      = RGBColor(21, 128, 80)
GREEN_SOFT = RGBColor(232, 248, 239)
RED        = RGBColor(180, 40, 40)
RED_SOFT   = RGBColor(252, 235, 235)
GOLD       = RGBColor(180, 115, 0)
GOLD_SOFT  = RGBColor(255, 245, 225)
WHITE      = RGBColor(255, 255, 255)
BLACK      = RGBColor(28, 35, 48)
GRAY       = RGBColor(108, 118, 132)
GRAY_LIGHT = RGBColor(243, 245, 248)
GRAY_MED   = RGBColor(200, 208, 218)
TEAL       = RGBColor(13, 148, 136)

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BL = prs.slide_layouts[6]

def shape(slide, l, t, w, h, fill, line_c=None, line_w=None, rounded=False):
    sh_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    r = slide.shapes.add_shape(sh_type, l, t, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = fill
    if line_c:
        r.line.color.rgb = line_c
        if line_w: r.line.width = line_w
    else:
        r.line.fill.background()
    return r

def text(slide, l, t, w, h, content, size=14, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = content
    p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = color; p.alignment = align
    return tf

def header(slide, title, subtitle="", slide_num=""):
    shape(slide, 0, 0, W, Inches(1.05), NAVY)
    shape(slide, 0, Inches(1.05), W, Inches(0.035), GOLD)
    if subtitle:
        text(slide, Inches(0.8), Inches(0.1), Inches(10), Inches(0.25), subtitle.upper(), 9, True, RGBColor(155, 185, 220))
    text(slide, Inches(0.8), Inches(0.38), Inches(11), Inches(0.55), title, 19, True, WHITE)
    if slide_num:
        text(slide, Inches(12.0), Inches(0.38), Inches(0.8), Inches(0.5), slide_num, 12, True, RGBColor(155, 185, 220), PP_ALIGN.RIGHT)
    shape(slide, 0, Inches(7.2), W, Inches(0.008), GRAY_MED)
    text(slide, Inches(0.8), Inches(7.22), Inches(11.7), Inches(0.25),
         "PTSC Quảng Ngãi  |  Báo cáo Data Platform & Trục tích hợp  |  Thực hiện Chỉ đạo số 9 BGĐ", 8, color=GRAY)

def numbered_label(slide, l, t, num, color=BLUE):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, Inches(0.4), Inches(0.4))
    c.fill.solid(); c.fill.fore_color.rgb = color; c.line.fill.background()
    tf = c.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = str(num)
    p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER

def arrow_shape(slide, l, t, direction="right", w=Inches(0.5), h=Inches(0.3)):
    sh_type = MSO_SHAPE.RIGHT_ARROW if direction == "right" else MSO_SHAPE.DOWN_ARROW
    if direction == "down": w, h = h, w
    ar = slide.shapes.add_shape(sh_type, l, t, w, h)
    ar.fill.solid(); ar.fill.fore_color.rgb = GOLD; ar.line.fill.background()

def note(slide, txt_content):
    slide.notes_slide.notes_text_frame.text = txt_content

# =========================================================================
# SLIDE 1: COVER
# =========================================================================
s1 = prs.slides.add_slide(BL)
shape(s1, 0, 0, W, H, NAVY)
shape(s1, 0, Inches(1.4), W, Inches(0.04), GOLD)
shape(s1, 0, Inches(5.8), W, Inches(0.015), RGBColor(60, 90, 140))

text(s1, Inches(1.0), Inches(1.8), Inches(11), Inches(0.4),
     "BÁO CÁO BAN GIÁM ĐỐC  —  THỰC HIỆN CHỈ ĐẠO MỤC 9", 13, True, RGBColor(180, 200, 225))
text(s1, Inches(1.0), Inches(2.5), Inches(11), Inches(1.6),
     "PHƯƠNG ÁN TRIỂN KHAI DATA PLATFORM\nVÀ TRỤC TÍCH HỢP DỮ LIỆU CÔNG TY – TCT", 30, True, WHITE)
shape(s1, Inches(1.0), Inches(4.4), Inches(2.5), Inches(0.04), GOLD)

scope_items = [
    "So sánh 2 Phương án: Thuê TCT  vs  Phát triển riêng đồng bộ TCT (có NCC tư vấn)",
    "Hạ tầng Nền tảng TCT đã có, Quy chế Quản trị 5 cấp và Quyền tự quyết dữ liệu của Đơn vị",
    "Mô hình Kiến trúc Trục tích hợp Hybrid (MinIO + Fabric) và Cơ chế Chi phí đầu tư"
]
for i, item in enumerate(scope_items):
    text(s1, Inches(1.0), Inches(4.7 + i * 0.32), Inches(11), Inches(0.3), "—  " + item, 12, color=RGBColor(200, 215, 235))

text(s1, Inches(1.0), Inches(6.1), Inches(11), Inches(0.3),
     "Đơn vị thực hiện:  Tổ Công tác CĐS & CNTT  –  PTSC Quảng Ngãi", 11.5, color=WHITE)
text(s1, Inches(1.0), Inches(6.45), Inches(11), Inches(0.3),
     "Kính trình:  Ban Giám đốc Công ty PTSC Quảng Ngãi", 11.5, True, RGBColor(253, 224, 71))

note(s1, "Kính thưa Ban Giám đốc, hôm nay bộ phận CNTT xin báo cáo toàn diện 4 nội dung "
     "theo đúng chỉ đạo số 9 của Ban Giám đốc, có đối chiếu trực tiếp với các tài liệu kỹ thuật và quy chế quản trị dữ liệu của Tổng công ty.")

# =========================================================================
# SLIDE 2: BỐI CẢNH — TẠI SAO PHẢI LÀM TRONG NĂM 2026?
# =========================================================================
s2 = prs.slides.add_slide(BL)
header(s2, "BỐI CẢNH BẮT BUỘC: NGHỊ QUYẾT TCT VÀ MỐC HẠN NĂM 2026", "CĂN CỨ CHIẾN LƯỢC — TÍNH CẤP BÁCH", "01")

phases = [
    ("GIAI ĐOẠN 1 (2024 – 2025)", "TCT ĐÃ XONG NỀN MÓNG", GREEN, GREEN_SOFT,
     ["TCT nghiệm thu Nền tảng Hybrid IMIP Data Platform",
      "Ban hành Quy chế Quản trị Dữ liệu & 29 Master Data",
      "Phạm vi: Mới chỉ áp dụng tại Cơ quan TCT, CHƯA CÓ ĐƠN VỊ THÀNH VIÊN"]),
    ("GIAI ĐOẠN 2 (2026 – 2027)", "BẮT BUỘC KẾT NỐI ĐƠN VỊ", GOLD, GOLD_SOFT,
     ["Đến lượt các Đơn vị thành viên (Quảng Ngãi bắt đầu)",
      "Chuẩn hóa dữ liệu nội bộ khớp 29 danh mục TCT",
      "Xây Trạm trung chuyển kết nối vào hạ tầng TCT Hub",
      "Tiêu chí KPI đánh giá mức độ CĐS của Đơn vị!"]),
    ("GIAI ĐOẠN 3 (2028+)", "KHAI THÁC TOÀN DIỆN", GRAY, GRAY_LIGHT,
     ["Phân tích dữ liệu chuyên sâu toàn Tổng công ty",
      "Ứng dụng AI/Machine Learning trong SXKD",
      "Kết nối IT và OT (vận hành thiết bị cảng/xưởng)"]),
]
for i, (period, status, col, bg_c, items) in enumerate(phases):
    x = Inches(0.8 + i * 4.1)
    shape(s2, x, Inches(1.5), Inches(3.7), Inches(0.65), col, rounded=True)
    text(s2, x, Inches(1.53), Inches(3.7), Inches(0.3), period, 14, True, WHITE, PP_ALIGN.CENTER)
    text(s2, x, Inches(1.82), Inches(3.7), Inches(0.25), status, 9, True, WHITE, PP_ALIGN.CENTER)
    shape(s2, x, Inches(2.25), Inches(3.7), Inches(2.6), bg_c, col, Pt(1), rounded=True)
    for j, item in enumerate(items):
        text(s2, x + Inches(0.25), Inches(2.45 + j * 0.58), Inches(3.2), Inches(0.55),
             "—  " + item, 11.5, color=BLACK)
    if i < 2:
        arrow_shape(s2, Inches(4.55 + i * 4.1), Inches(2.9))

shape(s2, Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.5), RED_SOFT, RED, Pt(1.5), rounded=True)
shape(s2, Inches(0.84), Inches(5.32), Inches(0.08), Inches(1.42), RED)
text(s2, Inches(1.2), Inches(5.4), Inches(10.8), Inches(0.35), "Tại sao PTSC Quảng Ngãi phải hành động ngay trong năm 2026?", 13.5, True, RED)
tf_urg = text(s2, Inches(1.2), Inches(5.82), Inches(10.8), Inches(0.35),
    "1. Nghị quyết 10/NQ-HĐQT-PTSC: Nền tảng Dữ liệu là trụ cột bắt buộc cho toàn bộ các đơn vị thành viên.", 11, color=BLACK)
p = tf_urg.add_paragraph()
p.text = "2. SOW Giai đoạn 1 của TCT ghi rõ chưa có dữ liệu đơn vị; 2026 là năm bản lề TCT yêu cầu mở rộng kết nối từ các đơn vị thành viên."
p.font.size = Pt(11); p.font.color.rgb = BLACK
p2 = tf_urg.add_paragraph()
p2.text = "3. Tuân thủ pháp lý mới: Luật Dữ liệu (01/7/2025) và Luật Bảo vệ DLCN 91/2025/QH15 (01/01/2026) bắt buộc dữ liệu phải có nguồn gốc và kiểm soát."
p2.font.size = Pt(11); p2.font.color.rgb = BLACK

note(s2, "TCT đã hoàn thành Giai đoạn 1 tại Văn phòng TCT. Năm 2026-2027 là thời hạn bắt buộc mở rộng "
     "đến các đơn vị thành viên như Quảng Ngãi. Nếu chậm trễ sẽ bị đánh giá KPI CĐS.")

# =========================================================================
# SLIDE 3: HIỆN TRẠNG TẠI QUẢNG NGÃI
# =========================================================================
s3 = prs.slides.add_slide(BL)
header(s3, "HIỆN TRẠNG DỮ LIỆU TẠI QUẢNG NGÃI: 3 NÚT THẮT CẦN THÁO GỠ", "THỰC TRẠNG — NGUYÊN NHÂN CHẬM TRỄ", "02")

silos = ["Phần mềm\nKẾ TOÁN", "Phần mềm\nNHÂN SỰ", "Phần mềm\nVẬT TƯ", "Quản lý\nDỰ ÁN"]
for i, name in enumerate(silos):
    x = Inches(0.8 + i * 2.95)
    shape(s3, x, Inches(1.5), Inches(2.55), Inches(1.25), WHITE, GRAY_MED, Pt(0.8), rounded=True)
    shape(s3, x + Inches(0.04), Inches(1.52), Inches(0.08), Inches(1.21), BLUE)
    text(s3, x + Inches(0.25), Inches(1.6), Inches(2.1), Inches(1.0), name, 12.5, True, NAVY, PP_ALIGN.CENTER)
    if i < 3:
        text(s3, Inches(3.1 + i * 2.95), Inches(1.8), Inches(0.5), Inches(0.4), "✕", 18, True, RED, PP_ALIGN.CENTER)

text(s3, Inches(0.8), Inches(3.15), Inches(11.7), Inches(0.35),
     "3 Nút thắt lớn cản trở công tác quản trị điều hành của Ban Giám đốc:", 13.5, True, RED)

knots = [
    ("DỮ LIỆU PHÂN MẢNH", RED, RED_SOFT,
     "4 phần mềm hoạt động riêng rẽ.\nChuyển số liệu bằng file Excel thủ công.\nMất 3–5 ngày tổng hợp báo cáo."),
    ("LỆCH 29 CHUẨN MASTER DATA", GOLD, GOLD_SOFT,
     "Mã nhân sự, vật tư, đối tác của QN\nchưa khớp 29 danh mục chuẩn của TCT.\nNếu đẩy lên ngay sẽ bị lỗi từ chối."),
    ("THIẾU KỸ SƯ DATA PIPELINE", BLUE, BLUE_SOFT,
     "IT nội bộ giỏi mạng, máy chủ và hỗ trợ,\nnhưng thiếu chuyên gia Data Engineering.\nTự mày mò: rủi ro trễ tiến độ 2026."),
]
for i, (title, col, bg_c, desc) in enumerate(knots):
    x = Inches(0.8 + i * 4.1)
    shape(s3, x, Inches(3.7), Inches(3.7), Inches(3.2), bg_c, col, Pt(1.5), rounded=True)
    shape(s3, x + Inches(0.04), Inches(3.72), Inches(0.08), Inches(3.14), col)
    numbered_label(s3, x + Inches(0.3), Inches(3.9), i + 1, col)
    text(s3, x + Inches(0.85), Inches(3.9), Inches(2.6), Inches(0.4), title, 13, True, col)
    shape(s3, x + Inches(0.3), Inches(4.5), Inches(3.1), Inches(0.01), GRAY_MED)
    text(s3, x + Inches(0.3), Inches(4.7), Inches(3.1), Inches(2.0), desc, 11.5, color=BLACK)

note(s3, "Hiện trạng QN: phần mềm cát cứ, chưa đồng bộ 29 danh mục Master Data, "
     "và IT nội bộ thiếu kinh nghiệm viết đường ống dữ liệu.")

# =========================================================================
# SLIDE 4: DATA PLATFORM LÀ GÌ?
# =========================================================================
s4 = prs.slides.add_slide(BL)
header(s4, "DATA PLATFORM: TỪ \"ỐC ĐẢO THỦ CÔNG\" ĐẾN \"ĐƯỜNG ỐNG TỰ ĐỘNG\"", "BẢN CHẤT GIẢI PHÁP — DỄ HIỂU CHO LÃNH ĐẠO", "03")

shape(s4, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.3), RED_SOFT, RED, Pt(1.5), rounded=True)
shape(s4, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5), RED)
text(s4, Inches(0.8), Inches(1.52), Inches(5.5), Inches(0.45), "TRƯỚC ĐÂY  —  Mô hình ốc đảo, phụ thuộc con người", 12, True, WHITE, PP_ALIGN.CENTER)

before_items = [
    "Dữ liệu nằm chết trong từng phần mềm riêng lẻ",
    "Nhân viên phải xuất Excel, copy-paste ghép file thủ công",
    "Mất 3–5 ngày mới có số liệu báo cáo tuần/tháng",
    "Kế toán báo một số, Vật tư báo số khác (lệch nguồn)",
    "Ban Giám đốc không có công cụ nhìn toàn cảnh tức thời"
]
for i, item in enumerate(before_items):
    text(s4, Inches(1.3), Inches(2.3 + i * 0.75), Inches(4.6), Inches(0.65),
         "✕  " + item, 11.5, color=RED if i == 0 else BLACK)

arrow_shape(s4, Inches(6.5), Inches(3.8))

shape(s4, Inches(7.2), Inches(1.5), Inches(5.3), Inches(5.3), GREEN_SOFT, GREEN, Pt(1.5), rounded=True)
shape(s4, Inches(7.2), Inches(1.5), Inches(5.3), Inches(0.5), GREEN)
text(s4, Inches(7.2), Inches(1.52), Inches(5.3), Inches(0.45), "SAU ĐÓ  —  Trục dữ liệu tự động (Data Platform)", 12, True, WHITE, PP_ALIGN.CENTER)

after_items = [
    "Dữ liệu tự động hút về một trạm xử lý tập trung",
    "Hệ thống tự làm sạch, quy đổi mã theo 29 chuẩn TCT",
    "Số liệu tự động cập nhật hàng ngày (Near real-time)",
    "Một nguồn sự thật duy nhất (Single Source of Truth)",
    "Ban Giám đốc xem Dashboard trực quan trên PC / Mobile"
]
for i, item in enumerate(after_items):
    text(s4, Inches(7.7), Inches(2.3 + i * 0.75), Inches(4.4), Inches(0.65),
         "✔  " + item, 11.5, color=GREEN if i == 0 else BLACK)

note(s4, "Data Platform giống hệ thống đường ống nước tự động. Gom dữ liệu, lọc sạch và tự bơm lên vòi "
     "của Ban Giám đốc mỗi sáng mà không cần người xách xô Excel thủ công.")

# =========================================================================
# SLIDE 5: VỊ THẾ QUẢNG NGÃI TRONG MÔ HÌNH HUB-SPOKE CỦA TCT
# =========================================================================
s5 = prs.slides.add_slide(BL)
header(s5, "VỊ THẾ QUẢNG NGÃI: LEVEL 3 TRONG MÔ HÌNH HUB-SPOKE CỦA TCT", "QUY HOẠCH CỦA TCT — 4 LEVEL TRIỂN KHAI", "04")

level_cards = [
    ("L1: CHI NHÁNH", "Vận hành như 1 ban của TCT", "Không có hạ tầng riêng, không dựng Spoke.", GRAY, GRAY_LIGHT),
    ("L2: ĐƠN VỊ NHỎ", "Tenant trên Hub TCT", "Chỉ dùng Agent hoặc API kéo dữ liệu, không có máy chủ.", GRAY, GRAY_LIGHT),
    ("L3: ĐƠN VỊ LỚN (PTSC QUẢNG NGÃI)", "Spoke trung chuyển dữ liệu",
     "• Được cấp phân vùng Workspace L3 riêng trên Hub TCT.\n• KHÔNG CẦN mua máy chủ dHCI đắt tiền tại đơn vị.\n• Tận dụng máy chủ hiện có để dựng Trạm trung chuyển.\n• Tự chủ dữ liệu nghiệp vụ, đồng bộ chuẩn lên TCT.",
     GOLD, GOLD_SOFT),
    ("L4: ĐƠN VỊ ĐẶC BIỆT LỚN", "Cụm dHCI riêng (PTSC M&C)", "Tự đầu tư cụm máy chủ dHCI riêng, tự vận hành hoàn toàn.", BLUE, BLUE_SOFT),
]
for i, (lvl_t, lvl_sub, lvl_desc, col, bg_c) in enumerate(level_cards):
    x = Inches(0.8 + i * 2.95)
    is_qngai = (i == 2)
    w_card = Inches(2.8)
    shape(s5, x, Inches(1.5), w_card, Inches(4.7), bg_c, col, Pt(2 if is_qngai else 0.8), rounded=True)
    shape(s5, x + Inches(0.04), Inches(1.52), Inches(0.08), Inches(4.64), col)
    text(s5, x + Inches(0.25), Inches(1.7), Inches(2.4), Inches(0.35), lvl_t, 12, True, col)
    text(s5, x + Inches(0.25), Inches(2.05), Inches(2.4), Inches(0.3), lvl_sub, 10, True, NAVY)
    shape(s5, x + Inches(0.25), Inches(2.4), Inches(2.3), Inches(0.01), GRAY_MED)
    text(s5, x + Inches(0.2), Inches(2.55), Inches(2.4), Inches(3.4), lvl_desc, 11 if is_qngai else 10.5, color=BLACK)

shape(s5, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.65), BLUE_SOFT, BLUE, Pt(1), rounded=True)
shape(s5, Inches(0.84), Inches(6.42), Inches(0.07), Inches(0.58), BLUE)
text(s5, Inches(1.1), Inches(6.45), Inches(11.2), Inches(0.55),
     "KẾT LUẬN CỦA TCT: PTSC Quảng Ngãi xếp cấp Level 3 — Tiết kiệm tối đa ngân sách vì KHÔNG phải mua máy chủ lưu trữ dHCI tiền tỷ, chỉ cần xây dựng Trạm trung chuyển dữ liệu tại chỗ.",
     11.5, True, NAVY, PP_ALIGN.CENTER)

note(s5, "TCT phân 4 Level. Quảng Ngãi thuộc Level 3 (Đơn vị lớn có trạm trung chuyển riêng). "
     "Chúng ta không phải đầu tư cụm dHCI hàng tỷ đồng như L4 (PTSC M&C).")

# =========================================================================
# SLIDE 6: [MỚI] TỔNG CÔNG TY ĐÃ XÂY DỰNG NHỮNG GÌ?
# =========================================================================
s6 = prs.slides.add_slide(BL)
header(s6, "TỔNG CÔNG TY ĐÃ XÂY DỰNG NHỮNG GÌ VÀ SẴN SÀNG RA SAO?", "NỀN MÓNG TCT ĐÃ NGHIỆM THU — CĂN CỨ TÀI LIỆU GIAI ĐOẠN 1", "05")

tct_assets = [
    ("HẠ TẦNG CLOUD (AZURE)", BLUE, BLUE_SOFT, [
        "Microsoft Fabric / OneLake: Khởi điểm 20 TB dung lượng phân tích tốc độ cao.",
        "Microsoft Purview: Quản trị danh mục và gán nhãn an toàn dữ liệu.",
        "Power BI Enterprise: Phục vụ trực quan hóa Dashboard toàn tổng.",
        "Cấp sẵn Workspace L3 riêng cho Quảng Ngãi đón nhận dữ liệu."
    ]),
    ("HẠ TẦNG ON-PREMISE (DATACENTER TCT)", TEAL, RGBColor(204, 251, 241), [
        "Cụm MinIO Lakehouse: Lưu trữ dữ liệu nội bộ on-prem an toàn.",
        "Trục tích hợp doanh nghiệp (ESB): Điều phối dữ liệu tập trung.",
        "Hệ thống MDM Golden Record: Chuẩn hóa dữ liệu chủ toàn PTSC.",
        "SIEM & SOC: Giám sát an ninh và ghi nhận nhật ký (Audit log) 24/7."
    ]),
    ("BỘ TIÊU CHUẨN ĐÃ HOÀN TẤT", GREEN, GREEN_SOFT, [
        "Bộ 29 Danh mục Master Data: Chuẩn hóa mã nhân sự, vật tư, đối tác, dự án...",
        "50 API Chuẩn: Đã xây dựng sẵn để thu thập và tích hợp dữ liệu.",
        "Đã kết nối thành công 8 phần mềm nguồn và 35 quy trình nội bộ TCT.",
        "Sẵn sàng mở rộng đấu nối cho các đơn vị thành viên trong năm 2026."
    ]),
]
for i, (title, col, bg_c, items) in enumerate(tct_assets):
    x = Inches(0.8 + i * 4.1)
    shape(s6, x, Inches(1.5), Inches(3.7), Inches(5.3), bg_c, col, Pt(1.5), rounded=True)
    shape(s6, x + Inches(0.04), Inches(1.52), Inches(0.08), Inches(5.24), col)
    text(s6, x + Inches(0.25), Inches(1.7), Inches(3.2), Inches(0.55), title, 12, True, col)
    shape(s6, x + Inches(0.25), Inches(2.3), Inches(3.1), Inches(0.01), GRAY_MED)
    for j, item in enumerate(items):
        text(s6, x + Inches(0.25), Inches(2.45 + j * 1.15), Inches(3.2), Inches(1.05),
             "✔  " + item, 11, color=BLACK)

note(s6, "TCT đã xong nền tảng Hybrid IMIP: Cloud có Fabric OneLake 20TB, On-prem có MinIO, ESB, MDM và 50 API chuẩn. "
     "TCT đã sẵn sàng hạ tầng, chỉ chờ đơn vị thành viên đấu nối.")

# =========================================================================
# SLIDE 7: [MỚI] QUY CHẾ QUẢN TRỊ DỮ LIỆU & CHỦ QUYỀN CỦA QUẢNG NGÃI
# =========================================================================
s7 = prs.slides.add_slide(BL)
header(s7, "QUY CHẾ QUẢN TRỊ & CHỦ QUYỀN DỮ LIỆU CỦA QUẢNG NGÃI", "QUẢN TRỊ 5 CẤP THEO CHUẨN PETROVIETNAM — TRẢ LỜI CÂU HỎI BAN GIÁM ĐỐC", "06")

levels_gov = [
    ("CẤP 1", "Hội đồng Quản trị Dữ liệu TCT", "Phê duyệt chính sách, tiêu chuẩn toàn tổng"),
    ("CẤP 2", "Hội đồng Dữ liệu khối / chuyên ngành", "Điều phối và đồng bộ tiêu chuẩn liên miền"),
    ("CẤP 3", "CHỦ QUẢN DỮ LIỆU — LÃNH ĐẠO ĐƠN VỊ (QUẢNG NGÃI)",
     "★ SỞ HỮU NGHIỆP VỤ & QUYỀN QUYẾT ĐỊNH (DATA OWNER):\n"
     "• Phê duyệt phân loại dữ liệu của đơn vị.\n"
     "• QUYẾT ĐỊNH DỮ LIỆU NÀO ĐƯỢC CHIA SẺ, DỮ LIỆU NÀO Ở LẠI NỘI BỘ.\n"
     "• Trách nhiệm giải trình cao nhất đặt tại cấp này!"),
    ("CẤP 4", "Quản trị miền dữ liệu (Data Stewards)", "Thực thi ánh xạ, từ điển và làm sạch dữ liệu"),
    ("CẤP 5", "Ban NCPT&CĐS TCT (Đơn vị vận hành)", "Chỉ vận hành kỹ thuật nền tảng — KHÔNG SỞ HỮU DỮ LIỆU"),
]
y_start = Inches(1.4)
for i, (lvl, title, role) in enumerate(levels_gov):
    is_c3 = (i == 2)
    h_box = Inches(1.6) if is_c3 else Inches(0.9)
    y = y_start if i == 0 else y_prev + h_prev + Inches(0.1)
    y_prev, h_prev = y, h_box
    
    col = GOLD if is_c3 else (NAVY if i < 2 else GRAY)
    bg_c = GOLD_SOFT if is_c3 else (WHITE if i < 2 else GRAY_LIGHT)
    
    shape(s7, Inches(0.8), y, Inches(11.733), h_box, bg_c, col, Pt(2 if is_c3 else 0.8), rounded=True)
    shape(s7, Inches(0.84), y + Inches(0.02), Inches(0.08), h_box - Inches(0.04), col)
    
    numbered_label(s7, Inches(1.1), y + (Inches(0.6) if is_c3 else Inches(0.25)), i + 1, col)
    text(s7, Inches(1.65), y + Inches(0.12), Inches(9.8), Inches(0.35), title, 12 if is_c3 else 11, True, col)
    text(s7, Inches(1.65), y + (Inches(0.5) if is_c3 else Inches(0.45)), Inches(10.5), Inches(1.0 if is_c3 else 0.4), role, 11 if is_c3 else 10, color=BLACK)

note(s7, "Trả lời trực tiếp câu hỏi của Sếp: Lãnh đạo Quảng Ngãi là Chủ quản Dữ liệu Cấp 3 (Data Owner). "
     "Sếp có quyền quyết định tối cao: dữ liệu nào được chuyển đi, dữ liệu nào giữ lại nội bộ. "
     "TCT ở Cấp 5 chỉ là thợ kỹ thuật vận hành máy, không có quyền lấy dữ liệu nghiệp vụ của mình.")

# =========================================================================
# SLIDE 8: [MỚI] NGUYÊN TẮC BẢO MẬT & DỮ LIỆU NÀO ĐƯỢC PHÉP CHUYỂN?
# =========================================================================
s8 = prs.slides.add_slide(BL)
header(s8, "DỮ LIỆU NÀO ĐƯỢC CHUYỂN ĐI? DỮ LIỆU NÀO Ở LẠI NỘI BỘ?", "BẢO MẬT & CHỦ QUYỀN DỮ LIỆU — THEO QUY CHẾ TỔNG CÔNG TY", "07")

shape(s8, Inches(0.8), Inches(1.5), Inches(5.6), Inches(4.7), BLUE_SOFT, BLUE, Pt(1.5), rounded=True)
shape(s8, Inches(0.8), Inches(1.5), Inches(5.6), Inches(0.5), BLUE)
text(s8, Inches(0.8), Inches(1.52), Inches(5.6), Inches(0.45), "DỮ LIỆU Ở LẠI NỘI BỘ QUẢNG NGÃI (ON-PREM)", 12, True, WHITE, PP_ALIGN.CENTER)

stay_items = [
    ("Dữ liệu chuyên ngành chi tiết:", "Số liệu sản xuất, nhật trình thi công, chi tiết thiết bị xưởng."),
    ("Dữ liệu nhạy cảm kinh doanh:", "Định mức giá thầu, chi phí riêng từng dự án, biên lợi nhuận."),
    ("Dữ liệu cá nhân chưa xử lý:", "Hồ sơ nhân sự chi tiết, lương thưởng, CCCD (theo Luật BVDLCN)."),
    ("Toàn bộ cơ sở dữ liệu gốc:", "Kế toán, Nhân sự, Vật tư giữ nguyên tại chỗ, bảo mật 100%.")
]
for i, (head_t, body_t) in enumerate(stay_items):
    y = Inches(2.2 + i * 1.0)
    text(s8, Inches(1.1), y, Inches(5.0), Inches(0.3), "🔒  " + head_t, 11, True, NAVY)
    text(s8, Inches(1.35), y + Inches(0.28), Inches(4.7), Inches(0.65), body_t, 10, color=BLACK)

shape(s8, Inches(6.8), Inches(1.5), Inches(5.7), Inches(4.7), GREEN_SOFT, GREEN, Pt(1.5), rounded=True)
shape(s8, Inches(6.8), Inches(1.5), Inches(5.7), Inches(0.5), GREEN)
text(s8, Inches(6.8), Inches(1.52), Inches(5.7), Inches(0.45), "DỮ LIỆU ĐƯỢC PHÉP ĐỒNG BỘ VỀ TCT", 12, True, WHITE, PP_ALIGN.CENTER)

send_items = [
    ("29 Danh mục Master Data:", "Mã khách hàng, đối tác, danh mục vật tư dùng chung (đã ánh xạ)."),
    ("Báo cáo số liệu tổng hợp:", "Doanh thu, chi phí tổng hợp, tồn kho chung phục vụ hợp nhất toàn PTSC."),
    ("Dữ liệu đã che mờ (Masking):", "Mã hóa thông tin nhạy cảm trước khi truyền (tuân thủ Nghị định 13)."),
    ("Đổ vào Tenant riêng cách ly:", "Chỉ Quảng Ngãi xem được số liệu chi tiết của mình trên Cloud TCT.")
]
for i, (head_t, body_t) in enumerate(send_items):
    y = Inches(2.2 + i * 1.0)
    text(s8, Inches(7.1), y, Inches(5.1), Inches(0.3), "✔  " + head_t, 11, True, GREEN)
    text(s8, Inches(7.35), y + Inches(0.28), Inches(4.8), Inches(0.65), body_t, 10, color=BLACK)

shape(s8, Inches(0.8), Inches(6.35), Inches(11.7), Inches(0.7), GOLD_SOFT, GOLD, Pt(1), rounded=True)
shape(s8, Inches(0.84), Inches(6.37), Inches(0.07), Inches(0.66), GOLD)
text(s8, Inches(1.1), Inches(6.4), Inches(11.2), Inches(0.6),
     "NGUYÊN TẮC QUY ĐỊNH CỦA TCT: Dữ liệu quan trọng & nhạy cảm ở lại On-prem tại Đơn vị; Cloud TCT chỉ nhận dữ liệu tổng hợp và danh mục chuẩn hóa. Quảng Ngãi giữ toàn quyền kiểm soát!",
     11.5, True, NAVY, PP_ALIGN.CENTER)

note(s8, "TCT quy định rất rõ: Dữ liệu chi tiết nhạy cảm và bí mật kinh doanh ở lại Quảng Ngãi. "
     "Chỉ dữ liệu tổng hợp và danh mục Master Data mới đẩy lên TCT sau khi đã che mờ (masking).")

# =========================================================================
# SLIDE 9: KIẾN TRÚC TRỤC TÍCH HỢP HYBRID (MINIO + FABRIC)
# =========================================================================
s9 = prs.slides.add_slide(BL)
header(s9, "TRỤC TÍCH HỢP HYBRID: KẾT NỐI TỪ QUẢNG NGÃI LÊN MINIO & FABRIC", "MÔ HÌNH KỸ THUẬT — ĐÁP ỨNG ĐÚNG CHUẨN HYBRID CỦA TCT", "08")

# Layer 1: QN Local
shape(s9, Inches(0.3), Inches(1.35), Inches(2.7), Inches(5.45), BLUE_SOFT, BLUE, Pt(1), rounded=True)
text(s9, Inches(0.3), Inches(1.38), Inches(2.7), Inches(0.35), "1. NGUỒN DỮ LIỆU QN", 9, True, BLUE, PP_ALIGN.CENTER)

sources = ["Kế toán – Tài chính", "Quản lý Nhân sự", "Vật tư – Kho xưởng", "Hợp đồng – Dự án"]
for i, name in enumerate(sources):
    y = Inches(1.95 + i * 1.2)
    shape(s9, Inches(0.55), y, Inches(2.2), Inches(0.85), WHITE, BLUE, Pt(0.8), rounded=True)
    text(s9, Inches(0.65), y + Inches(0.1), Inches(2.0), Inches(0.35), name, 10.5, True, NAVY, PP_ALIGN.CENTER)
    text(s9, Inches(0.65), y + Inches(0.45), Inches(2.0), Inches(0.3), "Cơ sở dữ liệu", 8.5, color=GRAY, align=PP_ALIGN.CENTER)
    arrow_shape(s9, Inches(2.8), y + Inches(0.25), "right", Inches(0.35), Inches(0.22))

# Layer 2: Integration Gateway
shape(s9, Inches(3.3), Inches(1.35), Inches(3.5), Inches(5.45), GOLD_SOFT, GOLD, Pt(1.5), rounded=True)
text(s9, Inches(3.3), Inches(1.38), Inches(3.5), Inches(0.35), "2. TRẠM TRUNG CHUYỂN (VM tại QN)", 9, True, GOLD, PP_ALIGN.CENTER)

gw_steps = [
    ("Trích xuất (Read-Only)", "Chỉ đọc bản sao dữ liệu\nKhông làm chậm hệ thống gốc"),
    ("Ánh xạ 29 Master Data", "Quy đổi mã nội bộ khớp\nchuẩn Petrovietnam / PTSC"),
    ("Che mờ (Data Masking)", "Bảo vệ thông tin cá nhân\ntheo Nghị định 13 & Luật BVDLCN"),
    ("Mã hóa & Đóng gói", "Chuẩn bảo mật AES-256\nSẵn sàng truyền qua VPN"),
]
for i, (st, sd) in enumerate(gw_steps):
    y = Inches(1.95 + i * 1.2)
    shape(s9, Inches(3.55), y, Inches(3.0), Inches(0.85), WHITE, GOLD, Pt(0.8), rounded=True)
    numbered_label(s9, Inches(3.65), y + Inches(0.22), i + 1, GOLD)
    text(s9, Inches(4.15), y + Inches(0.08), Inches(2.3), Inches(0.35), st, 10.5, True, NAVY)
    text(s9, Inches(4.15), y + Inches(0.42), Inches(2.3), Inches(0.4), sd, 8.5, color=GRAY)

arrow_shape(s9, Inches(6.9), Inches(3.65), "right", Inches(0.4), Inches(0.22))

# Layer 3: VPN Tunnel
shape(s9, Inches(7.4), Inches(2.3), Inches(1.5), Inches(3.0), GRAY_LIGHT, NAVY, Pt(1), rounded=True)
text(s9, Inches(7.4), Inches(2.6), Inches(1.5), Inches(0.35), "3. ĐƯỜNG TRUYỀN", 9, True, NAVY, PP_ALIGN.CENTER)
text(s9, Inches(7.4), Inches(2.95), Inches(1.5), Inches(0.35), "VPN Site-to-Site", 10, True, NAVY, PP_ALIGN.CENTER)
text(s9, Inches(7.4), Inches(3.5), Inches(1.5), Inches(0.9), "Kênh riêng mã hóa\nTuân thủ quy hoạch\n8 Zones TCT", 8.5, color=GRAY, align=PP_ALIGN.CENTER)

arrow_shape(s9, Inches(9.0), Inches(3.65), "right", Inches(0.4), Inches(0.22))

# Layer 4: TCT HYBRID HUB
shape(s9, Inches(9.5), Inches(1.35), Inches(3.5), Inches(5.45), GREEN_SOFT, GREEN, Pt(1), rounded=True)
text(s9, Inches(9.5), Inches(1.38), Inches(3.5), Inches(0.35), "4. HẠ TẦNG HYBRID CỦA TCT", 9, True, GREEN, PP_ALIGN.CENTER)

# On-prem MinIO
shape(s9, Inches(9.75), Inches(1.95), Inches(3.0), Inches(1.4), WHITE, TEAL, Pt(0.8), rounded=True)
text(s9, Inches(9.75), Inches(2.05), Inches(3.0), Inches(0.3), "ON-PREMISE (Datacenter TCT)", 9.5, True, TEAL, PP_ALIGN.CENTER)
text(s9, Inches(9.75), Inches(2.35), Inches(3.0), Inches(0.35), "MinIO Lakehouse + MDM", 10.5, True, NAVY, PP_ALIGN.CENTER)
text(s9, Inches(9.75), Inches(2.7), Inches(3.0), Inches(0.55), "Lưu trữ dữ liệu gốc an toàn, kiểm soát truy cập", 8.5, color=GRAY, align=PP_ALIGN.CENTER)

arrow_shape(s9, Inches(11.1), Inches(3.45), "down", Inches(0.22), Inches(0.3))

# Cloud Fabric
shape(s9, Inches(9.75), Inches(3.85), Inches(3.0), Inches(1.5), WHITE, BLUE, Pt(0.8), rounded=True)
text(s9, Inches(9.75), Inches(3.95), Inches(3.0), Inches(0.3), "CLOUD (Azure / Fabric 20TB)", 9.5, True, BLUE, PP_ALIGN.CENTER)
text(s9, Inches(9.75), Inches(4.25), Inches(3.0), Inches(0.35), "Workspace L3 Quảng Ngãi", 10.5, True, NAVY, PP_ALIGN.CENTER)
text(s9, Inches(9.75), Inches(4.6), Inches(3.0), Inches(0.65), "Không gian riêng cách ly hoàn toàn bằng RBAC", 8.5, color=GRAY, align=PP_ALIGN.CENTER)

arrow_shape(s9, Inches(11.1), Inches(5.45), "down", Inches(0.22), Inches(0.3))

# Dashboard
shape(s9, Inches(9.75), Inches(5.85), Inches(3.0), Inches(0.85), NAVY, rounded=True)
text(s9, Inches(9.75), Inches(5.95), Inches(3.0), Inches(0.3), "DASHBOARD POWER BI", 10, True, RGBColor(253, 224, 71), PP_ALIGN.CENTER)
text(s9, Inches(9.75), Inches(6.25), Inches(3.0), Inches(0.35), "Ban Giám đốc xem báo cáo trên PC / Mobile", 8.5, color=WHITE, align=PP_ALIGN.CENTER)

note(s9, "Đúng kiến trúc Hybrid của TCT: Dữ liệu từ QN đi qua đường hầm VPN, tiếp đất vào MinIO on-prem "
     "và Workspace L3 trên Microsoft Fabric Cloud để xuất ra Dashboard Power BI.")

# =========================================================================
# SLIDE 10: SO SÁNH 2 PHƯƠNG ÁN THEO CHỈ ĐẠO BAN GIÁM ĐỐC
# =========================================================================
s10 = prs.slides.add_slide(BL)
header(s10, "SO SÁNH 2 PHƯƠNG ÁN THEO CHỈ ĐẠO CỦA BAN GIÁM ĐỐC", "PHÂN TÍCH — TRẢ LỜI CÂU HỎI TRỌNG TÂM TRONG CHỈ ĐẠO SỐ 9", "09")

shape(s10, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3), RED_SOFT, RED, Pt(1.5), rounded=True)
shape(s10, Inches(0.8), Inches(1.5), Inches(5.6), Inches(0.5), RED)
text(s10, Inches(0.8), Inches(1.5), Inches(5.6), Inches(0.5),
     "PHƯƠNG ÁN 1:  Thuê TCT triển khai hộ từ A-Z", 12, True, WHITE, PP_ALIGN.CENTER)

pa1 = [
    ("TCT không đủ nhân lực hỗ trợ:", "Ban CNTT TCT chỉ làm hạ tầng lõi Hub; không có người xuống QN khảo sát từng CSDL phần mềm của mình."),
    ("Bị động và nguy cơ trễ hạn 2026:", "TCT phải mở rộng cho hàng chục đơn vị. Nếu xếp hàng chờ TCT, QN chắc chắn không kịp tiến độ."),
    ("Thiếu báo cáo quản trị riêng cho BGĐ:", "TCT chỉ xây báo cáo tài chính vĩ mô cho Tập đoàn; không giải quyết được bài toán tồn xưởng, nhân công dự án của QN."),
    ("Chi phí phân bổ không chủ động:", "Phải gánh phân bổ chi phí dịch vụ từ TCT mà không kiểm soát được hạng mục chi tiết.")
]
for i, (t, d) in enumerate(pa1):
    y = Inches(2.2 + i * 1.05)
    text(s10, Inches(1.1), y, Inches(5.0), Inches(0.3), "✕  " + t, 11, True, RED)
    text(s10, Inches(1.35), y + Inches(0.28), Inches(4.7), Inches(0.65), d, 9.5, color=BLACK)

shape(s10, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3), GREEN_SOFT, GREEN, Pt(1.5), rounded=True)
shape(s10, Inches(6.8), Inches(1.5), Inches(5.7), Inches(0.5), GREEN)
text(s10, Inches(6.8), Inches(1.5), Inches(5.7), Inches(0.5),
     "PHƯƠNG ÁN 2:  Phát triển riêng đồng bộ TCT (CÓ NCC)   ĐỀ XUẤT", 12, True, WHITE, PP_ALIGN.CENTER)

pa2 = [
    ("NCC làm việc trực tiếp tại Quảng Ngãi:", "Khảo sát, thiết kế và đấu nối trúng các phần mềm đặc thù mà công ty đang dùng hàng ngày."),
    ("Chủ động 100% tiến độ:", "Cam kết hoàn thành trạm trung chuyển trong 3–4 tháng, đảm bảo hoàn thành đúng mốc hạn 2026 của TCT."),
    ("Thiết kế Dashboard đo ni đóng giày cho BGĐ:", "Vừa đồng bộ chuẩn cho TCT, vừa có báo cáo quản trị chuyên sâu phục vụ điều hành của Ban Giám đốc QN."),
    ("Tối ưu chi phí & Chủ động kiểm soát:", "Tận dụng hạ tầng đám mây TCT đã mua sẵn; chi phí thuê NCC được đấu thầu cạnh tranh rõ ràng.")
]
for i, (t, d) in enumerate(pa2):
    y = Inches(2.2 + i * 1.05)
    text(s10, Inches(7.1), y, Inches(5.1), Inches(0.3), "✔  " + t, 11, True, GREEN)
    text(s10, Inches(7.35), y + Inches(0.28), Inches(4.8), Inches(0.65), d, 9.5, color=BLACK)

note(s10, "So sánh 2 phương án: TCT chỉ làm Hub chung, không đủ người làm thay QN. "
     "Chọn Phương án 2: Thuê NCC phát triển riêng đồng bộ TCT là phương án khả thi duy nhất.")

# =========================================================================
# SLIDE 11: CƠ CẤU ĐẦU TƯ & CẤU TRÚC CHI PHÍ (THEO PHƯƠNG ÁN 3 CỦA TCT)
# =========================================================================
s11 = prs.slides.add_slide(BL)
header(s11, "CƠ CHẾ ĐẦU TƯ & CẤU TRÚC CHI PHÍ THEO ĐỊNH HƯỚNG TCT", "TÀI CHÍNH — BÁM SÁT PHƯƠNG ÁN 3 CỦA TỔNG CÔNG TY", "10")

# Table
table_shape = s11.shapes.add_table(8, 4, Inches(0.8), Inches(1.4), Inches(11.733), Inches(4.0))
tbl = table_shape.table
tbl.columns[0].width = Inches(0.5)
tbl.columns[1].width = Inches(3.8)
tbl.columns[2].width = Inches(4.833)
tbl.columns[3].width = Inches(2.6)

headers_data = ["STT", "Cấu phần chi phí (Theo chuẩn TCT)", "Nội dung & Trách nhiệm thực hiện", "Dự toán ngân sách"]
for i, h in enumerate(headers_data):
    c = tbl.cell(0, i); c.fill.solid(); c.fill.fore_color.rgb = NAVY
    p = c.text_frame.paragraphs[0]; p.text = h
    p.font.size = Pt(10.5); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER

cost_structure = [
    ("1", "Hạ tầng phần cứng tại Đơn vị", "Tận dụng máy chủ ảo (VM) sẵn có của Quảng Ngãi", "0 VNĐ (Đơn vị tự có)"),
    ("2", "Hạ tầng Cloud Hub & Bản quyền khung", "Microsoft Fabric, MinIO, Purview, MDM do TCT đầu tư sẵn", "TCT ĐÃ ĐẦU TƯ"),
    ("3", "Dịch vụ: Khảo sát & Thiết kế kiến trúc", "Khảo sát CSDL 4 phần mềm + Lập bảng ánh xạ 29 Master Data", "[Chờ NCC báo giá sau khảo sát]"),
    ("4", "Dịch vụ: Xây dựng Trục tích hợp & ETL", "Lập trình đường ống trích xuất, làm sạch, mã hóa VPN lên TCT", "[Chờ NCC báo giá sau khảo sát]"),
    ("5", "Dịch vụ: Xây dựng Dashboard Power BI", "Xây dựng các bảng báo cáo quản trị phục vụ Ban Giám đốc QN", "[Chờ NCC báo giá sau khảo sát]"),
    ("6", "Đào tạo chuyển giao & Bảo hành", "Bàn giao mã nguồn, đào tạo đội ngũ IT Quảng Ngãi làm chủ 100%", "[Chờ NCC báo giá sau khảo sát]"),
    ("7", "Chi phí vận hành nền tảng hàng năm", "TCT phân bổ chi phí theo mức độ sử dụng thực tế (Usage-based)", "Theo cơ chế TCT ban hành")
]
for ri, row_vals in enumerate(cost_structure, 1):
    for ci, val in enumerate(row_vals):
        c = tbl.cell(ri, ci); c.fill.solid()
        c.fill.fore_color.rgb = WHITE if ri % 2 == 1 else GRAY_LIGHT
        p = c.text_frame.paragraphs[0]; p.text = val
        p.font.size = Pt(9.5); p.font.color.rgb = BLACK
        if ci == 0: p.alignment = PP_ALIGN.CENTER
        if ci == 3:
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            if "0 VNĐ" in val or "TCT ĐÃ" in val:
                p.font.color.rgb = GREEN
            elif "Chờ NCC" in val:
                p.font.color.rgb = GOLD
            else:
                p.font.color.rgb = BLUE

shape(s11, Inches(0.8), Inches(5.6), Inches(11.733), Inches(1.4), BLUE_SOFT, BLUE, Pt(1), rounded=True)
shape(s11, Inches(0.84), Inches(5.62), Inches(0.08), Inches(1.36), BLUE)
text(s11, Inches(1.15), Inches(5.68), Inches(11.0), Inches(0.3),
    "NGUYÊN TẮC MINH BẠCH VỀ CHI PHÍ TRÌNH BAN GIÁM ĐỐC:", 11, True, NAVY)
tf_c_note = text(s11, Inches(1.15), Inches(6.0), Inches(11.0), Inches(0.9),
    "• Phần cứng & Bản quyền: Quảng Ngãi tiết kiệm hàng tỷ đồng nhờ tận dụng hạ tầng TCT đã mua sẵn và máy chủ hiện hữu.\n"
    "• Chi phí thuê NCC triển khai: Sẽ được xác định chính xác sau khi hoàn thành Giai đoạn 1 (Khảo sát nội bộ) và nhận báo giá cạnh tranh từ các NCC uy tín.\n"
    "• Chi phí vận hành: Áp dụng Phương án 3 của TCT (phân bổ theo mức sử dụng; đơn vị chưa có nguồn TCT hỗ trợ đầu tư trước ghi nợ).", 9.5, color=BLACK)

note(s11, "Minh bạch chi phí: TCT đã gánh phần Cloud và phần mềm khung. Chi phí thuê NCC sẽ để trống "
     "và xác định chính xác sau khi khảo sát kỹ thuật ở Giai đoạn 1, tránh đưa con số cảm tính.")

# =========================================================================
# SLIDE 12: ĐẦU RA CỤ THỂ — BAN GIÁM ĐỐC ĐƯỢC GÌ?
# =========================================================================
s12 = prs.slides.add_slide(BL)
header(s12, "GIÁ TRỊ MANG LẠI CHO CÔNG TÁC ĐIỀU HÀNH BAN GIÁM ĐỐC", "ĐẦU RA CỤ THỂ — HIỆU QUẢ KHI DỰ ÁN HOÀN THÀNH", "11")

outputs = [
    ("BÁO CÁO ĐIỀU HÀNH\nTRỰC QUAN REAL-TIME", BLUE, BLUE_SOFT,
     "Dashboard Power BI cập nhật tự động mỗi ngày.\nBan Giám đốc xem trực quan trên máy tính hoặc iPad/di động mọi lúc mọi nơi."),
    ("DỮ LIỆU TỰ ĐỘNG\nKHÔNG PHỤ THUỘC EXCEL", TEAL, RGBColor(204, 251, 241),
     "Chấm dứt việc phòng ban mất 3–5 ngày làm báo cáo thủ công.\nSố liệu tự động trích xuất và đẩy về hệ thống."),
    ("MỘT NGUỒN SỰ THẬT\nCHUẨN XÁC DUY NHẤT", GREEN, GREEN_SOFT,
     "Kế toán, Vật tư, Nhân sự đều đồng bộ cùng 1 con số.\nKhông còn hiện tượng số liệu vênh nhau khi giải trình."),
    ("HOÀN THÀNH ĐÚNG HẠN\nCHỈ TIÊU CĐS TCT GIAO", GOLD, GOLD_SOFT,
     "Đáp ứng đúng hạn Nghị quyết số 10 HĐQT PTSC.\nTuân thủ đầy đủ Luật Dữ liệu và Luật Bảo vệ DLCN mới."),
]
for i, (title, col, bg_c, desc) in enumerate(outputs):
    x = Inches(0.8 + i * 3.05)
    shape(s12, x, Inches(1.5), Inches(2.7), Inches(5.2), bg_c, col, Pt(1.5), rounded=True)
    shape(s12, x + Inches(0.04), Inches(1.52), Inches(0.08), Inches(5.14), col)
    numbered_label(s12, x + Inches(0.25), Inches(1.7), i + 1, col)
    text(s12, x + Inches(0.15), Inches(2.3), Inches(2.4), Inches(0.8), title, 13, True, col, PP_ALIGN.CENTER)
    shape(s12, x + Inches(0.4), Inches(3.2), Inches(1.9), Inches(0.01), GRAY_MED)
    text(s12, x + Inches(0.2), Inches(3.4), Inches(2.3), Inches(3.0), desc, 11, color=BLACK, align=PP_ALIGN.CENTER)

note(s12, "4 giá trị cụ thể: Dashboard tức thời cho Sếp, hết phụ thuộc Excel, "
     "một nguồn sự thật duy nhất, và hoàn thành đúng hạn chỉ tiêu CĐS.")

# =========================================================================
# SLIDE 13: LỘ TRÌNH 3 GIAI ĐOẠN TRIỂN KHAI
# =========================================================================
s13 = prs.slides.add_slide(BL)
header(s13, "KẾ HOẠCH TRIỂN KHAI 3 GIAI ĐOẠN", "LỘ TRÌNH — TỪNG BƯỚC RÕ RÀNG, KIỂM SOÁT RỦI RO", "12")

timeline = [
    ("GIAI ĐOẠN 1", "Tháng 3 – 4 / 2026", "KHẢO SÁT & ĐỀ BÀI KỸ THUẬT", "Nội bộ tự chủ trì  •  Chi phí 0 đ",
     BLUE, BLUE_SOFT,
     ["Thành lập Tổ công tác Data Platform nội bộ",
      "Khảo sát CSDL 4 phần mềm nghiệp vụ",
      "Rà soát và lập bảng ánh xạ 29 Master Data TCT",
      "Lập Hồ sơ yêu cầu kỹ thuật (TOR) để mời thầu NCC",
      "*(Chi tiết xem file Kế hoạch Giai đoạn 1 đính kèm)*"]),
    ("GIAI ĐOẠN 2", "Tháng 5 – 7 / 2026", "LỰA CHỌN NCC & TRIỂN KHAI", "Giai đoạn chính  •  Thuê NCC",
     GOLD, GOLD_SOFT,
     ["Lựa chọn NCC có giải pháp tối ưu, giá cạnh tranh",
      "Cài đặt Trạm trung chuyển dữ liệu trên VM nội bộ",
      "Lập trình đường ống ETL tự động & mã hóa VPN",
      "Kiểm thử thông luồng dữ liệu lên MinIO & Fabric TCT"]),
    ("GIAI ĐOẠN 3", "Tháng 8 / 2026 →", "BÀN GIAO & KHAI THÁC", "IT làm chủ 100% vận hành",
     GREEN, GREEN_SOFT,
     ["Nghiệm thu hệ thống Trục tích hợp dữ liệu",
      "NCC bàn giao toàn bộ mã nguồn & tài liệu kiến trúc",
      "Đào tạo chuyển giao kỹ thuật cho IT Quảng Ngãi",
      "Xây dựng và đưa vào sử dụng Dashboard Power BI"]),
]
for i, (phase, period, title, tag, col, bg_c, items) in enumerate(timeline):
    x = Inches(0.8 + i * 4.1)
    shape(s13, x, Inches(1.5), Inches(3.7), Inches(5.2), bg_c, col, Pt(1.5), rounded=True)
    shape(s13, x + Inches(0.04), Inches(1.52), Inches(0.08), Inches(5.14), col)
    text(s13, x + Inches(0.25), Inches(1.65), Inches(3.2), Inches(0.3), phase, 10, True, col)
    text(s13, x + Inches(0.25), Inches(1.92), Inches(3.2), Inches(0.3), period, 13, True, NAVY)
    text(s13, x + Inches(0.25), Inches(2.25), Inches(3.2), Inches(0.4), title, 13, True, col)
    shape(s13, x + Inches(0.25), Inches(2.72), Inches(3.1), Inches(0.35), col, rounded=True)
    text(s13, x + Inches(0.25), Inches(2.72), Inches(3.1), Inches(0.35), tag, 9, True, WHITE, PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        text(s13, x + Inches(0.3), Inches(3.3 + j * 0.55), Inches(3.1), Inches(0.5),
             "—  " + item, 11, color=BLACK if not item.startswith("*") else GRAY)
    if i < 2:
        arrow_shape(s13, Inches(4.55 + i * 4.1), Inches(3.8))

note(s13, "Lộ trình 3 giai đoạn: GĐ1 nội bộ tự làm hoàn toàn miễn phí để ra đề bài; "
     "sau đó mới mời thầu NCC làm GĐ2; GĐ3 là nghiệm thu và làm chủ.")

# =========================================================================
# SLIDE 14: KIẾN NGHỊ BAN GIÁM ĐỐC PHÊ DUYỆT
# =========================================================================
s14 = prs.slides.add_slide(BL)
header(s14, "KIẾN NGHỊ BAN GIÁM ĐỐC XEM XÉT VÀ CHỈ ĐẠO", "ĐỀ XUẤT — 3 NỘI DUNG XIN CHỦ TRƯƠNG ĐỂ TRIỂN KHAI", "13")

proposals = [
    ("CHẤP THUẬN PHƯƠNG ÁN 2:",
     "Phê duyệt chủ trương: Công ty chủ động phát triển Trục tích hợp dữ liệu riêng đồng bộ với TCT (thuê Nhà cung cấp chuyên nghiệp tư vấn & triển khai) để đảm bảo tiến độ năm 2026.",
     BLUE, BLUE_SOFT),
    ("THÀNH LẬP TỔ CÔNG TÁC NỘI BỘ:",
     "Thành lập Tổ công tác Data Platform gồm Bộ phận CNTT (chủ trì kỹ thuật) và các cán bộ phụ trách dữ liệu (Key users) từ Kế toán, Vật tư, Nhân sự để bắt đầu ngay Giai đoạn 1 (chi phí 0 VNĐ).",
     GOLD, GOLD_SOFT),
    ("CHO PHÉP TIẾP XÚC NCC LẤY BÁO GIÁ:",
     "Cho phép Tổ công tác khảo sát thị trường, truyền đạt đúng khung kiến trúc và bài toán kỹ thuật để lấy báo giá cạnh tranh từ các NCC uy tín, hoàn thiện dự toán chi tiết trình Ban Giám đốc phê duyệt trước khi ký kết.",
     GREEN, GREEN_SOFT),
]
for i, (title, desc, col, bg_c) in enumerate(proposals):
    x = Inches(0.8 + i * 4.1)
    shape(s14, x, Inches(1.5), Inches(3.7), Inches(4.5), bg_c, col, Pt(2), rounded=True)
    shape(s14, x + Inches(0.04), Inches(1.52), Inches(0.08), Inches(4.44), col)
    numbered_label(s14, x + Inches(0.3), Inches(1.7), i + 1, col)
    text(s14, x + Inches(0.85), Inches(1.7), Inches(2.6), Inches(0.7), title, 13.5, True, col)
    shape(s14, x + Inches(0.3), Inches(2.55), Inches(3.1), Inches(0.01), GRAY_MED)
    text(s14, x + Inches(0.3), Inches(2.75), Inches(3.1), Inches(2.8), desc, 12, color=BLACK)

shape(s14, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.65), GREEN_SOFT, GREEN, Pt(1.5), rounded=True)
shape(s14, Inches(0.84), Inches(6.32), Inches(0.07), Inches(0.58), GREEN)
text(s14, Inches(1.1), Inches(6.3), Inches(11.2), Inches(0.65),
     "Kính trình Ban Giám đốc xem xét, thông qua chủ trương để đơn vị kịp tiến độ kết nối năm 2026 của Tổng công ty!",
     12, True, RGBColor(22, 100, 52), PP_ALIGN.CENTER)

note(s14, "3 đề xuất xin phê duyệt: 1 là chủ trương chọn PA2; 2 là thành lập tổ công tác nội bộ; "
     "3 là cho phép khảo sát lấy báo giá chính thức từ NCC để trình Sếp duyệt dự toán.")

# =========================================================================
# SLIDE 15: CLOSING
# =========================================================================
s15 = prs.slides.add_slide(BL)
shape(s15, 0, 0, W, H, NAVY)
shape(s15, 0, Inches(2.2), W, Inches(0.035), GOLD)
shape(s15, 0, Inches(5.0), W, Inches(0.015), RGBColor(60, 90, 140))

text(s15, Inches(1.0), Inches(2.6), Inches(11.3), Inches(1.0),
     "XIN TRÂN TRỌNG CẢM ƠN\nBAN GIÁM ĐỐC", 28, True, WHITE)
text(s15, Inches(1.0), Inches(4.0), Inches(11.3), Inches(0.5),
     "Kính mời Ban Giám đốc cho ý kiến chỉ đạo.", 15, color=RGBColor(200, 215, 235))

text(s15, Inches(1.0), Inches(5.3), Inches(11.3), Inches(0.4),
     "Tài liệu đính kèm:  Kế hoạch chi tiết Giai đoạn 1, 2, 3", 11, color=RGBColor(155, 185, 220))
text(s15, Inches(1.0), Inches(5.7), Inches(11.3), Inches(0.4),
     "Đơn vị thực hiện:  Tổ Công tác CĐS – Bộ phận CNTT PTSC Quảng Ngãi", 11, color=RGBColor(155, 185, 220))

note(s15, "Chúng tôi xin trân trọng cảm ơn và lắng nghe ý kiến chỉ đạo của Ban Giám đốc!")

# Save
output = r"d:\My Profiles\DataPlatform\bao_cao_dataplatform_ptsc_qn_v5.pptx"
prs.save(output)
print(f"SUCCESS: {output}")
print(f"Total slides: {len(prs.slides)}")
