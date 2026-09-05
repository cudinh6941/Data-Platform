"""
BÁO CÁO DATA PLATFORM PTSC QUẢNG NGÃI - VERSION 7 (FULL PROFESSIONAL DECK)
- Ngôn ngữ: Chuyên nghiệp nhưng dễ tiếp cận (không bình dân, không quá hàn lâm)
- Bổ sung chi tiết: Phần PTSC QN cần làm ở giai đoạn tới + Cách kết nối ETL với Hub
- Tổng: 18 slide
"""
import os, sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Corporate Color Palette
NAVY       = RGBColor(12, 35, 68)
NAVY_MID   = RGBColor(24, 59, 108)
BLUE       = RGBColor(30, 100, 190)
BLUE_SOFT  = RGBColor(230, 240, 252)
GREEN      = RGBColor(21, 128, 80)
GREEN_SOFT = RGBColor(232, 248, 239)
RED        = RGBColor(180, 40, 40)
RED_SOFT   = RGBColor(252, 235, 235)
GOLD       = RGBColor(180, 115, 0)
GOLD_SOFT  = RGBColor(255, 245, 225)
WHITE      = RGBColor(255, 255, 255)
BLACK      = RGBColor(28, 35, 48)
GRAY       = RGBColor(108, 118, 132)
GRAY_LIGHT = RGBColor(243, 245, 248)
GRAY_MED   = RGBColor(200, 208, 218)
TEAL       = RGBColor(13, 148, 136)
PURPLE     = RGBColor(109, 40, 217)

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BL = prs.slide_layouts[6]

def shape(slide, l, t, w, h, fill, line_c=None, line_w=None, rounded=False):
    sh_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    r = slide.shapes.add_shape(sh_type, l, t, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = fill
    if line_c:
        r.line.color.rgb = line_c
        if line_w: r.line.width = line_w
    else:
        r.line.fill.background()
    return r

def text(slide, l, t, w, h, content, size=14, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = content
    p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = color; p.alignment = align
    return tf

def header(slide, title, subtitle="", slide_num=""):
    shape(slide, 0, 0, W, Inches(1.05), NAVY)
    shape(slide, 0, Inches(1.05), W, Inches(0.035), GOLD)
    if subtitle:
        text(slide, Inches(0.8), Inches(0.1), Inches(10), Inches(0.25), subtitle.upper(), 9, True, RGBColor(155, 185, 220))
    text(slide, Inches(0.8), Inches(0.38), Inches(11), Inches(0.55), title, 18, True, WHITE)
    if slide_num:
        text(slide, Inches(12.0), Inches(0.38), Inches(0.8), Inches(0.5), slide_num, 12, True, RGBColor(155, 185, 220), PP_ALIGN.RIGHT)
    shape(slide, 0, Inches(7.2), W, Inches(0.008), GRAY_MED)
    text(slide, Inches(0.8), Inches(7.22), Inches(11.7), Inches(0.25),
         "PTSC Quảng Ngãi  ·  Báo cáo triển khai Nền tảng Dữ liệu & Trục tích hợp  ·  Thực hiện Chỉ đạo số 9 BGĐ", 8, color=GRAY)

def numbered_label(slide, l, t, num, color=BLUE):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, Inches(0.4), Inches(0.4))
    c.fill.solid(); c.fill.fore_color.rgb = color; c.line.fill.background()
    tf = c.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = str(num)
    p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER

def arrow_shape(slide, l, t, direction="right", w=Inches(0.5), h=Inches(0.3)):
    sh_type = MSO_SHAPE.RIGHT_ARROW if direction == "right" else MSO_SHAPE.DOWN_ARROW
    if direction == "down": w, h = h, w
    ar = slide.shapes.add_shape(sh_type, l, t, w, h)
    ar.fill.solid(); ar.fill.fore_color.rgb = GOLD; ar.line.fill.background()

def note(slide, txt_content):
    slide.notes_slide.notes_text_frame.text = txt_content

print("Building 18 slides (V7 Professional Edition)...")

# =========================================================================
# SLIDE 1: COVER
# =========================================================================
s1 = prs.slides.add_slide(BL)
shape(s1, 0, 0, W, H, NAVY)
shape(s1, 0, Inches(1.4), W, Inches(0.04), GOLD)
shape(s1, 0, Inches(5.8), W, Inches(0.015), RGBColor(60, 90, 140))

text(s1, Inches(1.0), Inches(1.8), Inches(11), Inches(0.4),
     "BÁO CÁO BAN GIÁM ĐỐC  —  THỰC HIỆN CHỈ ĐẠO MỤC 9", 13, True, RGBColor(180, 200, 225))
text(s1, Inches(1.0), Inches(2.5), Inches(11), Inches(1.6),
     "PHƯƠNG ÁN TRIỂN KHAI NỀN TẢNG DỮ LIỆU\nVÀ TRỤC TÍCH HỢP CÔNG TY – TỔNG CÔNG TY", 30, True, WHITE)
shape(s1, Inches(1.0), Inches(4.4), Inches(2.5), Inches(0.04), GOLD)

scope_items = [
    "Phần A — Tổng quan Nền tảng Dữ liệu: Bối cảnh, bản chất và tài nguyên sẵn có từ Tổng công ty",
    "Phần B — Hiện trạng tại PTSC Quảng Ngãi: Vị thế L3, thách thức nội bộ, chủ quyền dữ liệu",
    "Phần C — Kiến trúc kết nối & Phương án triển khai: Luồng tích hợp, so sánh phương án, dự toán",
    "Phần D — Kế hoạch hành động chi tiết: Các bước triển khai cụ thể, kết nối Hub và kiến nghị"
]
for i, item in enumerate(scope_items):
    text(s1, Inches(1.0), Inches(4.65 + i * 0.28), Inches(11), Inches(0.28), "—  " + item, 10.5, color=RGBColor(200, 215, 235))

text(s1, Inches(1.0), Inches(6.05), Inches(11), Inches(0.3),
     "Đơn vị thực hiện:  Tổ Công tác Chuyển đổi số & CNTT  –  PTSC Quảng Ngãi", 11.5, color=WHITE)
text(s1, Inches(1.0), Inches(6.4), Inches(11), Inches(0.3),
     "Kính trình:  Ban Giám đốc Công ty PTSC Quảng Ngãi", 11.5, True, RGBColor(253, 224, 71))

note(s1, "Kính thưa Ban Giám đốc, bộ phận CNTT xin báo cáo phương án triển khai Nền tảng Dữ liệu theo chỉ đạo số 9 của Ban Giám đốc. Bài báo cáo gồm 4 phần: Tổng quan giải pháp, hiện trạng Quảng Ngãi, kiến trúc kết nối và kế hoạch hành động chi tiết.")

# =========================================================================
# SLIDE 2: BỐI CẢNH — TẠI SAO TCT TRIỂN KHAI NỀN TẢNG DỮ LIỆU?
# =========================================================================
s2 = prs.slides.add_slide(BL)
header(s2, "BỐI CẢNH: TẠI SAO TỔNG CÔNG TY TRIỂN KHAI NỀN TẢNG DỮ LIỆU?", "PHẦN A: TỔNG QUAN GIẢI PHÁP — CĂN CỨ CHIẾN LƯỢC", "01")

shape(s2, Inches(0.8), Inches(1.4), Inches(5.6), Inches(5.4), RED_SOFT, RED, Pt(1.5), rounded=True)
shape(s2, Inches(0.8), Inches(1.4), Inches(5.6), Inches(0.5), RED)
text(s2, Inches(0.8), Inches(1.42), Inches(5.6), Inches(0.45), "THỰC TRẠNG CẦN GIẢI QUYẾT", 12, True, WHITE, PP_ALIGN.CENTER)

pains = [
    ("Dữ liệu phân tán tại các đơn vị thành viên", "Phần lớn số liệu sản xuất – kinh doanh phát sinh tại ĐVTV. Không tập hợp được thì giá trị nền tảng chỉ dừng ở Văn phòng TCT."),
    ("Báo cáo hợp nhất mang tính thủ công", "Tổng hợp tài chính, nhân sự, SXKD toàn Tổng công ty và báo cáo Tập đoàn PVN hiện chậm, thiếu nhất quán danh mục."),
    ("Dữ liệu chủ (Master Data) không thống nhất", "Mã khách hàng, nhà cung cấp, vật tư, dự án lệch nhau giữa các đơn vị, gây sai lệch hợp nhất và cản trở phân tích liên đơn vị."),
    ("Áp lực tuân thủ pháp lý mới", "Luật BVDLCN 91/2025/QH15, Luật An ninh mạng, Luật Dữ liệu, Quy chế Quản trị Dữ liệu Petrovietnam — yêu cầu quản trị thống nhất.")
]
for i, (title, desc) in enumerate(pains):
    y = Inches(2.1 + i * 1.1)
    shape(s2, Inches(1.0), y, Inches(5.2), Inches(0.95), WHITE, GRAY_MED, Pt(0.5), rounded=True)
    text(s2, Inches(1.2), y + Inches(0.06), Inches(4.8), Inches(0.28), title, 11, True, RED)
    text(s2, Inches(1.2), y + Inches(0.33), Inches(4.8), Inches(0.55), desc, 10, color=BLACK)

arrow_shape(s2, Inches(6.6), Inches(3.8))

shape(s2, Inches(7.3), Inches(1.4), Inches(5.2), Inches(5.4), GREEN_SOFT, GREEN, Pt(1.5), rounded=True)
shape(s2, Inches(7.3), Inches(1.4), Inches(5.2), Inches(0.5), GREEN)
text(s2, Inches(7.3), Inches(1.42), Inches(5.2), Inches(0.45), "MỤC TIÊU CHIẾN LƯỢC CỦA TCT", 12, True, WHITE, PP_ALIGN.CENTER)

goals = [
    ("Xây dựng Kho dữ liệu dùng chung (Data Lakehouse)", "Tập hợp dữ liệu toàn Tổng công ty vào một hệ thống trung tâm an toàn, phục vụ phân tích và ra quyết định."),
    ("Ban hành 29 Danh mục Dữ liệu chủ (Master Data)", "Thiết lập bộ mã thống nhất: 1 mã nhân viên, 1 mã vật tư, 1 mã dự án duy nhất xuyên suốt toàn Tổng công ty."),
    ("Tự động hóa báo cáo hợp nhất toàn Tổng", "Chuyển từ tổng hợp thủ công sang báo cáo quản trị (Dashboard) cập nhật tự động, phục vụ điều hành tức thời."),
    ("Nghị quyết 10/NQ-HĐQT: Mốc 2026–2027 bắt buộc", "Các đơn vị thành viên (bao gồm PTSC Quảng Ngãi) phải hoàn thành kết nối trong giai đoạn này — đây là KPI CĐS.")
]
for i, (title, desc) in enumerate(goals):
    y = Inches(2.1 + i * 1.1)
    shape(s2, Inches(7.5), y, Inches(4.8), Inches(0.95), WHITE, GRAY_MED, Pt(0.5), rounded=True)
    text(s2, Inches(7.7), y + Inches(0.06), Inches(4.4), Inches(0.28), title, 11, True, GREEN)
    text(s2, Inches(7.7), y + Inches(0.33), Inches(4.4), Inches(0.55), desc, 10, color=BLACK)

note(s2, "Tài liệu TCT nêu rõ 5 lý do bắt buộc mở rộng: Dữ liệu nằm ở đơn vị, báo cáo hợp nhất thủ công, Master Data không thống nhất, yêu cầu pháp lý mới, và thời điểm tối ưu sau GĐ1.")

# =========================================================================
# SLIDE 3: NỀN TẢNG DỮ LIỆU (DATA PLATFORM) — BẢN CHẤT KỸ THUẬT
# =========================================================================
s3 = prs.slides.add_slide(BL)
header(s3, "NỀN TẢNG DỮ LIỆU (DATA PLATFORM) — BẢN CHẤT KIẾN TRÚC 3 TẦNG", "PHẦN A: TỔNG QUAN GIẢI PHÁP — ĐỊNH NGHĨA & CẤU TRÚC", "02")

shape(s3, Inches(0.8), Inches(1.35), Inches(11.7), Inches(0.75), GOLD_SOFT, GOLD, Pt(1.5), rounded=True)
text(s3, Inches(1.1), Inches(1.4), Inches(11.1), Inches(0.65),
     "Nền tảng Dữ liệu không phải phần mềm mới thay thế hệ thống hiện hữu. Các phần mềm nghiệp vụ (Kế toán, Nhân sự, Vật tư) vẫn giữ nguyên.\n"
     "Bản chất: Hạ tầng tự động hóa thu thập, chuẩn hóa và trực quan hóa dữ liệu — vận hành ngầm, không ảnh hưởng đến quy trình làm việc hiện tại.",
     11, True, NAVY, PP_ALIGN.CENTER)

layers = [
    ("TẦNG 1: THU THẬP TỰ ĐỘNG", "(Data Ingestion)", BLUE, BLUE_SOFT,
     "Kết nối & trích xuất dữ liệu",
     ["Kết nối vào các CSDL nghiệp vụ hiện có (Kế toán, HRM, Vật tư, Dự án)",
      "Trích xuất bản sao chỉ đọc (Read-Only) — không can thiệp dữ liệu gốc",
      "Vận hành tự động theo lịch (batch hàng đêm hoặc near real-time)",
      "Công nghệ: Airbyte, Debezium hoặc API Connector tùy nguồn dữ liệu"]),
    ("TẦNG 2: CHUẨN HÓA & LƯU TRỮ", "(Data Lakehouse + MDM)", TEAL, RGBColor(230, 248, 246),
     "Làm sạch, quy đổi mã & lưu trữ",
     ["Gom dữ liệu vào Hồ dữ liệu tập trung (Lakehouse — MinIO + Apache Iceberg)",
      "Khử trùng lặp, phát hiện thiếu sót, kiểm tra chất lượng (dbt, Great Expectations)",
      "Ánh xạ tự động theo 29 Danh mục Dữ liệu chủ của TCT (MDM Golden Record)",
      "Kết quả: Một nguồn dữ liệu chuẩn duy nhất (Single Source of Truth)"]),
    ("TẦNG 3: PHÂN TÍCH & BÁO CÁO", "(Analytics & Power BI)", GREEN, GREEN_SOFT,
     "Trực quan hóa & hỗ trợ ra quyết định",
     ["Xây dựng mô hình dữ liệu phân tích (Star Schema) trên Microsoft Fabric",
      "Xuất ra Dashboard quản trị trực quan qua Power BI Enterprise",
      "Truy cập đa thiết bị: Laptop, iPad, Điện thoại di động",
      "Phân quyền truy cập theo vai trò (RBAC): Lãnh đạo, Phòng ban, Chuyên viên"])
]

for i, (title, sub, col, bg_c, role, items) in enumerate(layers):
    x = Inches(0.8 + i * 4.05)
    shape(s3, x, Inches(2.3), Inches(3.7), Inches(4.65), bg_c, col, Pt(1.5), rounded=True)
    shape(s3, x, Inches(2.3), Inches(3.7), Inches(0.65), col)
    text(s3, x, Inches(2.33), Inches(3.7), Inches(0.28), title, 11, True, WHITE, PP_ALIGN.CENTER)
    text(s3, x, Inches(2.6), Inches(3.7), Inches(0.22), sub, 9, False, RGBColor(220, 235, 255), PP_ALIGN.CENTER)

    shape(s3, x + Inches(0.15), Inches(3.05), Inches(3.4), Inches(0.35), WHITE, col, Pt(1.0), rounded=True)
    text(s3, x + Inches(0.15), Inches(3.08), Inches(3.4), Inches(0.28), role, 9.5, True, col, PP_ALIGN.CENTER)

    for j, item in enumerate(items):
        text(s3, x + Inches(0.2), Inches(3.55 + j * 0.8), Inches(3.3), Inches(0.75), "• " + item, 10, color=BLACK)

note(s3, "Kiến trúc 3 tầng theo chuẩn IMIP Data Platform của TCT: Thu thập (Airbyte/Debezium) → Chuẩn hóa (Lakehouse MinIO + MDM) → Phân tích (Fabric + Power BI). Các phần mềm nghiệp vụ hiện hữu giữ nguyên, nhân viên không cần thay đổi quy trình làm việc.")

# =========================================================================
# SLIDE 4: TCT ĐÃ ĐẦU TƯ NHỮNG GÌ?
# =========================================================================
s4 = prs.slides.add_slide(BL)
header(s4, "TÀI NGUYÊN SẴN CÓ: TỔNG CÔNG TY ĐÃ ĐẦU TƯ VÀ BÀN GIAO NHỮNG GÌ?", "PHẦN A: TỔNG QUAN GIẢI PHÁP — HẠ TẦNG HYBRID IMIP", "03")

shape(s4, Inches(0.8), Inches(1.4), Inches(5.6), Inches(4.3), BLUE_SOFT, BLUE, Pt(1.5), rounded=True)
shape(s4, Inches(0.8), Inches(1.4), Inches(5.6), Inches(0.5), BLUE)
text(s4, Inches(0.8), Inches(1.42), Inches(5.6), Inches(0.4), "HẠ TẦNG CLOUD — MICROSOFT AZURE", 12, True, WHITE, PP_ALIGN.CENTER)

cloud_items = [
    ("Microsoft Fabric / OneLake (20 TB khởi điểm)", "Kho dữ liệu đám mây tối ưu cho phân tích quy mô lớn, hỗ trợ AI/ML và mở rộng linh hoạt."),
    ("Workspace L3 riêng biệt cho Quảng Ngãi", "Phân vùng cách ly hoàn toàn — dữ liệu của QN không trộn lẫn với các đơn vị thành viên khác."),
    ("Bản quyền Power BI Enterprise (70 người dùng)", "Đã bao gồm license khai thác báo cáo quản trị cho lãnh đạo và cán bộ nghiệp vụ."),
    ("Microsoft Purview", "Quản trị danh mục dữ liệu (Data Catalog), phân loại và gắn nhãn bảo mật tự động.")
]
for i, (title, desc) in enumerate(cloud_items):
    y = Inches(2.05 + i * 0.88)
    shape(s4, Inches(1.0), y, Inches(5.2), Inches(0.78), WHITE, GRAY_MED, Pt(0.5), rounded=True)
    text(s4, Inches(1.15), y + Inches(0.04), Inches(4.9), Inches(0.25), title, 10.5, True, BLUE)
    text(s4, Inches(1.15), y + Inches(0.27), Inches(4.9), Inches(0.45), desc, 9.5, color=BLACK)

shape(s4, Inches(6.9), Inches(1.4), Inches(5.6), Inches(4.3), GRAY_LIGHT, NAVY, Pt(1.5), rounded=True)
shape(s4, Inches(6.9), Inches(1.4), Inches(5.6), Inches(0.5), NAVY)
text(s4, Inches(6.9), Inches(1.42), Inches(5.6), Inches(0.4), "HẠ TẦNG ON-PREMISE — DATACENTER TCT", 12, True, WHITE, PP_ALIGN.CENTER)

onprem_items = [
    ("Hồ dữ liệu MinIO Lakehouse (Apache Iceberg)", "Lưu trữ on-premise tại Datacenter TCT, tuân thủ Luật Dữ liệu — dữ liệu gốc nằm trong lãnh thổ Việt Nam."),
    ("Trục tích hợp doanh nghiệp ESB (WSO2)", "Truyền nhận thông điệp thời gian thực giữa 8 phần mềm nguồn và 35 quy trình liên phòng ban."),
    ("Hệ thống MDM — 29 Danh mục & 50 API chuẩn", "Quản trị Bản ghi vàng (Golden Record); 50 API chuẩn đã sẵn sàng để đơn vị kết nối."),
    ("Trung tâm Giám sát ATTT (SIEM & SOC)", "Giám sát an toàn thông tin 24/7, tuân thủ kiến trúc 8 Vùng mạng (Zones) của TCT.")
]
for i, (title, desc) in enumerate(onprem_items):
    y = Inches(2.05 + i * 0.88)
    shape(s4, Inches(7.1), y, Inches(5.2), Inches(0.78), WHITE, GRAY_MED, Pt(0.5), rounded=True)
    text(s4, Inches(7.25), y + Inches(0.04), Inches(4.9), Inches(0.25), title, 10.5, True, NAVY)
    text(s4, Inches(7.25), y + Inches(0.27), Inches(4.9), Inches(0.45), desc, 9.5, color=BLACK)

shape(s4, Inches(0.8), Inches(5.9), Inches(11.7), Inches(0.95), GOLD_SOFT, GOLD, Pt(1.5), rounded=True)
text(s4, Inches(1.1), Inches(5.95), Inches(11.1), Inches(0.85),
     "KẾT LUẬN: Tổng công ty đã hoàn tất đầu tư hạ tầng Hybrid (Cloud + On-premise) và bộ tiêu chuẩn kỹ thuật.\n"
     "PTSC Quảng Ngãi không phải mua sắm máy chủ lưu trữ — chỉ cần xây dựng Trạm trung chuyển (Spoke) để đấu nối vào hệ sinh thái TCT.",
     11, True, NAVY, PP_ALIGN.CENTER)

note(s4, "Theo tài liệu Hội thảo TCT: Nền tảng IMIP đã tích hợp 8 phần mềm nguồn, 35 quy trình liên phòng ban và 50 API chuẩn. Quảng Ngãi kế thừa toàn bộ bộ triển khai mẫu này.")

# =========================================================================
# SLIDE 5: VỊ THẾ QUẢNG NGÃI — LEVEL 3 SPOKE
# =========================================================================
s5 = prs.slides.add_slide(BL)
header(s5, "VỊ THẾ PTSC QUẢNG NGÃI: LEVEL 3 TRONG MÔ HÌNH HUB-SPOKE", "PHẦN B: HIỆN TRẠNG TẠI QUẢNG NGÃI — PHÂN LOẠI ĐƠN VỊ", "04")

levels = [
    ("L1 — CHI NHÁNH", "Vận hành như một Ban TCT.\nDùng chung toàn bộ hệ thống,\nkhông dựng Spoke riêng.", GRAY, GRAY_LIGHT, False),
    ("L2 — ĐƠN VỊ NHỎ", "Tenant trên Hub (không dựng Spoke).\nChỉ cần Agent thu thập hoặc\nkéo dữ liệu trực tiếp qua API.", GRAY, GRAY_LIGHT, False),
    ("L3 — ĐƠN VỊ LỚN\n(PTSC QUẢNG NGÃI)", "Được cấp Workspace L3 riêng biệt.\nKhông cần đầu tư hạ tầng dHCI.\nSpoke tái sử dụng nền tảng IMIP\n(cấu hình quy mô nhỏ) — đồng nhất\ncông nghệ và kỹ năng vận hành.", BLUE, BLUE_SOFT, True),
    ("L4 — ĐẶC BIỆT LỚN\n(PTSC M&C)", "Tự đầu tư cụm dHCI riêng tại chỗ.\nVận hành Data Platform đầy đủ\ntại đơn vị.", GRAY, GRAY_LIGHT, False),
]

for i, (lvl, desc, col, bg_c, highlight) in enumerate(levels):
    x = Inches(0.8 + i * 2.95)
    border_w = Pt(2.5) if highlight else Pt(1.0)
    shape(s5, x, Inches(1.5), Inches(2.75), Inches(4.5), bg_c, col, border_w, rounded=True)
    shape(s5, x, Inches(1.5), Inches(2.75), Inches(0.55), col)
    text(s5, x, Inches(1.55), Inches(2.75), Inches(0.45), lvl.split('\n')[0], 12, True, WHITE, PP_ALIGN.CENTER)
    if '\n' in lvl:
        text(s5, x + Inches(0.1), Inches(2.15), Inches(2.55), Inches(0.5), lvl.split('\n')[1], 11, True, col, PP_ALIGN.CENTER)
        text(s5, x + Inches(0.1), Inches(2.75), Inches(2.55), Inches(3.0), desc, 10, color=BLACK)
    else:
        text(s5, x + Inches(0.1), Inches(2.2), Inches(2.55), Inches(3.5), desc, 10, color=BLACK)

shape(s5, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.75), GREEN_SOFT, GREEN, Pt(1.5), rounded=True)
text(s5, Inches(1.1), Inches(6.25), Inches(11.1), Inches(0.65),
     "Là Level 3, PTSC Quảng Ngãi vừa có không gian lưu trữ độc lập (Workspace L3), vừa không phát sinh chi phí đầu tư\n"
     "phần cứng (CAPEX). Chi phí vận hành nền tảng phân bổ theo mức sử dụng thực tế (Usage-based).",
     11, True, GREEN, PP_ALIGN.CENTER)

note(s5, "Theo slide 7 của tài liệu TCT: Spoke tái sử dụng chính nền tảng IMIP Data Platform ở cấu hình quy mô nhỏ — đồng nhất công nghệ, kỹ năng vận hành và bộ triển khai mẫu khi nhân rộng. L3 không phát sinh CAPEX hạ tầng riêng.")

# =========================================================================
# SLIDE 6: HIỆN TRẠNG DỮ LIỆU TẠI QUẢNG NGÃI — 3 THÁCH THỨC
# =========================================================================
s6 = prs.slides.add_slide(BL)
header(s6, "HIỆN TRẠNG DỮ LIỆU TẠI PTSC QUẢNG NGÃI: 3 THÁCH THỨC CẦN GIẢI QUYẾT", "PHẦN B: HIỆN TRẠNG TẠI QUẢNG NGÃI — PHÂN TÍCH NỘI BỘ", "05")

silos = ["Phần mềm\nKẾ TOÁN", "Phần mềm\nNHÂN SỰ", "Phần mềm\nVẬT TƯ", "Quản lý\nDỰ ÁN"]
for i, name in enumerate(silos):
    x = Inches(0.8 + i * 2.95)
    shape(s6, x, Inches(1.4), Inches(2.55), Inches(1.0), WHITE, GRAY_MED, Pt(0.8), rounded=True)
    shape(s6, x + Inches(0.04), Inches(1.42), Inches(0.08), Inches(0.96), BLUE)
    text(s6, x + Inches(0.2), Inches(1.48), Inches(2.1), Inches(0.8), name, 12, True, NAVY, PP_ALIGN.CENTER)
    if i < 3:
        text(s6, Inches(3.1 + i * 2.95), Inches(1.6), Inches(0.5), Inches(0.4), "✕", 16, True, RED, PP_ALIGN.CENTER)

text(s6, Inches(0.8), Inches(2.6), Inches(11.7), Inches(0.3),
     "Ba thách thức cốt lõi cản trở việc kết nối với Nền tảng Dữ liệu Tổng công ty:", 12.5, True, RED)

knots = [
    ("DỮ LIỆU PHÂN MẢNH (DATA SILO)", RED, RED_SOFT,
     "4 hệ thống nghiệp vụ vận hành độc lập.\nTrao đổi số liệu bằng file Excel thủ công.\nThời gian tổng hợp báo cáo: 3–5 ngày.\nSố liệu giữa các phòng ban thường lệch nhau."),
    ("CHƯA ĐỒNG BỘ 29 DANH MỤC DỮ LIỆU CHỦ", GOLD, GOLD_SOFT,
     "Hệ thống mã nhân sự, vật tư, đối tác nội bộ\nchưa ánh xạ với 29 danh mục chuẩn TCT.\nNếu truyền dữ liệu trực tiếp lên Hub,\nhệ thống MDM của TCT sẽ từ chối tiếp nhận."),
    ("THIẾU CHUYÊN MÔN DATA ENGINEERING", BLUE, BLUE_SOFT,
     "Đội ngũ IT nội bộ có năng lực về hạ tầng\nmạng và quản trị hệ thống máy chủ.\nTuy nhiên chưa có kinh nghiệm xây dựng\nluồng dữ liệu tự động (ETL Pipeline)."),
]
for i, (title, col, bg_c, desc) in enumerate(knots):
    x = Inches(0.8 + i * 4.05)
    shape(s6, x, Inches(3.05), Inches(3.7), Inches(3.85), bg_c, col, Pt(1.5), rounded=True)
    shape(s6, x + Inches(0.04), Inches(3.07), Inches(0.08), Inches(3.79), col)
    numbered_label(s6, x + Inches(0.3), Inches(3.25), i + 1, col)
    text(s6, x + Inches(0.85), Inches(3.25), Inches(2.6), Inches(0.4), title, 11.5, True, col)
    shape(s6, x + Inches(0.3), Inches(3.8), Inches(3.1), Inches(0.01), GRAY_MED)
    text(s6, x + Inches(0.3), Inches(3.95), Inches(3.1), Inches(2.7), desc, 10.5, color=BLACK)

note(s6, "Ba thách thức cần giải quyết: phần mềm cát cứ, chưa đồng bộ Master Data, và đội IT chưa có chuyên môn Data Engineering. Đây là lý do cần thuê NCC chuyên nghiệp hỗ trợ.")

# =========================================================================
# SLIDE 7: QUY CHẾ QUẢN TRỊ DỮ LIỆU & CHỦ QUYỀN ĐƠN VỊ
# =========================================================================
s7 = prs.slides.add_slide(BL)
header(s7, "QUY CHẾ QUẢN TRỊ DỮ LIỆU PTSC: MÔ HÌNH 5 CẤP VÀ QUYỀN CỦA ĐƠN VỊ", "PHẦN B: HIỆN TRẠNG TẠI QUẢNG NGÃI — CHỦ QUYỀN DỮ LIỆU", "06")

gov_levels = [
    ("CẤP 1–2: HỘI ĐỒNG QUẢN TRỊ DỮ LIỆU & HỘI ĐỒNG DỮ LIỆU KHỐI", "Định hướng chiến lược dữ liệu toàn Tổng công ty. Phê duyệt chính sách, tiêu chuẩn. Điều phối liên Ban, đồng bộ tiêu chuẩn giữa các miền.", GRAY, GRAY_LIGHT, False),
    ("CẤP 3: CHỦ QUẢN DỮ LIỆU — LÃNH ĐẠO ĐƠN VỊ THÀNH VIÊN (PTSC QUẢNG NGÃI)",
     "• Quản lý nghiệp vụ dữ liệu của miền; phê duyệt phân loại, chất lượng, truy cập, chia sẻ.\n"
     "• TRÁCH NHIỆM GIẢI TRÌNH ĐẶT Ở CẤP NÀY — Chất lượng dữ liệu là tiêu chí đánh giá.\n"
     "• Quyết định dữ liệu nào được chia sẻ, dữ liệu nào giữ nội bộ. Có quyền dừng chia sẻ nếu vi phạm quy định.\n"
     "• Phê duyệt bản ghi chuẩn khi dữ liệu giữa các hệ thống mâu thuẫn.",
     GREEN, GREEN_SOFT, True),
    ("CẤP 4: QUẢN TRỊ MIỀN DỮ LIỆU (DATA STEWARDS)", "Siêu dữ liệu, từ điển nghiệp vụ, quy tắc chất lượng; chuẩn hóa, ánh xạ; xử lý vấn đề dữ liệu tại cấp kỹ thuật.", GRAY, GRAY_LIGHT, False),
    ("CẤP 5: ĐƠN VỊ VẬN HÀNH HỆ THỐNG DỮ LIỆU (BAN NCPT&CĐS TCT)",
     "Vận hành nền tảng, phân quyền, giám sát, sao lưu. KHÔNG SỞ HỮU DỮ LIỆU NGHIỆP VỤ CỦA ĐƠN VỊ.", RED, RED_SOFT, False),
]

y_pos = Inches(1.4)
for title, desc, col, bg_c, highlight in gov_levels:
    h = Inches(2.0) if highlight else Inches(1.0)
    border_w = Pt(2.5) if highlight else Pt(1.0)
    shape(s7, Inches(0.8), y_pos, Inches(11.7), h, bg_c, col, border_w, rounded=True)
    shape(s7, Inches(0.84), y_pos + Inches(0.04), Inches(0.1), h - Inches(0.08), col)
    if highlight:
        text(s7, Inches(1.2), y_pos + Inches(0.1), Inches(11.0), Inches(0.4), title, 12, True, col)
        text(s7, Inches(1.2), y_pos + Inches(0.55), Inches(11.0), Inches(1.35), desc, 10.5, color=BLACK)
    else:
        text(s7, Inches(1.2), y_pos + Inches(0.08), Inches(11.0), Inches(0.3), title, 11, True, col)
        text(s7, Inches(1.2), y_pos + Inches(0.4), Inches(11.0), Inches(0.55), desc, 10, color=BLACK)
    y_pos += h + Inches(0.12)

note(s7, "Theo Quy chế Quản trị Dữ liệu PTSC (7 chương, 30 điều, 13 phụ lục): Chủ quản Dữ liệu Cấp 3 (Lãnh đạo Đơn vị) là cấp phê duyệt 10/12 quy trình số hóa trên Data Platform. Ban NCPT&CĐS Cấp 5 chỉ vận hành kỹ thuật.")

# =========================================================================
# SLIDE 8: BẢO MẬT — DỮ LIỆU NÀO Ở LẠI, NÀO ĐI?
# =========================================================================
s8 = prs.slides.add_slide(BL)
header(s8, "NGUYÊN TẮC BẢO MẬT: PHÂN LOẠI DỮ LIỆU NỘI BỘ VÀ DỮ LIỆU ĐỒNG BỘ", "PHẦN B: HIỆN TRẠNG TẠI QUẢNG NGÃI — AN TOÀN THÔNG TIN", "07")

shape(s8, Inches(0.8), Inches(1.4), Inches(5.6), Inches(5.4), RED_SOFT, RED, Pt(1.5), rounded=True)
shape(s8, Inches(0.8), Inches(1.4), Inches(5.6), Inches(0.5), RED)
text(s8, Inches(0.8), Inches(1.42), Inches(5.6), Inches(0.4), "DỮ LIỆU LƯU GIỮ 100% TẠI NỘI BỘ ĐƠN VỊ", 11.5, True, WHITE, PP_ALIGN.CENTER)

stay = [
    ("Dữ liệu kinh doanh mang tính cạnh tranh", "Định mức đơn giá thầu, biên lợi nhuận từng dự án, chiến lược thương mại, chi phí chi tiết từng hạng mục."),
    ("Dữ liệu vận hành chuyên ngành chi tiết", "Nhật trình thi công, nhật ký xưởng, thông số bảo dưỡng thiết bị, tiến độ thi công hàng ngày."),
    ("Dữ liệu cá nhân chưa qua xử lý", "Bảng lương chi tiết, thông tin tài chính cá nhân, hồ sơ sức khỏe CBNV — theo Luật BVDLCN 91/2025."),
    ("Toàn bộ CSDL sản xuất gốc (Production DB)", "Cơ sở dữ liệu đang vận hành đặt tại phòng máy chủ đơn vị. TCT không truy cập trực tiếp.")
]
for i, (title, desc) in enumerate(stay):
    y = Inches(2.1 + i * 1.1)
    shape(s8, Inches(1.0), y, Inches(5.2), Inches(0.95), WHITE, GRAY_MED, Pt(0.5), rounded=True)
    text(s8, Inches(1.2), y + Inches(0.06), Inches(4.8), Inches(0.28), title, 10.5, True, RED)
    text(s8, Inches(1.2), y + Inches(0.33), Inches(4.8), Inches(0.55), desc, 9.5, color=BLACK)

shape(s8, Inches(6.9), Inches(1.4), Inches(5.6), Inches(5.4), GREEN_SOFT, GREEN, Pt(1.5), rounded=True)
shape(s8, Inches(6.9), Inches(1.4), Inches(5.6), Inches(0.5), GREEN)
text(s8, Inches(6.9), Inches(1.42), Inches(5.6), Inches(0.4), "DỮ LIỆU ĐƯỢC PHÉP ĐỒNG BỘ LÊN HUB TCT", 11.5, True, WHITE, PP_ALIGN.CENTER)

send = [
    ("29 Danh mục Dữ liệu chủ (Master Data)", "Mã phòng ban, mã chức danh, mã vật tư, mã dự án... đã ánh xạ theo chuẩn thống nhất toàn Tổng công ty."),
    ("Số liệu tổng hợp (Aggregated Data)", "Doanh thu tổng, sản lượng tổng, số lượng lao động — phục vụ hợp nhất báo cáo toàn Tổng."),
    ("Dữ liệu đã che mờ (Data Masking)", "Thông tin cá nhân (CCCD, SĐT) được mã hóa/làm mờ tự động trước khi truyền, tuân thủ NĐ 13."),
    ("Phân vùng Tenant L3 cách ly hoàn toàn", "Dữ liệu QN chuyển vào Workspace riêng; các ĐVTV khác không có quyền truy cập chéo.")
]
for i, (title, desc) in enumerate(send):
    y = Inches(2.1 + i * 1.1)
    shape(s8, Inches(7.1), y, Inches(5.2), Inches(0.95), WHITE, GRAY_MED, Pt(0.5), rounded=True)
    text(s8, Inches(7.3), y + Inches(0.06), Inches(4.8), Inches(0.28), title, 10.5, True, GREEN)
    text(s8, Inches(7.3), y + Inches(0.33), Inches(4.8), Inches(0.55), desc, 9.5, color=BLACK)

note(s8, "Phân loại 6 chiều theo Quy chế TCT: Tính chất chia sẻ, mức độ quan trọng, tính chất bí mật, dữ liệu cá nhân, nguồn gốc, vòng đời. Mỗi tài sản dữ liệu mang 6 nhãn, Cấp 3 (Lãnh đạo đơn vị) phê duyệt phần lớn.")

# =========================================================================
# SLIDE 9: KIẾN TRÚC TRỤC TÍCH HỢP — LUỒNG DỮ LIỆU 4 BƯỚC
# =========================================================================
s9 = prs.slides.add_slide(BL)
header(s9, "KIẾN TRÚC TRỤC TÍCH HỢP: LUỒNG DỮ LIỆU TỪ QUẢNG NGÃI ĐẾN HUB TCT", "PHẦN C: KIẾN TRÚC KẾT NỐI — SƠ ĐỒ KỸ THUẬT TỔNG THỂ", "08")

pipe = [
    ("BƯỚC 1\nNGUỒN DỮ LIỆU", "NỘI BỘ QUẢNG NGÃI", BLUE, BLUE_SOFT,
     ["4 Phần mềm nghiệp vụ:", "• Kế toán (FAST/Bravo)", "• Nhân sự (HRM)", "• Quản lý Vật tư", "• Quản lý Dự án", "", "Giữ nguyên hiện trạng,", "CBNV làm việc như cũ."]),
    ("BƯỚC 2\nTRẠM TRUNG CHUYỂN", "VM TẠI DATACENTER QN", TEAL, RGBColor(230, 248, 246),
     ["Xử lý tự động (ETL):", "• Trích xuất bản sao (Read-Only)", "• Làm sạch, khử trùng lặp", "• Ánh xạ 29 Danh mục MD", "• Che mờ dữ liệu cá nhân", "• Mã hóa AES-256", "", "Chạy ngầm theo lịch (23h00)."]),
    ("BƯỚC 3\nKÊNH TRUYỀN", "VPN SITE-TO-SITE MÃ HÓA", GOLD, GOLD_SOFT,
     ["Hạ tầng mạng bảo mật:", "• Đường hầm IPSec/SSL", "• DC Quảng Ngãi ↔ DC TCT", "• Tuân thủ 8 Zones ATTT", "• Không lộ IP ra Internet", "• Tự động đồng bộ định kỳ", "", "IT QN + IT TCT phối hợp."]),
    ("BƯỚC 4\nHUB TỔNG CÔNG TY", "WORKSPACE L3 & POWER BI", GREEN, GREEN_SOFT,
     ["Tiếp nhận & khai thác:", "• Lưu trữ: MinIO On-prem", "• Phân tích: Fabric OneLake", "• MDM kiểm tra Golden Record", "• Kích hoạt Dashboard", "  Power BI tự động", "", "Ban Giám đốc xem tức thời."])
]
for i, (st_name, st_sub, col, bg_c, lines) in enumerate(pipe):
    x = Inches(0.6 + i * 3.1)
    shape(s9, x, Inches(1.5), Inches(2.75), Inches(5.4), bg_c, col, Pt(1.5), rounded=True)
    shape(s9, x, Inches(1.5), Inches(2.75), Inches(0.65), col)
    text(s9, x, Inches(1.53), Inches(2.75), Inches(0.55), st_name, 10.5, True, WHITE, PP_ALIGN.CENTER)
    text(s9, x, Inches(2.15), Inches(2.75), Inches(0.22), st_sub, 8.5, True, RGBColor(220, 235, 255), PP_ALIGN.CENTER)
    for j, l in enumerate(lines):
        bold_line = True if j == 0 or j == len(lines)-1 else False
        c = col if bold_line else BLACK
        text(s9, x + Inches(0.12), Inches(2.5 + j * 0.45), Inches(2.5), Inches(0.4), l, 10, bold_line, c)
    if i < 3:
        arrow_shape(s9, Inches(3.42 + i * 3.1), Inches(4.0), "right", Inches(0.3), Inches(0.22))

note(s9, "Đồng bộ qua ESB / API chuẩn (near real-time hoặc batch) theo kiến trúc TCT. Spoke tái sử dụng bộ triển khai mẫu IMIP Data Platform.")

# =========================================================================
# SLIDE 10: SO SÁNH 2 PHƯƠNG ÁN
# =========================================================================
s10 = prs.slides.add_slide(BL)
header(s10, "SO SÁNH HAI PHƯƠNG ÁN TRIỂN KHAI THEO CHỈ ĐẠO BAN GIÁM ĐỐC", "PHẦN C: KIẾN TRÚC KẾT NỐI — LỰA CHỌN TỐI ƯU", "09")

shape(s10, Inches(0.8), Inches(1.4), Inches(5.6), Inches(5.4), RED_SOFT, RED, Pt(1.5), rounded=True)
shape(s10, Inches(0.8), Inches(1.4), Inches(5.6), Inches(0.5), RED)
text(s10, Inches(0.8), Inches(1.42), Inches(5.6), Inches(0.4), "PHƯƠNG ÁN 1: CHỜ TỔNG CÔNG TY TRIỂN KHAI", 12, True, WHITE, PP_ALIGN.CENTER)

pa1 = [
    ("Khả thi về nhân lực", "KHÔNG ĐẢM BẢO", RED, "Ban NCPT&CĐS TCT vận hành Hub chung, không đủ nhân sự cắm chốt tại từng đơn vị để bóc tách phần mềm nội bộ."),
    ("Tiến độ thực hiện", "RỦI RO TRỄ HẠN CAO", RED, "Hơn 10 ĐVTV đang xếp hàng. Không thể đảm bảo mốc KPI năm 2026 do TCT giao nếu chờ lần lượt."),
    ("Đáp ứng nghiệp vụ nội bộ", "CHỈ PHỤC VỤ BÁO CÁO TCT", RED, "TCT chỉ xây báo cáo hợp nhất phục vụ TCT. Không xây dựng Dashboard quản trị chi tiết cho Ban Giám đốc QN."),
    ("Mức độ tự chủ", "PHỤ THUỘC HOÀN TOÀN", RED, "Không kiểm soát được tiến độ, mã nguồn và lịch nâng cấp. Mọi điều chỉnh phải chờ TCT phê duyệt.")
]
for i, (crit, val, v_col, desc) in enumerate(pa1):
    y = Inches(2.1 + i * 1.1)
    shape(s10, Inches(1.0), y, Inches(5.2), Inches(0.95), WHITE, GRAY_MED, Pt(0.5), rounded=True)
    text(s10, Inches(1.15), y + Inches(0.05), Inches(2.5), Inches(0.25), crit, 10, True, BLACK)
    text(s10, Inches(3.5), y + Inches(0.05), Inches(2.5), Inches(0.25), val, 10, True, v_col, PP_ALIGN.RIGHT)
    text(s10, Inches(1.15), y + Inches(0.32), Inches(4.9), Inches(0.55), desc, 9.5, color=BLACK)

shape(s10, Inches(6.9), Inches(1.4), Inches(5.6), Inches(5.4), GREEN_SOFT, GREEN, Pt(2.0), rounded=True)
shape(s10, Inches(6.9), Inches(1.4), Inches(5.6), Inches(0.5), GREEN)
text(s10, Inches(6.9), Inches(1.42), Inches(5.6), Inches(0.4), "PHƯƠNG ÁN 2 (ĐỀ XUẤT): CHỦ ĐỘNG + NCC CHUYÊN MÔN", 11, True, WHITE, PP_ALIGN.CENTER)

pa2 = [
    ("Khả thi về nhân lực", "NCC CHUYÊN MÔN TẬN NƠI", GREEN, "NCC có kinh nghiệm Data Engineering cử kỹ sư cắm chốt tại QN, phối hợp trực tiếp với IT và các phòng ban."),
    ("Tiến độ thực hiện", "3–4 THÁNG, CHỦ ĐỘNG", GREEN, "Chủ động kiểm soát tiến độ hoàn toàn. Kịp mốc KPI năm 2026 và không phụ thuộc lịch trình của TCT."),
    ("Đáp ứng nghiệp vụ nội bộ", "THIẾT KẾ RIÊNG CHO QN", GREEN, "Vừa đồng bộ dữ liệu lên Hub TCT, vừa xây Dashboard quản trị phục vụ trực tiếp nhu cầu điều hành nội bộ."),
    ("Mức độ tự chủ", "BÀN GIAO MÃ NGUỒN 100%", GREEN, "Bàn giao toàn bộ source code, tài liệu kiến trúc. Đào tạo IT QN làm chủ vận hành. Tận dụng license miễn phí từ TCT.")
]
for i, (crit, val, v_col, desc) in enumerate(pa2):
    y = Inches(2.1 + i * 1.1)
    shape(s10, Inches(7.1), y, Inches(5.2), Inches(0.95), WHITE, GRAY_MED, Pt(0.5), rounded=True)
    text(s10, Inches(7.25), y + Inches(0.05), Inches(2.4), Inches(0.25), crit, 10, True, BLACK)
    text(s10, Inches(9.5), y + Inches(0.05), Inches(2.7), Inches(0.25), val, 10, True, v_col, PP_ALIGN.RIGHT)
    text(s10, Inches(7.25), y + Inches(0.32), Inches(4.9), Inches(0.55), desc, 9.5, color=BLACK)

note(s10, "So sánh 2 phương án: PA1 phụ thuộc TCT — bất khả thi về tiến độ và không có Dashboard riêng. PA2 chủ động thuê NCC — kiểm soát tiến độ, thiết kế riêng cho QN, bàn giao mã nguồn.")

# =========================================================================
# SLIDE 11: CƠ CẤU CHI PHÍ
# =========================================================================
s11 = prs.slides.add_slide(BL)
header(s11, "CƠ CẤU ĐẦU TƯ VÀ DỰ TOÁN NGÂN SÁCH (THEO KHUNG CHI PHÍ CỦA TCT)", "PHẦN C: KIẾN TRÚC KẾT NỐI — TÀI CHÍNH", "10")

cost_rows = [
    ("1", "Hạ tầng máy chủ tại đơn vị", "Tận dụng máy chủ ảo (VM) sẵn có tại Datacenter QN", "0 VNĐ (Tự có)", GREEN),
    ("2", "Hạ tầng Cloud Hub & Bản quyền khung", "Microsoft Fabric 20TB, MinIO, Power BI Enterprise, Purview", "TCT đã đầu tư", BLUE),
    ("3", "Dịch vụ: Khảo sát & Thiết kế kiến trúc", "Khảo sát CSDL 4 phần mềm, lập bảng ánh xạ 29 Master Data", "[Chờ NCC báo giá]", GOLD),
    ("4", "Dịch vụ: Xây dựng Trục tích hợp (ETL)", "Lập trình luồng trích xuất, chuẩn hóa, mã hóa và truyền VPN", "[Chờ NCC báo giá]", GOLD),
    ("5", "Dịch vụ: Xây dựng Dashboard Power BI", "Thiết kế mô hình dữ liệu và các bảng báo cáo quản trị", "[Chờ NCC báo giá]", GOLD),
    ("6", "Đào tạo chuyển giao & Bảo hành", "Bàn giao mã nguồn, đào tạo IT QN vận hành độc lập", "[Chờ NCC báo giá]", GOLD),
    ("7", "Chi phí vận hành nền tảng hàng năm", "TCT phân bổ theo mức sử dụng thực tế (Usage-based)", "Theo cơ chế TCT", GRAY),
]

shape(s11, Inches(0.8), Inches(1.4), Inches(11.7), Inches(0.4), NAVY)
text(s11, Inches(0.9), Inches(1.45), Inches(0.6), Inches(0.3), "STT", 10, True, WHITE, PP_ALIGN.CENTER)
text(s11, Inches(1.6), Inches(1.45), Inches(3.2), Inches(0.3), "CẤU PHẦN CHI PHÍ", 10, True, WHITE)
text(s11, Inches(4.9), Inches(1.45), Inches(4.5), Inches(0.3), "NỘI DUNG THỰC HIỆN", 10, True, WHITE)
text(s11, Inches(9.5), Inches(1.45), Inches(2.9), Inches(0.3), "DỰ TOÁN", 10, True, WHITE, PP_ALIGN.CENTER)

for i, (stt, comp, desc, cost, c_col) in enumerate(cost_rows):
    y = Inches(1.85 + i * 0.58)
    bg = WHITE if i % 2 == 0 else GRAY_LIGHT
    shape(s11, Inches(0.8), y, Inches(11.7), Inches(0.53), bg, GRAY_MED, Pt(0.5))
    text(s11, Inches(0.9), y + Inches(0.1), Inches(0.6), Inches(0.3), stt, 10, True, BLACK, PP_ALIGN.CENTER)
    text(s11, Inches(1.6), y + Inches(0.1), Inches(3.2), Inches(0.3), comp, 10, True, BLACK)
    text(s11, Inches(4.9), y + Inches(0.1), Inches(4.5), Inches(0.3), desc, 9.5, color=BLACK)
    text(s11, Inches(9.5), y + Inches(0.1), Inches(2.9), Inches(0.3), cost, 10, True, c_col, PP_ALIGN.CENTER)

shape(s11, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.9), GOLD_SOFT, GOLD, Pt(1.0), rounded=True)
text(s11, Inches(1.1), Inches(6.05), Inches(11.1), Inches(0.8),
     "NGUYÊN TẮC TÀI CHÍNH: Tiết kiệm chi phí hạ tầng và bản quyền nhờ tận dụng đầu tư sẵn có của TCT.\n"
     "Kinh phí dịch vụ NCC sẽ được xác định chính xác sau khi hoàn thành Giai đoạn 1 (khảo sát kỹ thuật nội bộ) và nhận báo giá cạnh tranh.",
     10.5, True, NAVY, PP_ALIGN.CENTER)

note(s11, "7 cấu phần chi phí theo khung TCT. L3 không phát sinh CAPEX hạ tầng riêng. Hạng mục 3–6 chờ NCC báo giá sau khảo sát GĐ1.")

# =========================================================================
# SLIDE 12: KẾ HOẠCH GIAI ĐOẠN 1 CHI TIẾT — PTSC QN CẦN LÀM GÌ
# =========================================================================
s12 = prs.slides.add_slide(BL)
header(s12, "GIAI ĐOẠN 1 CHI TIẾT: KHẢO SÁT, CHUẨN HÓA DỮ LIỆU CHỦ & LẬP ĐỀ BÀI", "PHẦN D: KẾ HOẠCH HÀNH ĐỘNG — GĐ1 NỘI BỘ TỰ CHỦ TRÌ (CHI PHÍ 0 VNĐ)", "11")

gd1_tasks = [
    ("1.1", "Thành lập Tổ Công tác Nền tảng Dữ liệu", BLUE, BLUE_SOFT,
     ["IT chủ trì kỹ thuật; Key Users từ các phòng Kế toán, Nhân sự, Vật tư, Dự án",
      "Ban Giám đốc ký Quyết định thành lập — là cơ sở để yêu cầu phối hợp liên phòng ban",
      "Giao Tổ công tác lập kế hoạch và điều phối toàn bộ Giai đoạn 1"]),
    ("1.2", "Lập Danh mục nguồn dữ liệu (Data Catalog)", TEAL, RGBColor(230, 248, 246),
     ["Thống kê tất cả phần mềm đang vận hành: tên, CSDL (SQL Server/MySQL/Oracle), vị trí máy chủ",
      "Xác minh quyền truy cập: Có tài khoản Admin vào DB không? Có cần liên hệ NCC phần mềm mở API?",
      "Kết quả: File Excel Data Catalog hoàn chỉnh, sẵn sàng bàn giao cho NCC ở Giai đoạn 2"]),
    ("1.3", "Lập Bảng ánh xạ 29 Danh mục Dữ liệu chủ", GOLD, GOLD_SOFT,
     ["IT xuất (export) toàn bộ bảng mã hiện hành ra Excel: mã nhân viên, mã vật tư, mã phòng ban...",
      "Key Users đối chiếu từng mã với tài liệu 29 danh mục TCT ban hành, lập quy tắc quy đổi",
      "VD: Mã phòng ban QN đang là 'P.TCKT' → Chuẩn TCT: 'QN-BAN-TCKT'. Bảng này là Từ điển cho ETL"]),
    ("1.4", "Phân loại dữ liệu nhạy cảm & Lập TOR", NAVY, GRAY_LIGHT,
     ["Rà soát các trường dữ liệu cần che mờ (Masking): lương, SĐT, CCCD... → Đánh dấu [Cần Masking]",
      "Tổng hợp kết quả 1.2 + 1.3 + phân loại → Lập Hồ sơ yêu cầu kỹ thuật (TOR) chuyên nghiệp",
      "Sử dụng TOR để mời 2–3 NCC gửi đề xuất kỹ thuật và báo giá cạnh tranh trước khi trình BGĐ"])
]

for i, (num, title, col, bg_c, items) in enumerate(gd1_tasks):
    y = Inches(1.4 + i * 1.4)
    shape(s12, Inches(0.8), y, Inches(11.7), Inches(1.3), bg_c, col, Pt(1.2), rounded=True)
    shape(s12, Inches(0.84), y + Inches(0.04), Inches(0.1), Inches(1.22), col)
    text(s12, Inches(1.15), y + Inches(0.08), Inches(1.0), Inches(0.3), f"Bước {num}", 11, True, col)
    text(s12, Inches(2.2), y + Inches(0.08), Inches(10.0), Inches(0.3), title, 11.5, True, BLACK)
    for j, item in enumerate(items):
        text(s12, Inches(1.15), y + Inches(0.4 + j * 0.28), Inches(11.1), Inches(0.25), "• " + item, 9.5, color=BLACK)

note(s12, "GĐ1 nội bộ tự làm, chi phí 0 VNĐ, thời gian 6–8 tuần. Kết quả đầu ra: Data Catalog, Bảng ánh xạ 29 MD, Danh sách trường cần Masking, và TOR để mời NCC báo giá.")

# =========================================================================
# SLIDE 13: KẾ HOẠCH GIAI ĐOẠN 2 CHI TIẾT — XÂY DỰNG ETL PIPELINE
# =========================================================================
s13 = prs.slides.add_slide(BL)
header(s13, "GIAI ĐOẠN 2 CHI TIẾT: XÂY DỰNG TRỤC TÍCH HỢP DỮ LIỆU (ETL PIPELINE)", "PHẦN D: KẾ HOẠCH HÀNH ĐỘNG — GĐ2 TRIỂN KHAI CÓ NCC CHUYÊN MÔN", "12")

gd2_tasks = [
    ("2.1", "Cấp phát hạ tầng & Thiết lập môi trường", BLUE, BLUE_SOFT,
     ["Tạo 1–2 VM chuyên dụng tại Datacenter QN (8 Core, 32GB RAM, 500GB SSD)",
      "Đặt VM vào phân vùng mạng Internal Service Zone (tuân thủ 8 Zones ATTT)",
      "Cài đặt bộ công cụ ETL: Airbyte (thu thập), dbt (biến đổi), Apache Airflow (điều phối)"]),
    ("2.2", "Xây dựng luồng Trích xuất dữ liệu (Extract)", TEAL, RGBColor(230, 248, 246),
     ["Cấu hình Connector kết nối vào từng CSDL nghiệp vụ (JDBC/ODBC hoặc API)",
      "Thiết lập chế độ CHỈ ĐỌC (Read-Only Replica) — không ảnh hưởng hệ thống sản xuất",
      "Lập lịch trích xuất tự động: Batch hàng đêm (23h00) hoặc CDC near real-time tùy nguồn"]),
    ("2.3", "Xây dựng luồng Chuẩn hóa & Bảo mật (Transform)", GOLD, GOLD_SOFT,
     ["Nạp Bảng ánh xạ (GĐ1) vào hệ thống → Tự động quy đổi mã QN sang mã chuẩn TCT",
      "Áp dụng quy tắc làm sạch: Khử trùng lặp, kiểm tra tính toàn vẹn (dbt + Great Expectations)",
      "Thực thi Masking: Băm (hash) hoặc cắt các trường dữ liệu cá nhân nhạy cảm trước khi truyền"]),
    ("2.4", "Kiểm thử thông luồng End-to-End & Giám sát", GREEN, GREEN_SOFT,
     ["Chạy thử toàn bộ luồng E-T-L trên dữ liệu thực, đối soát kết quả với Key Users nghiệp vụ",
      "Thiết lập hệ thống cảnh báo tự động: Gửi thông báo qua Email/Zalo nếu lỗi trích xuất hoặc đứt mạng",
      "Vận hành thử nghiệm liên tục ≥ 30 ngày — đảm bảo ổn định trước khi chuyển sang kết nối Hub"])
]
for i, (num, title, col, bg_c, items) in enumerate(gd2_tasks):
    y = Inches(1.4 + i * 1.4)
    shape(s13, Inches(0.8), y, Inches(11.7), Inches(1.3), bg_c, col, Pt(1.2), rounded=True)
    shape(s13, Inches(0.84), y + Inches(0.04), Inches(0.1), Inches(1.22), col)
    text(s13, Inches(1.15), y + Inches(0.08), Inches(1.0), Inches(0.3), f"Bước {num}", 11, True, col)
    text(s13, Inches(2.2), y + Inches(0.08), Inches(10.0), Inches(0.3), title, 11.5, True, BLACK)
    for j, item in enumerate(items):
        text(s13, Inches(1.15), y + Inches(0.4 + j * 0.28), Inches(11.1), Inches(0.25), "• " + item, 9.5, color=BLACK)

note(s13, "GĐ2 có NCC chuyên môn cắm chốt. IT QN giám sát và tiếp nhận chuyển giao kỹ thuật song song. Thời gian: 8–12 tuần. Kết quả: ETL Pipeline vận hành ổn định ≥ 30 ngày trên dữ liệu thực.")

# =========================================================================
# SLIDE 14: CHI TIẾT KẾT NỐI ETL VỚI HUB TCT
# =========================================================================
s14 = prs.slides.add_slide(BL)
header(s14, "CHI TIẾT KỸ THUẬT: KẾT NỐI TỪ TRẠM ETL TẠI QN LÊN HUB TỔNG CÔNG TY", "PHẦN D: KẾ HOẠCH HÀNH ĐỘNG — QUY TRÌNH ĐỒNG BỘ LÊN HUB", "13")

shape(s14, Inches(0.8), Inches(1.35), Inches(11.7), Inches(0.55), NAVY_MID)
text(s14, Inches(0.8), Inches(1.38), Inches(11.7), Inches(0.5),
     "Quy trình kết nối Hub TCT theo chuẩn IMIP: Đồng bộ dữ liệu dùng chung + dữ liệu tổng hợp qua ESB / API chuẩn (near real-time / batch)",
     10.5, True, WHITE, PP_ALIGN.CENTER)

conn_steps = [
    ("A", "THIẾT LẬP KÊNH TRUYỀN BẢO MẬT", BLUE,
     ["Phối hợp IT TCT cấu hình VPN Site-to-Site (IPSec) hoặc SD-WAN",
      "Endpoint: Datacenter QN (Internal Service Zone) ↔ Datacenter TCT (Tầng 10 CQTCT)",
      "Kiểm tra: Bandwidth tối thiểu 10 Mbps, Latency ≤ 50ms, mã hóa AES-256",
      "Firewall Rules: Chỉ mở cổng kết nối cần thiết (443/HTTPS, cổng MinIO S3 API)"]),
    ("B", "ĐĂNG KÝ ENDPOINT & XÁC THỰC VỚI HUB", TEAL,
     ["IT TCT cấp Service Account + API Key cho Trạm trung chuyển QN trên hệ thống ESB (WSO2)",
      "Đăng ký Endpoint trên Microsoft Purview: Khai báo metadata, schema và lịch đồng bộ",
      "Cấu hình xác thực OAuth 2.0 / mTLS giữa Spoke QN và Hub TCT",
      "Kiểm tra phân quyền: QN chỉ ghi (Write) vào Landing Zone của Workspace L3 riêng"]),
    ("C", "TRUYỀN DỮ LIỆU VÀ KIỂM CHỨNG", GOLD,
     ["Cơ chế đẩy: Trạm ETL QN đẩy file Parquet/Delta đã mã hóa lên MinIO Lakehouse (S3 API)",
      "Hoặc gọi API chuẩn trên ESB (WSO2) để truyền bản ghi Master Data theo batch hoặc CDC",
      "Hệ thống MDM Hub tự động kiểm tra Golden Record: Nếu mã hợp lệ → Chấp nhận; Sai → Reject + Log",
      "Dashboard Power BI tự động refresh dữ liệu mới từ Workspace L3 theo lịch (mỗi sáng 06h00)"]),
    ("D", "GIÁM SÁT VẬN HÀNH & XỬ LÝ SỰ CỐ", GREEN,
     ["Hệ thống SIEM/SOC TCT giám sát luồng truyền 24/7 — phát hiện bất thường tự động",
      "Airflow DAG tại QN giám sát trạng thái từng bước: Extract → Transform → Load → Verify",
      "Nếu Load thất bại: Tự động retry 3 lần → Gửi cảnh báo qua Email/Zalo cho IT QN và IT TCT",
      "Báo cáo vận hành hàng tuần: Số bản ghi đồng bộ thành công, tỷ lệ lỗi, thời gian xử lý"])
]

for i, (letter, title, col, items) in enumerate(conn_steps):
    y = Inches(2.05 + i * 1.27)
    shape(s14, Inches(0.8), y, Inches(11.7), Inches(1.18), GRAY_LIGHT, col, Pt(1.2), rounded=True)
    shape(s14, Inches(0.84), y + Inches(0.04), Inches(0.1), Inches(1.1), col)
    text(s14, Inches(1.15), y + Inches(0.06), Inches(0.5), Inches(0.3), letter, 13, True, col)
    text(s14, Inches(1.65), y + Inches(0.06), Inches(10.5), Inches(0.3), title, 11, True, BLACK)
    for j, item in enumerate(items):
        text(s14, Inches(1.15), y + Inches(0.35 + j * 0.2), Inches(11.1), Inches(0.18), "• " + item, 9, color=BLACK)

note(s14, "Chi tiết kỹ thuật kết nối Hub: VPN Site-to-Site + đăng ký Endpoint trên ESB + đẩy file Parquet lên MinIO S3 API hoặc gọi API chuẩn WSO2 + MDM tự động kiểm tra Golden Record + SIEM giám sát 24/7.")

# =========================================================================
# SLIDE 15: GIAI ĐOẠN 3 — BÀN GIAO & KHAI THÁC
# =========================================================================
s15 = prs.slides.add_slide(BL)
header(s15, "GIAI ĐOẠN 3 CHI TIẾT: BÀN GIAO, KHAI THÁC BÁO CÁO & ĐÀO TẠO CHUYỂN GIAO", "PHẦN D: KẾ HOẠCH HÀNH ĐỘNG — GĐ3 IT QUẢNG NGÃI LÀM CHỦ 100%", "14")

gd3_tasks = [
    ("3.1", "Tiếp nhận Workspace L3 & Thiết lập phân quyền (RBAC)", BLUE, BLUE_SOFT,
     ["Phối hợp IT TCT kích hoạt SSO (Single Sign-On) — CBNV dùng email @ptsc.com.vn đăng nhập Power BI",
      "Tiếp nhận quyền Admin Workspace L3 trên Microsoft Fabric",
      "Phân quyền theo vai trò: Ban GĐ xem toàn bộ; TP Kế toán chỉ xem Tài chính; TP Nhân sự chỉ xem HR"]),
    ("3.2", "Thiết kế mô hình dữ liệu & Xây dựng Dashboard Power BI", TEAL, RGBColor(230, 248, 246),
     ["Xây dựng mô hình phân tích (Star Schema) trên Fabric Lakehouse từ dữ liệu đã chuẩn hóa",
      "Thiết kế Dashboard: Tổng hợp SXKD, Doanh thu–Chi phí–Dòng tiền, Tiến độ Dự án, Biến động Nhân sự",
      "Cấu hình tự động refresh: Power BI kéo dữ liệu mới nhất từ Workspace L3 mỗi sáng 06h00"]),
    ("3.3", "Nghiệm thu kỹ thuật & Bàn giao mã nguồn", GREEN, GREEN_SOFT,
     ["Chạy nghiệm thu toàn bộ hệ thống End-to-End: Nguồn QN → ETL → VPN → Hub → Dashboard",
      "NCC bàn giao 100%: Mã nguồn script ETL (Git), Tài liệu kiến trúc, Từ điển dữ liệu (Data Dictionary)",
      "Cam kết bảo hành 6–12 tháng và hỗ trợ kỹ thuật từ xa sau bàn giao"]),
    ("3.4", "Đào tạo chuyển giao & Đưa vào vận hành chính thức", GOLD, GOLD_SOFT,
     ["Đào tạo IT QN: Vận hành ETL Pipeline, xử lý sự cố, bổ sung nguồn dữ liệu mới",
      "Hướng dẫn Ban GĐ cài đặt Power BI Mobile trên điện thoại — truy cập báo cáo mọi lúc mọi nơi",
      "Báo cáo TCT: Hoàn thành chỉ tiêu KPI Chuyển đổi số năm 2026, kèm hồ sơ nghiệm thu"])
]
for i, (num, title, col, bg_c, items) in enumerate(gd3_tasks):
    y = Inches(1.4 + i * 1.4)
    shape(s15, Inches(0.8), y, Inches(11.7), Inches(1.3), bg_c, col, Pt(1.2), rounded=True)
    shape(s15, Inches(0.84), y + Inches(0.04), Inches(0.1), Inches(1.22), col)
    text(s15, Inches(1.15), y + Inches(0.08), Inches(1.0), Inches(0.3), f"Bước {num}", 11, True, col)
    text(s15, Inches(2.2), y + Inches(0.08), Inches(10.0), Inches(0.3), title, 11.5, True, BLACK)
    for j, item in enumerate(items):
        text(s15, Inches(1.15), y + Inches(0.4 + j * 0.28), Inches(11.1), Inches(0.25), "• " + item, 9.5, color=BLACK)

note(s15, "GĐ3: Nghiệm thu, bàn giao mã nguồn 100%, đào tạo IT QN và đưa vào sử dụng chính thức. Dashboard phục vụ giao ban Ban Giám đốc hàng sáng.")

# =========================================================================
# SLIDE 16: TỔNG HỢP TIMELINE 3 GIAI ĐOẠN
# =========================================================================
s16 = prs.slides.add_slide(BL)
header(s16, "LỘ TRÌNH TỔNG THỂ: TIMELINE 3 GIAI ĐOẠN TRIỂN KHAI", "PHẦN D: KẾ HOẠCH HÀNH ĐỘNG — TIẾN ĐỘ TỔNG THỂ", "15")

phases_summary = [
    ("GIAI ĐOẠN 1", "T3 – T4/2026", "KHẢO SÁT & LẬP ĐỀ BÀI KỸ THUẬT", "NỘI BỘ TỰ CHỦ TRÌ · CHI PHÍ 0 VNĐ", GOLD, GOLD_SOFT,
     ["Thành lập Tổ Công tác Nền tảng Dữ liệu nội bộ",
      "Lập Data Catalog 4 phần mềm nghiệp vụ",
      "Lập Bảng ánh xạ 29 Danh mục Dữ liệu chủ",
      "Phân loại dữ liệu nhạy cảm cần che mờ",
      "Hoàn thiện TOR, mời 2–3 NCC gửi báo giá cạnh tranh"]),
    ("GIAI ĐOẠN 2", "T5 – T8/2026", "XÂY DỰNG TRỤC TÍCH HỢP & KẾT NỐI HUB", "THUÊ NCC CHUYÊN MÔN · IT QN GIÁM SÁT", BLUE, BLUE_SOFT,
     ["Lựa chọn NCC qua xét duyệt báo giá cạnh tranh",
      "Cấp phát VM, cài đặt bộ công cụ ETL (Airbyte + dbt + Airflow)",
      "Xây dựng luồng Extract → Transform → Load hoàn chỉnh",
      "Thiết lập VPN, đăng ký Endpoint và kết nối thông luồng lên Hub TCT",
      "Vận hành thử nghiệm ≥ 30 ngày trên dữ liệu thực"]),
    ("GIAI ĐOẠN 3", "T9/2026+", "BÀN GIAO, KHAI THÁC & ĐÀO TẠO", "IT QUẢNG NGÃI LÀM CHỦ 100%", GREEN, GREEN_SOFT,
     ["Tiếp nhận Workspace L3, thiết lập SSO và phân quyền RBAC",
      "Thiết kế Dashboard Power BI phục vụ giao ban Ban Giám đốc",
      "Nghiệm thu kỹ thuật, bàn giao 100% mã nguồn và tài liệu",
      "Đào tạo IT QN vận hành độc lập, hướng dẫn BGĐ dùng Power BI Mobile",
      "Báo cáo TCT hoàn thành KPI Chuyển đổi số năm 2026"])
]

for i, (phase, timeline, name, cost, col, bg_c, tasks) in enumerate(phases_summary):
    x = Inches(0.8 + i * 4.05)
    shape(s16, x, Inches(1.45), Inches(3.7), Inches(5.4), bg_c, col, Pt(1.5), rounded=True)
    shape(s16, x, Inches(1.45), Inches(3.7), Inches(0.75), col)
    text(s16, x, Inches(1.48), Inches(3.7), Inches(0.28), f"{phase}  ({timeline})", 11, True, WHITE, PP_ALIGN.CENTER)
    text(s16, x, Inches(1.76), Inches(3.7), Inches(0.22), name, 9.5, True, RGBColor(220, 235, 255), PP_ALIGN.CENTER)

    shape(s16, x + Inches(0.15), Inches(2.3), Inches(3.4), Inches(0.32), WHITE, col, Pt(1.0), rounded=True)
    text(s16, x + Inches(0.15), Inches(2.33), Inches(3.4), Inches(0.25), cost, 9, True, col, PP_ALIGN.CENTER)

    for j, t in enumerate(tasks):
        text(s16, x + Inches(0.15), Inches(2.8 + j * 0.78), Inches(3.4), Inches(0.72), "• " + t, 9.5, color=BLACK)

note(s16, "Lộ trình 3 giai đoạn: GĐ1 nội bộ tự làm (2 tháng, 0 VNĐ) → GĐ2 NCC triển khai ETL + kết nối Hub (4 tháng) → GĐ3 bàn giao, đào tạo, khai thác (2 tháng+).")

# =========================================================================
# SLIDE 17: CAM KẾT NGUỒN LỰC & AN TOÀN
# =========================================================================
s17 = prs.slides.add_slide(BL)
header(s17, "CAM KẾT NGUỒN LỰC CÁC PHÒNG BAN VÀ AN TOÀN THÔNG TIN KHI THUÊ NCC", "PHẦN D: KẾ HOẠCH HÀNH ĐỘNG — GIẢI TỎA CÁC MỐI QUAN NGẠI", "16")

shape(s17, Inches(0.8), Inches(1.4), Inches(5.6), Inches(5.4), GREEN_SOFT, GREEN, Pt(1.5), rounded=True)
shape(s17, Inches(0.8), Inches(1.4), Inches(5.6), Inches(0.5), GREEN)
text(s17, Inches(0.8), Inches(1.42), Inches(5.6), Inches(0.4), "CAM KẾT VỚI CÁC PHÒNG BAN NGHIỆP VỤ", 11.5, True, WHITE, PP_ALIGN.CENTER)

dept = [
    ("Không phát sinh công việc nhập liệu mới", "Hệ thống tự động đọc dữ liệu từ phần mềm hiện tại. CBNV các phòng ban không phải thay đổi quy trình hay gõ thêm dữ liệu."),
    ("Thời lượng tham gia tối thiểu", "Key Users mỗi phòng chỉ cần tham gia 2–3 buổi phỏng vấn nghiệp vụ (mỗi buổi 1–2 giờ) để giải thích cấu trúc dữ liệu và đối soát kết quả."),
    ("90% khối lượng kỹ thuật do IT & NCC đảm nhiệm", "Toàn bộ công việc lập trình, cấu hình hạ tầng, xây dựng luồng dữ liệu do bộ phận CNTT và kỹ sư NCC chịu trách nhiệm."),
    ("Các phòng ban được hưởng lợi trực tiếp", "Sau hoàn thành: Các phòng được cấp tài khoản xem báo cáo tự động của phòng mình, giảm thiểu thời gian tổng hợp Excel.")
]
for i, (title, desc) in enumerate(dept):
    y = Inches(2.1 + i * 1.1)
    shape(s17, Inches(1.0), y, Inches(5.2), Inches(0.95), WHITE, GRAY_MED, Pt(0.5), rounded=True)
    text(s17, Inches(1.2), y + Inches(0.06), Inches(4.8), Inches(0.28), title, 10.5, True, GREEN)
    text(s17, Inches(1.2), y + Inches(0.33), Inches(4.8), Inches(0.55), desc, 9.5, color=BLACK)

shape(s17, Inches(6.9), Inches(1.4), Inches(5.6), Inches(5.4), BLUE_SOFT, BLUE, Pt(1.5), rounded=True)
shape(s17, Inches(6.9), Inches(1.4), Inches(5.6), Inches(0.5), BLUE)
text(s17, Inches(6.9), Inches(1.42), Inches(5.6), Inches(0.4), "AN TOÀN THÔNG TIN & CHỐNG PHỤ THUỘC NCC", 11.5, True, WHITE, PP_ALIGN.CENTER)

sec = [
    ("Chỉ cấp quyền Đọc trên bản sao dữ liệu", "NCC được tiếp cận Read-Only Replica. Tuyệt đối không cấp quyền truy cập vào CSDL sản xuất gốc đang vận hành."),
    ("Ký Thỏa thuận bảo mật (NDA) ràng buộc pháp lý", "Cam kết bảo mật theo Luật BVDLCN 91/2025/QH15. Nghiêm cấm sao chép, chuyển giao dữ liệu ra bên ngoài."),
    ("Bàn giao 100% mã nguồn và tài liệu kiến trúc", "Hợp đồng bắt buộc: Bàn giao toàn bộ script ETL, Data Dictionary, tài liệu vận hành. IT QN hoàn toàn tự chủ sau bàn giao."),
    ("Thiết kế hệ thống 'Hai nấc độc lập'", "Dashboard nội bộ phục vụ BGĐ QN hoạt động độc lập. Nếu TCT chậm cấp cổng kết nối, tiến độ QN không bị ảnh hưởng.")
]
for i, (title, desc) in enumerate(sec):
    y = Inches(2.1 + i * 1.1)
    shape(s17, Inches(7.1), y, Inches(5.2), Inches(0.95), WHITE, GRAY_MED, Pt(0.5), rounded=True)
    text(s17, Inches(7.3), y + Inches(0.06), Inches(4.8), Inches(0.28), title, 10.5, True, BLUE)
    text(s17, Inches(7.3), y + Inches(0.33), Inches(4.8), Inches(0.55), desc, 9.5, color=BLACK)

note(s17, "Giải tỏa 2 mối lo: Các phòng ban không bị quá tải (chỉ 2-3 buổi phỏng vấn), và an toàn tuyệt đối khi thuê NCC (Read-only, NDA, bàn giao mã nguồn 100%).")

# =========================================================================
# SLIDE 18: KIẾN NGHỊ BAN GIÁM ĐỐC PHÊ DUYỆT
# =========================================================================
s18 = prs.slides.add_slide(BL)
header(s18, "KIẾN NGHỊ BAN GIÁM ĐỐC XEM XÉT VÀ PHÊ DUYỆT", "PHẦN D: KẾ HOẠCH HÀNH ĐỘNG — ĐỀ XUẤT", "17")

proposals = [
    ("1", "PHÊ DUYỆT CHỦ TRƯƠNG PHƯƠNG ÁN 2: CHỦ ĐỘNG TRIỂN KHAI CÓ NCC CHUYÊN MÔN",
     "Đồng ý chủ trương: PTSC Quảng Ngãi chủ động xây dựng Trục tích hợp dữ liệu đồng bộ với Nền tảng TCT, có sự hỗ trợ của Nhà cung cấp "
     "chuyên môn về Data Engineering. Đảm bảo vừa phục vụ điều hành nội bộ, vừa hoàn thành KPI Chuyển đổi số năm 2026.", BLUE),
    ("2", "THÀNH LẬP TỔ CÔNG TÁC NỀN TẢNG DỮ LIỆU NỘI BỘ (BẮT ĐẦU GIAI ĐOẠN 1)",
     "Gồm Bộ phận CNTT chủ trì kỹ thuật và Key Users các phòng ban (Kế toán, Nhân sự, Vật tư, Dự án). Triển khai ngay Giai đoạn 1: "
     "Khảo sát hiện trạng, lập Bảng ánh xạ 29 Danh mục Dữ liệu chủ và hoàn thiện Hồ sơ yêu cầu kỹ thuật (TOR) — chi phí 0 VNĐ.", GOLD),
    ("3", "CHO PHÉP TIẾP XÚC NHÀ CUNG CẤP ĐỂ LẤY ĐỀ XUẤT VÀ BÁO GIÁ CẠNH TRANH",
     "Cho phép Tổ Công tác tiếp xúc 2–3 Nhà cung cấp chuyên nghiệp để nhận đề xuất giải pháp kỹ thuật và báo giá cạnh tranh. "
     "Hoàn thiện dự toán chi tiết trình Ban Giám đốc xem xét phê duyệt trước khi ký kết hợp đồng.", GREEN)
]
for i, (num, title, desc, col) in enumerate(proposals):
    y = Inches(1.5 + i * 1.75)
    shape(s18, Inches(0.8), y, Inches(11.7), Inches(1.55), WHITE, col, Pt(1.5), rounded=True)
    shape(s18, Inches(0.8), y, Inches(0.6), Inches(1.55), col)
    numbered_label(s18, Inches(0.9), y + Inches(0.55), num, col)
    text(s18, Inches(1.7), y + Inches(0.1), Inches(10.5), Inches(0.35), title, 12, True, col)
    shape(s18, Inches(1.7), y + Inches(0.48), Inches(10.5), Inches(0.01), GRAY_MED)
    text(s18, Inches(1.7), y + Inches(0.58), Inches(10.5), Inches(0.85), desc, 10.5, color=BLACK)

note(s18, "Kính thưa Ban Giám đốc, bộ phận CNTT kính đề xuất 3 chủ trương trên để triển khai ngay Giai đoạn 1 nội bộ, đảm bảo tiến độ chung của Tổng công ty. Trân trọng cảm ơn Ban Giám đốc!")

# Save
out_path = r"d:\My Profiles\DataPlatform\bao_cao_dataplatform_ptsc_qn_v7.pptx"
prs.save(out_path)
print(f"V7 saved: {out_path}")
print("Total slides: 18")
