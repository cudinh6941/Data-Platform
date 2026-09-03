"""
BÁO CÁO DATA PLATFORM PTSC QUẢNG NGÃI - VERSION 4
Phong cách: Corporate Executive — nghiêm túc, chuyên nghiệp, không lạm dụng icon
Bám sát Chỉ đạo số 9 của Ban Giám đốc
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ─────────────────────── COLOR PALETTE (Corporate) ───────────────────────
NAVY       = RGBColor(12, 35, 68)
NAVY_MID   = RGBColor(24, 59, 108)
BLUE       = RGBColor(30, 100, 190)
BLUE_SOFT  = RGBColor(230, 240, 252)
GREEN      = RGBColor(21, 128, 80)
GREEN_SOFT = RGBColor(232, 248, 239)
RED        = RGBColor(180, 40, 40)
RED_SOFT   = RGBColor(252, 235, 235)
GOLD       = RGBColor(180, 115, 0)
GOLD_SOFT  = RGBColor(255, 245, 225)
WHITE      = RGBColor(255, 255, 255)
BLACK      = RGBColor(28, 35, 48)
GRAY       = RGBColor(108, 118, 132)
GRAY_LIGHT = RGBColor(243, 245, 248)
GRAY_MED   = RGBColor(200, 208, 218)
TEAL       = RGBColor(13, 148, 136)

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BL = prs.slide_layouts[6]

# ─────────────── HELPER FUNCTIONS ───────────────
def shape(slide, l, t, w, h, fill, line_c=None, line_w=None, rounded=False):
    sh_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    r = slide.shapes.add_shape(sh_type, l, t, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = fill
    if line_c:
        r.line.color.rgb = line_c
        if line_w: r.line.width = line_w
    else:
        r.line.fill.background()
    return r

def text(slide, l, t, w, h, content, size=14, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = content
    p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = color; p.alignment = align
    return tf

def header(slide, title, subtitle="", slide_num=""):
    shape(slide, 0, 0, W, Inches(1.05), NAVY)
    shape(slide, 0, Inches(1.05), W, Inches(0.035), GOLD)
    if subtitle:
        text(slide, Inches(0.8), Inches(0.1), Inches(10), Inches(0.25), subtitle.upper(), 9, True, RGBColor(155, 185, 220))
    text(slide, Inches(0.8), Inches(0.38), Inches(11), Inches(0.55), title, 19, True, WHITE)
    if slide_num:
        text(slide, Inches(12.0), Inches(0.38), Inches(0.8), Inches(0.5), slide_num, 12, True, RGBColor(155, 185, 220), PP_ALIGN.RIGHT)
    # Footer line
    shape(slide, 0, Inches(7.2), W, Inches(0.008), GRAY_MED)
    text(slide, Inches(0.8), Inches(7.22), Inches(11.7), Inches(0.25),
         "PTSC Quảng Ngãi  |  Báo cáo Data Platform  |  Thực hiện Chỉ đạo số 9 Ban Giám đốc", 8, color=GRAY)

def card(slide, l, t, w, h, accent_color=BLUE, bg=WHITE, border=GRAY_MED):
    shape(slide, l, t, w, h, bg, border, Pt(0.8), True)
    shape(slide, l + Inches(0.04), t + Inches(0.02), Inches(0.08), h - Inches(0.04), accent_color)
    return l, t

def note(slide, txt_content):
    slide.notes_slide.notes_text_frame.text = txt_content

def numbered_label(slide, l, t, num, color=BLUE):
    """A small numbered circle — professional, not emoji."""
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, Inches(0.4), Inches(0.4))
    c.fill.solid(); c.fill.fore_color.rgb = color; c.line.fill.background()
    tf = c.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = str(num)
    p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER

def arrow_shape(slide, l, t, direction="right", w=Inches(0.5), h=Inches(0.3)):
    sh_type = MSO_SHAPE.RIGHT_ARROW if direction == "right" else MSO_SHAPE.DOWN_ARROW
    if direction == "down": w, h = h, w
    ar = slide.shapes.add_shape(sh_type, l, t, w, h)
    ar.fill.solid(); ar.fill.fore_color.rgb = GOLD; ar.line.fill.background()

# ═══════════════════════════════════════════════════════════════
# SLIDE 1: TRANG BÌA
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
shape(s, 0, 0, W, H, NAVY)
# Subtle accent stripe
shape(s, 0, Inches(1.4), W, Inches(0.04), GOLD)
shape(s, 0, Inches(5.8), W, Inches(0.015), RGBColor(60, 90, 140))

text(s, Inches(1.0), Inches(1.8), Inches(11), Inches(0.4),
     "BÁO CÁO BAN GIÁM ĐỐC  —  THỰC HIỆN CHỈ ĐẠO MỤC 9", 13, True, RGBColor(180, 200, 225))
text(s, Inches(1.0), Inches(2.5), Inches(11), Inches(1.6),
     "PHƯƠNG ÁN TRIỂN KHAI DATA PLATFORM\nVÀ TRỤC TÍCH HỢP DỮ LIỆU CÔNG TY – TCT", 30, True, WHITE)
# Divider
shape(s, Inches(1.0), Inches(4.4), Inches(2.5), Inches(0.04), GOLD)
# Scope points
scope_items = [
    "So sánh 2 Phương án: Thuê TCT  vs  Phát triển riêng đồng bộ TCT (có NCC tư vấn)",
    "Mô hình Kiến trúc tổng quan và Trục tích hợp phần mềm Công ty – TCT (Ưu tiên cao)",
    "Khung Chi phí Dự toán và Kế hoạch triển khai cụ thể"
]
for i, item in enumerate(scope_items):
    text(s, Inches(1.0), Inches(4.7 + i * 0.32), Inches(11), Inches(0.3), "—  " + item, 12, color=RGBColor(200, 215, 235))

text(s, Inches(1.0), Inches(6.1), Inches(11), Inches(0.3),
     "Đơn vị thực hiện:  Tổ Công tác CĐS & CNTT  –  PTSC Quảng Ngãi", 11.5, color=WHITE)
text(s, Inches(1.0), Inches(6.45), Inches(11), Inches(0.3),
     "Kính trình:  Ban Giám đốc Công ty PTSC Quảng Ngãi", 11.5, True, RGBColor(253, 224, 71))

note(s, "Kính thưa Ban Giám đốc, hôm nay bộ phận CNTT xin báo cáo 4 nội dung trọng tâm "
     "theo đúng chỉ đạo số 9 của Ban Giám đốc.")

# ═══════════════════════════════════════════════════════════════
# SLIDE 2: VẤN ĐỀ — DỮ LIỆU PHÂN MẢNH
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
header(s, "HIỆN TRẠNG: DỮ LIỆU ĐANG PHÂN MẢNH VÀ THỦ CÔNG", "BỐI CẢNH — VẤN ĐỀ CẦN GIẢI QUYẾT", "01")

# 4 silo boxes
silos = ["Phần mềm\nKẾ TOÁN", "Phần mềm\nNHÂN SỰ", "Phần mềm\nVẬT TƯ", "Quản lý\nDỰ ÁN"]
for i, name in enumerate(silos):
    x = Inches(0.8 + i * 2.95)
    card(s, x, Inches(1.5), Inches(2.55), Inches(1.3), BLUE)
    text(s, x + Inches(0.25), Inches(1.6), Inches(2.1), Inches(1.0), name, 13, True, NAVY, PP_ALIGN.CENTER)
    # X marks between silos
    if i < 3:
        text(s, Inches(3.1 + i * 2.95), Inches(1.85), Inches(0.5), Inches(0.4), "✕", 18, True, RED, PP_ALIGN.CENTER)

# Pain points
text(s, Inches(0.8), Inches(3.2), Inches(11.7), Inches(0.35),
     "Hệ quả trực tiếp đối với công tác điều hành:", 14, True, RED)

pains = [
    ("Chậm trễ báo cáo", "Các phòng ban mất 3–5 ngày\ntổng hợp số liệu thủ công\nbằng Excel mỗi tháng"),
    ("Số liệu không thống nhất", "Kế toán và Vật tư thường\nbáo cáo lệch số do\nlấy dữ liệu khác nguồn"),
    ("Phụ thuộc con người", "Nếu cán bộ phụ trách nghỉ,\nbáo cáo bị đình trệ.\nKhông có cơ chế tự động"),
    ("Thiếu tầm nhìn toàn cảnh", "Ban Giám đốc muốn xem\ntổng thể phải chờ rất lâu,\nkhông có công cụ trực quan"),
]
for i, (title, desc) in enumerate(pains):
    x = Inches(0.8 + i * 3.05)
    card(s, x, Inches(3.8), Inches(2.7), Inches(2.7), RED)
    numbered_label(s, x + Inches(0.2), Inches(3.9), i + 1, RED)
    text(s, x + Inches(0.7), Inches(3.9), Inches(1.8), Inches(0.35), title, 11.5, True, RED)
    text(s, x + Inches(0.25), Inches(4.45), Inches(2.2), Inches(1.8), desc, 10.5, color=BLACK)

note(s, "Đây là bức tranh thực tế: 4 phần mềm hoạt động rời rạc, không kết nối. "
     "Hậu quả là BGĐ muốn xem số liệu tổng thể phải chờ rất lâu, và số liệu thường lệch giữa các phòng.")

# ═══════════════════════════════════════════════════════════════
# SLIDE 3: DATA PLATFORM LÀ GÌ? (Trước – Sau)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
header(s, "DATA PLATFORM: TỪ \"ỐC ĐẢO DỮ LIỆU\" ĐẾN \"ĐƯỜNG ỐNG TỰ ĐỘNG\"", "GIẢI PHÁP — GIẢI THÍCH ĐƠN GIẢN", "02")

# LEFT: Before
shape(s, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.0), RED_SOFT, RED, Pt(1.5), True)
shape(s, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5), RED)
text(s, Inches(0.8), Inches(1.52), Inches(5.5), Inches(0.45), "TRƯỚC ĐÂY  —  Không có Data Platform", 12, True, WHITE, PP_ALIGN.CENTER)

before_items = [
    "Mỗi phần mềm là một \"ốc đảo\" riêng biệt",
    "Muốn có số liệu → nhân viên xuất Excel, copy-paste ghép file",
    "Mất 3–5 ngày mới có báo cáo tổng hợp",
    "Phòng Kế toán báo một số, Phòng Vật tư báo số khác",
    "Ban Giám đốc không có công cụ nhìn toàn cảnh tức thời",
]
for i, item in enumerate(before_items):
    text(s, Inches(1.3), Inches(2.25 + i * 0.7), Inches(4.6), Inches(0.6),
         "✕  " + item, 11.5, color=RED if i == 0 else BLACK)

# Arrow
arrow_shape(s, Inches(6.5), Inches(3.7))

# RIGHT: After
shape(s, Inches(7.2), Inches(1.5), Inches(5.3), Inches(5.0), GREEN_SOFT, GREEN, Pt(1.5), True)
shape(s, Inches(7.2), Inches(1.5), Inches(5.3), Inches(0.5), GREEN)
text(s, Inches(7.2), Inches(1.52), Inches(5.3), Inches(0.45), "SAU KHI CÓ  —  Data Platform", 12, True, WHITE, PP_ALIGN.CENTER)

after_items = [
    "Tất cả dữ liệu được gom tự động về một nơi duy nhất",
    "Hệ thống tự trích xuất, làm sạch, chuẩn hóa mỗi ngày",
    "Báo cáo Dashboard cập nhật tức thời, không chờ đợi",
    "Một nguồn sự thật duy nhất — mọi phòng ban đều nhìn cùng một số",
    "Ban Giám đốc xem trực quan trên máy tính hoặc điện thoại",
]
for i, item in enumerate(after_items):
    text(s, Inches(7.7), Inches(2.25 + i * 0.7), Inches(4.4), Inches(0.6),
         "✔  " + item, 11.5, color=GREEN if i == 0 else BLACK)

note(s, "Data Platform giống hệ thống đường ống nước tự động. Trước đây mỗi phòng ban là một cái giếng, "
     "muốn dùng nước phải xách xô – giờ nước sạch tự chảy thẳng tới vòi của Ban Giám đốc.")

# ═══════════════════════════════════════════════════════════════
# SLIDE 4: MÔ HÌNH HUB – SPOKE
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
header(s, "VỊ THẾ CỦA QUẢNG NGÃI TRONG HỆ SINH THÁI DATA TCT", "MÔ HÌNH — HUB & SPOKE (TRUNG TÂM – VỆ TINH)", "03")

# Central Hub
shape(s, Inches(4.2), Inches(2.0), Inches(4.9), Inches(2.8), NAVY, rounded=True)
text(s, Inches(4.2), Inches(2.2), Inches(4.9), Inches(0.45), "TỔNG CÔNG TY PTSC  (HUB)", 15, True, WHITE, PP_ALIGN.CENTER)
# Divider inside hub
shape(s, Inches(4.7), Inches(2.7), Inches(3.9), Inches(0.015), RGBColor(60, 90, 140))
tf_hub = text(s, Inches(4.5), Inches(2.85), Inches(4.5), Inches(1.6),
    "Đầu tư máy chủ Data Lakehouse", 11, color=RGBColor(186, 215, 245), align=PP_ALIGN.CENTER)
for item in ["Mua bản quyền Microsoft Fabric, Power BI", "Ban hành 29 chuẩn Master Data", "Cấp sẵn Workspace L3 cho Quảng Ngãi"]:
    p = tf_hub.add_paragraph()
    p.text = item; p.font.size = Pt(11); p.font.color.rgb = RGBColor(186, 215, 245); p.alignment = PP_ALIGN.CENTER

# Spoke: Quang Ngai (highlighted)
shape(s, Inches(0.5), Inches(2.6), Inches(3.3), Inches(1.8), GOLD_SOFT, GOLD, Pt(2), True)
shape(s, Inches(0.54), Inches(2.62), Inches(0.08), Inches(1.72), GOLD)
text(s, Inches(0.8), Inches(2.7), Inches(2.7), Inches(0.4), "PTSC QUẢNG NGÃI", 12, True, NAVY)
text(s, Inches(0.8), Inches(3.05), Inches(2.7), Inches(0.3), "Level 3 Spoke", 10, True, GOLD)
text(s, Inches(0.8), Inches(3.45), Inches(2.7), Inches(0.8),
     "Chỉ cần dựng\n\"Trạm trung chuyển\" dữ liệu\nrồi bơm lên TCT", 10.5, color=BLACK)

arrow_shape(s, Inches(3.85), Inches(3.25), "right", Inches(0.3), Inches(0.25))

# Other spokes
spoke_positions = [
    (Inches(9.5), Inches(2.4), "PTSC M&C"),
    (Inches(9.5), Inches(4.0), "PTSC Thanh Hóa"),
    (Inches(4.7), Inches(5.3), "PTSC Đình Vũ"),
    (Inches(1.5), Inches(5.3), "Đơn vị khác..."),
]
for x, y, lb in spoke_positions:
    shape(s, x, y, Inches(2.8), Inches(0.9), GRAY_LIGHT, GRAY_MED, Pt(0.8), True)
    text(s, x + Inches(0.15), y + Inches(0.15), Inches(2.5), Inches(0.55), lb, 10, color=GRAY, align=PP_ALIGN.CENTER)

# Key message bar
shape(s, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.55), BLUE_SOFT, BLUE, Pt(1))
text(s, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.55),
     "Quảng Ngãi không phải tự mua máy chủ đắt tiền  —  chỉ cần xây \"đường dẫn kết nối\" lên hạ tầng TCT đã sẵn có.",
     11.5, True, NAVY, PP_ALIGN.CENTER)

note(s, "TCT là trung tâm, đã đầu tư hàng triệu đô xây hạ tầng. "
     "Quảng Ngãi là vệ tinh, chỉ cần xây đường ống nối lên là tận dụng được toàn bộ.")

# ═══════════════════════════════════════════════════════════════
# SLIDE 5: ÁP LỰC TIẾN ĐỘ
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
header(s, "TCT ĐÃ HOÀN THÀNH NỀN MÓNG — ĐẾN LƯỢT QUẢNG NGÃI", "ÁP LỰC — LỘ TRÌNH CHIẾN LƯỢC TCT", "04")

phases = [
    ("2024 – 2025", "ĐÃ HOÀN THÀNH", GREEN, GREEN_SOFT,
     ["TCT xây xong Data Lakehouse", "Ban hành 29 Master Data", "Cấp Workspace L3 cho đơn vị"]),
    ("2026 – 2027", "ĐANG THỰC HIỆN", GOLD, GOLD_SOFT,
     ["Quảng Ngãi phải kết nối!", "Dựng Trạm trung chuyển", "Bơm dữ liệu lên TCT"]),
    ("2028 – 2030", "TƯƠNG LAI", GRAY, GRAY_LIGHT,
     ["Khai thác toàn diện", "AI/ML phân tích dự báo", "Báo cáo thông minh"]),
]
for i, (period, status, col, bg_c, items) in enumerate(phases):
    x = Inches(0.8 + i * 4.1)
    # Period header
    shape(s, x, Inches(1.5), Inches(3.7), Inches(0.65), col, rounded=True)
    text(s, x, Inches(1.53), Inches(3.7), Inches(0.3), period, 15, True, WHITE, PP_ALIGN.CENTER)
    text(s, x, Inches(1.82), Inches(3.7), Inches(0.25), status, 9, True, WHITE, PP_ALIGN.CENTER)
    # Body card
    shape(s, x, Inches(2.25), Inches(3.7), Inches(2.6), bg_c, col, Pt(1), True)
    for j, item in enumerate(items):
        text(s, x + Inches(0.25), Inches(2.45 + j * 0.55), Inches(3.2), Inches(0.5),
             "—  " + item, 12, color=BLACK)
    if i < 2:
        arrow_shape(s, Inches(4.55 + i * 4.1), Inches(2.9))

# Urgency callout
shape(s, Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.5), RED_SOFT, RED, Pt(1.5), True)
shape(s, Inches(0.84), Inches(5.32), Inches(0.08), Inches(1.42), RED)
text(s, Inches(1.2), Inches(5.4), Inches(10.8), Inches(0.35), "Tại sao phải triển khai ngay?", 14, True, RED)
tf_urg = text(s, Inches(1.2), Inches(5.85), Inches(10.8), Inches(0.35),
    "Nghị quyết số 10/NQ-HĐQT-PTSC: Data Platform là trụ cột bắt buộc cho toàn bộ đơn vị thành viên", 11, color=BLACK)
p = tf_urg.add_paragraph()
p.text = "Giai đoạn 2026–2027 là mốc deadline do TCT giao — chậm trễ sẽ ảnh hưởng đến đánh giá KPI chuyển đổi số của đơn vị"
p.font.size = Pt(11); p.font.color.rgb = BLACK

note(s, "TCT đã đi trước 2 năm, xây xong sân ga trung tâm và mở sẵn làn đường cho mình. "
     "2026-2027 là deadline bắt buộc.")

# ═══════════════════════════════════════════════════════════════
# SLIDE 6: 3 NÚT THẮT
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
header(s, "3 THÁCH THỨC CẦN GIẢI QUYẾT TẠI QUẢNG NGÃI", "HIỆN TRẠNG — ĐÁNH GIÁ NĂNG LỰC NỘI BỘ", "05")

knots = [
    ("DỮ LIỆU PHÂN MẢNH", GOLD, GOLD_SOFT,
     ["4 phần mềm hoạt động\nđộc lập, không kết nối", "Chuyển dữ liệu bằng\nExcel thủ công", "Chưa có cơ chế\ntập trung dữ liệu"]),
    ("LỆCH CHUẨN MASTER DATA", BLUE, BLUE_SOFT,
     ["Mã vật tư, nhân sự, dự án\ncủa QN chưa khớp chuẩn TCT", "Cần ánh xạ toàn bộ\nvới 29 danh mục Master Data", "Nếu không mapping trước,\nTCT sẽ từ chối tiếp nhận"]),
    ("THIẾU CHUYÊN GIA\nDATA ENGINEERING", RED, RED_SOFT,
     ["IT giỏi hạ tầng mạng\nvà hỗ trợ người dùng", "Chưa có kinh nghiệm\nData Pipeline / ETL", "Tự nghiên cứu:\nmất 8–12 tháng, rủi ro trễ"]),
]

for i, (title, col, bg_c, items) in enumerate(knots):
    x = Inches(0.8 + i * 4.1)
    shape(s, x, Inches(1.5), Inches(3.7), Inches(5.2), bg_c, col, Pt(1.5), True)
    # Accent bar
    shape(s, x + Inches(0.04), Inches(1.52), Inches(0.08), Inches(5.14), col)
    # Number
    numbered_label(s, x + Inches(0.3), Inches(1.7), i + 1, col)
    # Title
    text(s, x + Inches(0.85), Inches(1.72), Inches(2.6), Inches(0.7), title, 13, True, col)
    # Items
    for j, item in enumerate(items):
        y = Inches(2.65 + j * 1.4)
        shape(s, x + Inches(0.3), y, Inches(3.0), Inches(1.1), WHITE, GRAY_MED, Pt(0.5), True)
        text(s, x + Inches(0.5), y + Inches(0.1), Inches(2.6), Inches(0.9), item, 11, color=BLACK)

note(s, "3 thách thức: dữ liệu rời rạc, lệch chuẩn TCT, và thiếu kỹ sư Data chuyên sâu.")

# ═══════════════════════════════════════════════════════════════
# SLIDE 7: SƠ ĐỒ LUỒNG DỮ LIỆU (Diagram chính — ƯU TIÊN CAO)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
header(s, "TRỤC TÍCH HỢP: DỮ LIỆU CHẢY TỪ QUẢNG NGÃI LÊN TCT NHƯ THẾ NÀO?", "KIẾN TRÚC — NỘI DUNG ƯU TIÊN CAO THEO CHỈ ĐẠO", "06")

# Layer 1: Sources
shape(s, Inches(0.3), Inches(1.35), Inches(2.7), Inches(5.45), BLUE_SOFT, BLUE, Pt(1), True)
text(s, Inches(0.3), Inches(1.38), Inches(2.7), Inches(0.35), "NGUỒN DỮ LIỆU NỘI BỘ", 9, True, BLUE, PP_ALIGN.CENTER)

sources = ["Kế toán – Tài chính", "Quản lý Nhân sự", "Vật tư – Thiết bị", "Hợp đồng – Dự án"]
for i, name in enumerate(sources):
    y = Inches(1.95 + i * 1.2)
    shape(s, Inches(0.55), y, Inches(2.2), Inches(0.85), WHITE, BLUE, Pt(0.8), True)
    text(s, Inches(0.65), y + Inches(0.1), Inches(2.0), Inches(0.35), name, 10.5, True, NAVY, PP_ALIGN.CENTER)
    text(s, Inches(0.65), y + Inches(0.45), Inches(2.0), Inches(0.3), "Cơ sở dữ liệu", 8.5, color=GRAY, align=PP_ALIGN.CENTER)
    arrow_shape(s, Inches(2.8), y + Inches(0.25), "right", Inches(0.35), Inches(0.22))

# Layer 2: Gateway
shape(s, Inches(3.3), Inches(1.35), Inches(3.7), Inches(5.45), GOLD_SOFT, GOLD, Pt(1.5), True)
text(s, Inches(3.3), Inches(1.38), Inches(3.7), Inches(0.35), "TRẠM TRUNG CHUYỂN (đặt tại QN)", 9, True, GOLD, PP_ALIGN.CENTER)

gw_steps = [
    ("1.  Trích xuất", "Chế độ CHỈ ĐỌC\nKhông ảnh hưởng PM gốc"),
    ("2.  Chuẩn hóa", "Ánh xạ theo\n29 Master Data TCT"),
    ("3.  Kiểm tra chất lượng", "Lọc dữ liệu bẩn\nĐảm bảo chính xác"),
    ("4.  Mã hóa bảo mật", "Đóng gói AES-256\nSẵn sàng truyền đi"),
]
for i, (step_t, step_d) in enumerate(gw_steps):
    y = Inches(1.95 + i * 1.2)
    shape(s, Inches(3.55), y, Inches(3.2), Inches(0.85), WHITE, GOLD, Pt(0.8), True)
    numbered_label(s, Inches(3.65), y + Inches(0.22), i + 1, GOLD)
    text(s, Inches(4.15), y + Inches(0.08), Inches(2.5), Inches(0.35), step_t, 10.5, True, NAVY)
    text(s, Inches(4.15), y + Inches(0.42), Inches(2.5), Inches(0.4), step_d, 8.5, color=GRAY)

# Arrow Gateway → VPN
arrow_shape(s, Inches(7.1), Inches(3.65), "right", Inches(0.4), Inches(0.22))

# Layer 3: VPN
shape(s, Inches(7.6), Inches(2.3), Inches(1.4), Inches(3.0), GRAY_LIGHT, NAVY, Pt(1), True)
text(s, Inches(7.6), Inches(2.65), Inches(1.4), Inches(0.35), "Đường hầm", 9, True, NAVY, PP_ALIGN.CENTER)
text(s, Inches(7.6), Inches(2.95), Inches(1.4), Inches(0.35), "VPN bảo mật", 10, True, NAVY, PP_ALIGN.CENTER)
text(s, Inches(7.6), Inches(3.5), Inches(1.4), Inches(0.8), "Mã hóa\nSSL/IPSec\nSite-to-Site", 8.5, color=GRAY, align=PP_ALIGN.CENTER)

# Arrow VPN → TCT
arrow_shape(s, Inches(9.1), Inches(3.65), "right", Inches(0.4), Inches(0.22))

# Layer 4: TCT
shape(s, Inches(9.6), Inches(1.35), Inches(3.3), Inches(5.45), GREEN_SOFT, GREEN, Pt(1), True)
text(s, Inches(9.6), Inches(1.38), Inches(3.3), Inches(0.35), "TỔNG CÔNG TY (CLOUD)", 9, True, GREEN, PP_ALIGN.CENTER)

# Workspace
shape(s, Inches(9.85), Inches(1.95), Inches(2.8), Inches(1.6), WHITE, GREEN, Pt(0.8), True)
text(s, Inches(9.85), Inches(2.05), Inches(2.8), Inches(0.35), "Workspace L3", 11, True, NAVY, PP_ALIGN.CENTER)
text(s, Inches(9.85), Inches(2.4), Inches(2.8), Inches(0.3), "Quảng Ngãi", 10, True, BLUE, PP_ALIGN.CENTER)
text(s, Inches(9.85), Inches(2.75), Inches(2.8), Inches(0.7), "Lưu trữ trên\nMicrosoft Fabric\n(Cloud TCT)", 9, color=GRAY, align=PP_ALIGN.CENTER)

arrow_shape(s, Inches(11.08), Inches(3.65), "down", Inches(0.22), Inches(0.35))

# Dashboard
shape(s, Inches(9.85), Inches(4.2), Inches(2.8), Inches(2.3), WHITE, BLUE, Pt(1.5), True)
text(s, Inches(9.85), Inches(4.35), Inches(2.8), Inches(0.35), "Dashboard Power BI", 11, True, BLUE, PP_ALIGN.CENTER)
shape(s, Inches(10.1), Inches(4.75), Inches(2.3), Inches(0.015), GRAY_MED)
text(s, Inches(9.85), Inches(4.95), Inches(2.8), Inches(1.4),
     "Báo cáo trực quan\ncập nhật tự động\n\nBan Giám đốc xem\ntrên máy tính\nhoặc điện thoại", 9.5, color=BLACK, align=PP_ALIGN.CENTER)

note(s, "Trục tích hợp hoạt động 4 bước tự động: Hút → Chuẩn hóa → Kiểm tra → Mã hóa, "
     "rồi bơm qua VPN sang Cloud TCT. Chế độ chỉ đọc, không bao giờ ảnh hưởng phần mềm gốc.")

# ═══════════════════════════════════════════════════════════════
# SLIDE 8: CAM KẾT AN TOÀN
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
header(s, "3 NGUYÊN TẮC AN TOÀN KHI VẬN HÀNH TRỤC TÍCH HỢP", "AN TOÀN — ĐẢM BẢO CHO HỆ THỐNG HIỆN TẠI", "07")

safes = [
    ("KHÔNG ẢNH HƯỞNG\nPHẦN MỀM HIỆN TẠI", GREEN, GREEN_SOFT,
     ["Trạm trung chuyển chạy trên\nmáy chủ ảo riêng biệt",
      "Chỉ đọc bản sao dữ liệu,\nkhông can thiệp hệ thống gốc",
      "Chạy vào ban đêm hoặc\ngiờ thấp điểm"]),
    ("KHÔNG RÒ RỈ\nDỮ LIỆU RA NGOÀI", NAVY, BLUE_SOFT,
     ["Dữ liệu chỉ đi trong\nđường hầm VPN riêng của TCT",
      "Không mở cổng internet\ncông cộng",
      "Mã hóa AES-256\n(tiêu chuẩn ngân hàng)"]),
    ("TỰ ĐỘNG PHÁT HIỆN\nVÀ CẢNH BÁO LỖI", GOLD, GOLD_SOFT,
     ["Nếu nghẽn mạng hoặc\nsai lệch dữ liệu",
      "Hệ thống tự gửi email\ncảnh báo cho IT ngay",
      "Toàn bộ hoạt động\nđược ghi log kiểm toán"]),
]

for i, (title, col, bg_c, items) in enumerate(safes):
    x = Inches(0.8 + i * 4.1)
    shape(s, x, Inches(1.5), Inches(3.7), Inches(5.2), bg_c, col, Pt(1.5), True)
    shape(s, x + Inches(0.04), Inches(1.52), Inches(0.08), Inches(5.14), col)
    numbered_label(s, x + Inches(0.3), Inches(1.7), i + 1, col)
    text(s, x + Inches(0.85), Inches(1.7), Inches(2.6), Inches(0.7), title, 13, True, col)
    for j, item in enumerate(items):
        y = Inches(2.65 + j * 1.4)
        shape(s, x + Inches(0.3), y, Inches(3.0), Inches(1.1), WHITE, GRAY_MED, Pt(0.5), True)
        text(s, x + Inches(0.5), y + Inches(0.12), Inches(2.6), Inches(0.85), item, 11, color=BLACK)

note(s, "3 cam kết an toàn: chạy máy riêng không ảnh hưởng PM kế toán, đi đường truyền riêng không lộ ra internet, "
     "tự cảnh báo khi có lỗi.")

# ═══════════════════════════════════════════════════════════════
# SLIDE 9: SO SÁNH 2 PHƯƠNG ÁN
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
header(s, "SO SÁNH: THUÊ TCT LÀM HỘ  VS  PHÁT TRIỂN RIÊNG CÓ NCC", "PHÂN TÍCH — CÂU HỎI TRỌNG TÂM CỦA BAN GIÁM ĐỐC", "08")

# Column 1: PA1
shape(s, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3), RED_SOFT, RED, Pt(1.5), True)
shape(s, Inches(0.8), Inches(1.5), Inches(5.6), Inches(0.5), RED)
text(s, Inches(0.8), Inches(1.5), Inches(5.6), Inches(0.5),
     "PHƯƠNG ÁN 1:  Thuê dịch vụ của TCT", 12, True, WHITE, PP_ALIGN.CENTER)

pa1 = [
    ("TCT không đủ nguồn lực", "Ban CNTT TCT chỉ phụ trách hạ tầng lõi Hub.\nKhông có người xuống QN khảo sát từng CSDL."),
    ("Bị động về tiến độ", "Phải xếp hàng chờ TCT triển khai cho hàng chục\nđơn vị thành viên. Chắc chắn trễ hạn 2026."),
    ("Không có báo cáo riêng", "TCT chỉ xây Dashboard tài chính vĩ mô.\nKhông phục vụ bài toán quản trị tại QN."),
    ("Phải trả phí duy trì", "Gánh phân bổ chi phí phần mềm từ TCT,\nkhó kiểm soát ngân sách dài hạn."),
]
for i, (t, d) in enumerate(pa1):
    y = Inches(2.2 + i * 1.05)
    text(s, Inches(1.1), y, Inches(5.0), Inches(0.3), "✕  " + t, 11, True, RED)
    text(s, Inches(1.35), y + Inches(0.28), Inches(4.7), Inches(0.65), d, 9.5, color=BLACK)

# Column 2: PA2
shape(s, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3), GREEN_SOFT, GREEN, Pt(1.5), True)
shape(s, Inches(6.8), Inches(1.5), Inches(5.7), Inches(0.5), GREEN)
text(s, Inches(6.8), Inches(1.5), Inches(5.7), Inches(0.5),
     "PHƯƠNG ÁN 2:  Phát triển riêng + NCC   ĐỀ XUẤT", 12, True, WHITE, PP_ALIGN.CENTER)

pa2 = [
    ("NCC làm việc trực tiếp tại QN", "Đấu nối đúng từng phần mềm đặc thù\nmà công ty đang dùng hàng ngày."),
    ("Hoàn thành trong 3–4 tháng", "Đơn vị chủ động 100% tiến độ.\nĐảm bảo kịp mốc hạn TCT giao."),
    ("Có Dashboard riêng cho BGĐ", "Vừa đồng bộ chuẩn cho TCT, vừa có\nbáo cáo tùy biến phục vụ điều hành tại QN."),
    ("Chi 1 lần, tận dụng Cloud TCT", "Tiết kiệm tiền tỷ nhờ dùng miễn phí\nhạ tầng Cloud mà TCT đã mua."),
]
for i, (t, d) in enumerate(pa2):
    y = Inches(2.2 + i * 1.05)
    text(s, Inches(7.1), y, Inches(5.1), Inches(0.3), "✔  " + t, 11, True, GREEN)
    text(s, Inches(7.35), y + Inches(0.28), Inches(4.8), Inches(0.65), d, 9.5, color=BLACK)

note(s, "Câu hỏi trọng tâm: Tại sao không nhờ TCT? Vì TCT chỉ xây Hub lõi, không có người đi sâu vào PM của QN. "
     "Phương án 2 thuê NCC là tối ưu: làm nhanh, có báo cáo riêng cho Sếp, và chủ động tiến độ.")

# ═══════════════════════════════════════════════════════════════
# SLIDE 10: KẾT QUẢ — BGĐ ĐƯỢC GÌ?
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
header(s, "KẾT QUẢ ĐẠT ĐƯỢC KHI HOÀN THÀNH DỰ ÁN", "ĐẦU RA — GIÁ TRỊ CHO BAN GIÁM ĐỐC VÀ ĐƠN VỊ", "09")

outputs = [
    ("BÁO CÁO ĐIỀU HÀNH\nTRỰC QUAN", BLUE, BLUE_SOFT,
     "Dashboard cập nhật mỗi ngày.\nXem trên máy tính hoặc\nđiện thoại, mọi lúc mọi nơi."),
    ("DỮ LIỆU TỰ ĐỘNG\nKHÔNG PHỤ THUỘC EXCEL", TEAL, RGBColor(204, 251, 241),
     "Phòng ban không cần\nxuất file thủ công.\nSố liệu tự chảy về hệ thống."),
    ("MỘT NGUỒN SỰ THẬT\nDUY NHẤT", GREEN, GREEN_SOFT,
     "Kế toán, Vật tư, Nhân sự\nđều nhìn cùng một con số.\nHết tranh cãi lệch số liệu."),
    ("ĐẠT CHỈ TIÊU CĐS\nTCT GIAO", GOLD, GOLD_SOFT,
     "Hoàn thành đúng hạn\nnhiệm vụ kết nối Data Platform\ntheo Nghị quyết HĐQT."),
]

for i, (title, col, bg_c, desc) in enumerate(outputs):
    x = Inches(0.8 + i * 3.05)
    shape(s, x, Inches(1.5), Inches(2.7), Inches(5.0), bg_c, col, Pt(1.5), True)
    shape(s, x + Inches(0.04), Inches(1.52), Inches(0.08), Inches(4.94), col)
    numbered_label(s, x + Inches(0.25), Inches(1.7), i + 1, col)
    text(s, x + Inches(0.15), Inches(2.3), Inches(2.4), Inches(0.8), title, 13, True, col, PP_ALIGN.CENTER)
    # Divider
    shape(s, x + Inches(0.4), Inches(3.2), Inches(1.9), Inches(0.01), GRAY_MED)
    text(s, x + Inches(0.2), Inches(3.4), Inches(2.3), Inches(2.8), desc, 11.5, color=BLACK, align=PP_ALIGN.CENTER)

note(s, "4 đầu ra cụ thể: Dashboard tức thời, tự động hóa, một nguồn sự thật, và đạt KPI chuyển đổi số.")

# ═══════════════════════════════════════════════════════════════
# SLIDE 11: DỰ TOÁN CHI PHÍ
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
header(s, "KHUNG DỰ TOÁN CHI PHÍ SƠ BỘ", "NGÂN SÁCH — CHỦ ĐỘNG TRƯỚC KHI LÀM VIỆC VỚI NCC", "10")

# Table
table_shape = s.shapes.add_table(6, 4, Inches(0.8), Inches(1.5), Inches(11.733), Inches(3.5))
tbl = table_shape.table
tbl.columns[0].width = Inches(0.5)
tbl.columns[1].width = Inches(4.0)
tbl.columns[2].width = Inches(5.0)
tbl.columns[3].width = Inches(2.233)

headers_data = ["", "Hạng mục", "Phạm vi thực hiện", "Dự toán"]
for i, h in enumerate(headers_data):
    c = tbl.cell(0, i); c.fill.solid(); c.fill.fore_color.rgb = NAVY
    p = c.text_frame.paragraphs[0]; p.text = h
    p.font.size = Pt(11); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER

rows_data = [
    ("1", "Hạ tầng máy chủ (VM)", "Tận dụng máy chủ hiện có của Công ty", "0 đ  (Tự có)"),
    ("2", "Khảo sát & Thiết kế kiến trúc", "Khảo sát CSDL + Ánh xạ 29 Master Data TCT", "80 – 120 Triệu"),
    ("3", "Xây dựng Trục tích hợp & ETL", "Đường ống trích xuất, làm sạch, mã hóa, bơm lên TCT", "150 – 250 Triệu"),
    ("4", "Dashboard Power BI", "3–5 bảng báo cáo quản trị cho Ban Giám đốc", "100 – 150 Triệu"),
    ("5", "Đào tạo + Bảo hành 12 tháng", "Bàn giao mã nguồn, đào tạo IT làm chủ 100%", "50 – 80 Triệu"),
]
for ri, (num, item, scope, cost) in enumerate(rows_data, 1):
    row = [num, item, scope, cost]
    for ci, val in enumerate(row):
        c = tbl.cell(ri, ci); c.fill.solid()
        c.fill.fore_color.rgb = WHITE if ri % 2 == 1 else GRAY_LIGHT
        p = c.text_frame.paragraphs[0]; p.text = val
        p.font.size = Pt(10.5); p.font.color.rgb = BLACK
        if ci == 0 or ci == 3: p.alignment = PP_ALIGN.CENTER
        if ci == 3: p.font.bold = True; p.font.color.rgb = BLUE

# Total bar
shape(s, Inches(0.8), Inches(5.25), Inches(11.733), Inches(0.7), NAVY, rounded=True)
text(s, Inches(1.3), Inches(5.3), Inches(5), Inches(0.35), "TỔNG KHÁI TOÁN DỰ KIẾN", 13, True, WHITE)
text(s, Inches(8.0), Inches(5.25), Inches(4.3), Inches(0.4), "380 – 600 Triệu VNĐ", 20, True, RGBColor(253, 224, 71), PP_ALIGN.RIGHT)
text(s, Inches(8.0), Inches(5.6), Inches(4.3), Inches(0.3), "Thanh toán phân kỳ theo mốc nghiệm thu", 9.5, color=RGBColor(186, 215, 245), align=PP_ALIGN.RIGHT)

# Note box
shape(s, Inches(0.8), Inches(6.15), Inches(11.733), Inches(0.85), BLUE_SOFT, BLUE, Pt(1), True)
shape(s, Inches(0.84), Inches(6.17), Inches(0.07), Inches(0.78), BLUE)
tf_note = text(s, Inches(1.15), Inches(6.2), Inches(11.0), Inches(0.3),
    "Tiết kiệm hàng tỷ đồng nhờ tận dụng hạ tầng Cloud và bản quyền Power BI mà TCT đã đầu tư sẵn", 10.5, True, NAVY)
p = tf_note.add_paragraph()
p.text = "Cơ chế giải ngân: GĐ1 chuẩn bị nội bộ (0 đ)  →  GĐ2 kết nối xong (70%)  →  GĐ3 bàn giao (30%)"
p.font.size = Pt(10.5); p.font.color.rgb = BLACK

note(s, "Tổng chi phí khoảng 380-600 triệu. So với tự mua Cloud Server và bản quyền thì tiết kiệm hàng tỷ. "
     "Giải ngân theo mốc nghiệm thu, không trả 1 cục.")

# ═══════════════════════════════════════════════════════════════
# SLIDE 12: LỘ TRÌNH 3 GIAI ĐOẠN
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
header(s, "KẾ HOẠCH TRIỂN KHAI 3 GIAI ĐOẠN", "LỘ TRÌNH — PHÂN KỲ RÕ RÀNG, KIỂM SOÁT CHẶT CHẼ", "11")

timeline = [
    ("GIAI ĐOẠN 1", "Tháng 3 – 4 / 2026", "KHẢO SÁT NỘI BỘ", "Nội bộ tự chủ trì  •  Chi phí 0 đ",
     BLUE, BLUE_SOFT,
     ["Thành lập Tổ công tác", "Khảo sát CSDL 4 phần mềm", "Ánh xạ 29 Master Data", "Lập đề bài kỹ thuật cho NCC"]),
    ("GIAI ĐOẠN 2", "Tháng 5 – 7 / 2026", "THUÊ NCC TRIỂN KHAI", "Giai đoạn chính  •  Thuê NCC",
     GOLD, GOLD_SOFT,
     ["Lựa chọn NCC tối ưu", "Dựng Trạm trung chuyển", "Lập trình đường ống ETL", "Kiểm thử kết nối TCT"]),
    ("GIAI ĐOẠN 3", "Tháng 8 / 2026 →", "BÀN GIAO & KHAI THÁC", "IT làm chủ 100%",
     GREEN, GREEN_SOFT,
     ["Nghiệm thu hệ thống", "NCC bàn giao mã nguồn", "Đào tạo IT vận hành", "Xây Dashboard cho BGĐ"]),
]

for i, (phase, period, title, tag, col, bg_c, items) in enumerate(timeline):
    x = Inches(0.8 + i * 4.1)
    shape(s, x, Inches(1.5), Inches(3.7), Inches(5.2), bg_c, col, Pt(1.5), True)
    shape(s, x + Inches(0.04), Inches(1.52), Inches(0.08), Inches(5.14), col)

    text(s, x + Inches(0.25), Inches(1.65), Inches(3.2), Inches(0.3), phase, 10, True, col)
    text(s, x + Inches(0.25), Inches(1.92), Inches(3.2), Inches(0.3), period, 13, True, NAVY)
    text(s, x + Inches(0.25), Inches(2.3), Inches(3.2), Inches(0.35), title, 14, True, col)

    # Tag
    shape(s, x + Inches(0.25), Inches(2.75), Inches(3.1), Inches(0.35), col, rounded=True)
    text(s, x + Inches(0.25), Inches(2.75), Inches(3.1), Inches(0.35), tag, 9, True, WHITE, PP_ALIGN.CENTER)

    for j, item in enumerate(items):
        text(s, x + Inches(0.4), Inches(3.35 + j * 0.55), Inches(3.0), Inches(0.45),
             "—  " + item, 11.5, color=BLACK)

    if i < 2:
        arrow_shape(s, Inches(4.55 + i * 4.1), Inches(3.8))

note(s, "3 bước: Bước 1 nội bộ tự làm miễn phí 2 tháng; Bước 2 thuê NCC dựng trạm; "
     "Bước 3 nghiệm thu và đào tạo IT làm chủ.")

# ═══════════════════════════════════════════════════════════════
# SLIDE 13: KIẾN NGHỊ PHÊ DUYỆT
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
header(s, "KIẾN NGHỊ BAN GIÁM ĐỐC PHÊ DUYỆT", "ĐỀ XUẤT — 3 NỘI DUNG XIN CHỦ TRƯƠNG", "12")

proposals = [
    ("CHẤP THUẬN\nPHƯƠNG ÁN 2",
     "Phát triển riêng Trục tích hợp\nđồng bộ với TCT, có thuê\nNhà cung cấp chuyên nghiệp\ntư vấn và triển khai.",
     BLUE, BLUE_SOFT),
    ("THÀNH LẬP\nTỔ CÔNG TÁC",
     "IT chủ trì kỹ thuật +\nKey Users từ Kế toán,\nVật tư, Nhân sự\nbắt đầu Giai đoạn 1 ngay.",
     GOLD, GOLD_SOFT),
    ("CHO PHÉP\nLIÊN HỆ NCC",
     "Khảo sát thị trường,\nlấy báo giá cạnh tranh\ntheo khung kiến trúc\nđã xây dựng sẵn.",
     GREEN, GREEN_SOFT),
]

for i, (title, desc, col, bg_c) in enumerate(proposals):
    x = Inches(0.8 + i * 4.1)
    shape(s, x, Inches(1.5), Inches(3.7), Inches(4.5), bg_c, col, Pt(2), True)
    shape(s, x + Inches(0.04), Inches(1.52), Inches(0.08), Inches(4.44), col)

    numbered_label(s, x + Inches(0.3), Inches(1.7), i + 1, col)
    text(s, x + Inches(0.85), Inches(1.7), Inches(2.6), Inches(0.7), title, 14, True, col)
    # Divider
    shape(s, x + Inches(0.3), Inches(2.55), Inches(3.1), Inches(0.01), GRAY_MED)
    text(s, x + Inches(0.3), Inches(2.75), Inches(3.1), Inches(2.8), desc, 12.5, color=BLACK)

# Bottom CTA
shape(s, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.65), GREEN_SOFT, GREEN, Pt(1.5), True)
shape(s, Inches(0.84), Inches(6.32), Inches(0.07), Inches(0.58), GREEN)
text(s, Inches(1.1), Inches(6.3), Inches(11.2), Inches(0.65),
     "Kính trình Ban Giám đốc xem xét, thông qua chủ trương để đơn vị kịp tiến độ kết nối năm 2026 của TCT.",
     12, True, RGBColor(22, 100, 52), PP_ALIGN.CENTER)

note(s, "3 đề xuất: Chọn phương án 2, thành lập tổ công tác, và cho phép liên hệ NCC lấy báo giá.")

# ═══════════════════════════════════════════════════════════════
# SLIDE 14: TRANG KẾT THÚC
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
shape(s, 0, 0, W, H, NAVY)
shape(s, 0, Inches(2.2), W, Inches(0.035), GOLD)
shape(s, 0, Inches(5.0), W, Inches(0.015), RGBColor(60, 90, 140))

text(s, Inches(1.0), Inches(2.6), Inches(11.3), Inches(1.0),
     "XIN TRÂN TRỌNG CẢM ƠN\nBAN GIÁM ĐỐC", 28, True, WHITE)
text(s, Inches(1.0), Inches(4.0), Inches(11.3), Inches(0.5),
     "Kính mời Ban Giám đốc cho ý kiến chỉ đạo.", 15, color=RGBColor(200, 215, 235))

text(s, Inches(1.0), Inches(5.3), Inches(11.3), Inches(0.4),
     "Tài liệu đính kèm:  Kế hoạch chi tiết Giai đoạn 1, 2, 3", 11, color=RGBColor(155, 185, 220))
text(s, Inches(1.0), Inches(5.7), Inches(11.3), Inches(0.4),
     "Đơn vị thực hiện:  Tổ Công tác CĐS – Bộ phận CNTT PTSC Quảng Ngãi", 11, color=RGBColor(155, 185, 220))

note(s, "Chúng tôi xin trân trọng cảm ơn và lắng nghe ý kiến chỉ đạo của Ban Giám đốc!")

# ═══════════════ SAVE ═══════════════
output = r"d:\My Profiles\DataPlatform\bao_cao_dataplatform_ptsc_qn_v4.pptx"
prs.save(output)
print(f"SUCCESS: {output}")
print(f"Total slides: {len(prs.slides)}")
